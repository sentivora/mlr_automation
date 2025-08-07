"""
PowerPoint presentation generator using python-pptx library.
Creates slides with folder names as titles and arranges images on each slide.
"""

import os
import tempfile
import logging
import math
from datetime import datetime
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from .image_processor import ImageProcessor


class PresentationGenerator:
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        self.image_processor = ImageProcessor(self.logger)
        self.folder_mapping = {
            'ott': 'OTT',
            'vdxdesktopexpandable': 'DESKTOP EXPANDABLE',
            'vdxdesktopinframe': 'DESKTOP IN-FRAME',
            'vdxinstream': 'VDX DESKTOP INSTREAM',
            'vdxdesktopinstream': 'DESKTOP INSTREAM',
            'vdxmobileexpandable': 'MOBILE EXPANDABLE',
            'vdxmobileinframe': 'MOBILE IN-FRAME',
            'vdxmobileinstream': 'MOBILE INSTREAM',
            'ctv': 'CTV'
        }
    
    def _sort_images_exact_priority(self, img_path):
        """Sort images with exact filename priority for regular slides (excluding consolidated teaser slides and FULL ISI):
        1. teaser.png (exact filename)
        2. mainunit.png (exact filename)  
        3. All other images in sequential order
        If no teaser.png or mainunit.png exist, use sequential order
        """
        filename = os.path.basename(img_path).lower()
        if filename == 'teaser.png':
            return (0, filename)  # Exact teaser.png gets highest priority
        elif filename == 'mainunit.png':
            return (1, filename)  # Exact mainunit.png gets second priority
        else:
            return (2, filename)  # All other images in alphabetical order
    
    def _sort_images_engaged_priority(self, img_path):
        """Sort images for engaged folder with specific priority:
        1. vmp.png (if available)
        2. mainunit.png (if available) 
        3. All other images in sequential order
        Excludes mainunit-disclaimer.png as it's used in FULL ISI slide
        """
        filename = os.path.basename(img_path).lower()
        if filename == 'vmp.png':
            return (0, filename)  # vmp.png gets highest priority
        elif filename == 'mainunit.png':
            return (1, filename)  # mainunit.png gets second priority
        elif filename == 'mainunit-disclaimer.png':
            return (999, filename)  # Exclude disclaimer files
        else:
            return (2, filename)  # All other images in sequential order
    
    def _sort_images_desktop_expandable_priority(self, img_path, image_list):
        """Sort images for Desktop Expandable slides with special handling:
        If exactly 3 images with vpm and mainunit (no teaser):
        1. Image with 'mainunit' in filename (first)
        2. Image with 'vpm' in filename (second)
        3. Other image (third)
        Otherwise, use normal priority:
        1. teaser.png
        2. mainunit.png
        3. All other images in sequential order
        """
        filename = os.path.basename(img_path).lower()
        
        # Check if we have exactly 3 images with vpm and mainunit, but no teaser
        filenames = [os.path.basename(img).lower() for img in image_list]
        has_vpm = any('vpm' in fname for fname in filenames)
        has_mainunit = any('mainunit' in fname for fname in filenames)
        has_teaser = any('teaser' in fname for fname in filenames)
        
        if len(image_list) == 3 and has_vpm and has_mainunit and not has_teaser:
            # Special case: mainunit first, VPM second, other third
            if 'mainunit' in filename:
                return (0, filename)
            elif 'vpm' in filename:
                return (1, filename)
            else:
                return (2, filename)
        else:
            # Normal priority: teaser first, mainunit second, others alphabetical
            if filename == 'teaser.png':
                return (0, filename)
            elif filename == 'mainunit.png':
                return (1, filename)
            else:
                return (2, filename)
    
    def _format_folder_name(self, folder_name):
        """Format folder name according to the specified logic."""
        # Split by both forward and backward slashes
        parts = folder_name.replace('\\', '/').split('/')
        
        if len(parts) == 1:
            # Single folder name
            base_name = parts[0].lower()
            return self.folder_mapping.get(base_name, folder_name)
        elif len(parts) == 2:
            # Folder with subfolder (e.g., "vdxdesktopexpandable/300x250")
            base_name = parts[0].lower()
            sub_name = parts[1].lower()
            
            # Special cases where we don't want the subfolder name
            if base_name in ['ott', 'ctv', 'vdxinstream', 'vdxdesktopinstream', 'vdxmobileinstream'] and sub_name == '1x10':
                return self.folder_mapping.get(base_name, parts[0])
            
            formatted_base = self.folder_mapping.get(base_name, parts[0])
            return f"{formatted_base} - {sub_name}"
        else:
            # Handle nested folders by using the last two parts
            base_name = parts[-2].lower()
            sub_name = parts[-1].lower()
            
            # Special cases where we don't want the subfolder name
            if base_name in ['ott', 'ctv', 'vdxinstream', 'vdxdesktopinstream', 'vdxmobileinstream'] and sub_name == '1x10':
                return self.folder_mapping.get(base_name, parts[-2])
            
            formatted_base = self.folder_mapping.get(base_name, parts[-2])
            return f"{formatted_base} - {sub_name}"
    
    def generate_from_folder(self, temp_dir, annotation_option='with_annos', implement_video_frames=False, video_position_params=None, original_filename=None):
        """
        Generate presentation from a folder containing images.
        
        Args:
            temp_dir (str): Path to temporary directory containing extracted files
            annotation_option (str): Either 'with_annos' or 'no_annos'
            implement_video_frames (bool): Whether to implement video frames
            video_position_params (dict): Video positioning parameters
            original_filename (str): Original uploaded file name (without extension)
            
        Returns:
            str: Path to the generated presentation file
        """
        # Organize folder structure
        folder_structure = self._organize_folder_structure(temp_dir)
        
        # Create outputs directory if it doesn't exist
        output_dir = 'outputs'
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate presentation
        ppt_path, slide_count, video_folder_found = self.create_presentation(
            folder_structure=folder_structure,
            output_dir=output_dir,
            annotation_option=annotation_option,
            is_multi_tab=False,  # Always use Manual tab logic
            implement_video_frames=implement_video_frames,
            video_position_params=video_position_params,
            original_filename=original_filename
        )
        
        return ppt_path
    
    def _organize_folder_structure(self, temp_dir):
        """
        Organize files from temp directory into folder structure.
        
        Args:
            temp_dir (str): Path to temporary directory
            
        Returns:
            dict: Dictionary with folder names as keys and image paths as values
        """
        folder_structure = {}
        
        # Walk through all subdirectories
        for root, dirs, files in os.walk(temp_dir):
            # Skip if no image files
            image_files = [f for f in files if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp'))]
            if not image_files:
                continue
                
            # Get relative path from temp_dir
            rel_path = os.path.relpath(root, temp_dir)
            if rel_path == '.':
                # Files in root directory
                folder_name = 'root'
            else:
                folder_name = rel_path
            
            # Store full paths to image files
            image_paths = [os.path.join(root, f) for f in image_files]
            folder_structure[folder_name] = image_paths
        
        return folder_structure

    def create_presentation(self, folder_structure, output_dir, annotation_option='with_annos', is_multi_tab=False, implement_video_frames=False, video_position_params=None, original_filename=None):
        """
        Create a PowerPoint presentation from folder structure.
        
        Args:
            folder_structure (dict): Dictionary with folder names as keys and image paths as values
            output_dir (str): Directory to save the presentation
            annotation_option (str): Either 'with_annos' or 'no_annos' to control annotation display
            is_multi_tab (bool): Legacy parameter - always uses Manual tab logic now
            implement_video_frames (bool): Whether to implement video frames for all units
            
        Returns:
            tuple: (str, int, bool) Path to the created presentation file, slide count, and video folder found
        """
        self.logger.info("=== PRESENTATION GENERATOR ENTRY ===")
        self.logger.info(f"PARAMS: annotation_option={annotation_option}")
        self.logger.info(f"PARAMS: implement_video_frames={implement_video_frames}")
        self.logger.info(f"PARAMS: video_position_params={video_position_params}")



        
        # Store video position parameters for use in video frames functions
        self.video_position_params = video_position_params or {}
        try:
            # Create a new presentation
            prs = Presentation()
            
            # Set slide size to 16:9 widescreen
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            # Title slide removed as per user request
            
            # Collect all mainunit-disclaimer.png files
            mainunit_disclaimer_files = []
            
            # Process slides in specific order
            all_desktop_teaser_images = []
            all_mobile_teaser_images = []
            
            # First, collect mainunit-disclaimer.png files and teaser images
            for folder_name, image_paths in folder_structure.items():
                # Find mainunit-disclaimer.png files in this folder
                for img_path in image_paths:
                    if os.path.basename(img_path).lower() == 'mainunit-disclaimer.png':
                        mainunit_disclaimer_files.append(img_path)
                
                # Collect teaser image from vdxdesktopexpandable folders for consolidation
                if 'vdxdesktopexpandable' in folder_name.lower():
                    # Filter out disclaimer and engaged images
                    filtered_images = [
                        img_path for img_path in image_paths 
                        if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                               os.path.basename(img_path).lower() == 'mainunit-disclaimer.png' or
                               'engaged' in img_path.lower())
                    ]
                    # Prioritize images with "teaser" in filename, otherwise use first image
                    if filtered_images:
                        # First, look for images with "teaser" in the filename
                        teaser_images = [img for img in filtered_images if 'teaser' in os.path.basename(img).lower()]
                        
                        if teaser_images:
                            # Use the first teaser image found
                            teaser_images.sort(key=lambda x: os.path.basename(x).lower())
                            all_desktop_teaser_images.append(teaser_images[0])
                        else:
                            # No teaser images found, use the first image regardless of name
                            filtered_images.sort(key=lambda x: os.path.basename(x).lower())
                            all_desktop_teaser_images.append(filtered_images[0])
                
                # Collect teaser images from vdxmobileexpandable folders for consolidation
                if 'vdxmobileexpandable' in folder_name.lower():
                    # Filter out disclaimer and engaged images
                    filtered_images = [
                        img_path for img_path in image_paths 
                        if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                               os.path.basename(img_path).lower() == 'mainunit-disclaimer.png' or
                               'engaged' in img_path.lower())
                    ]
                    
                    # Prioritize images with "teaser" in filename, otherwise use all images
                    if filtered_images:
                        # First, look for images with "teaser" in the filename
                        teaser_images = [img for img in filtered_images if 'teaser' in os.path.basename(img).lower()]
                        
                        if teaser_images:
                            # Use all teaser images found
                            teaser_images.sort(key=lambda x: os.path.basename(x).lower())
                            all_mobile_teaser_images.extend(teaser_images)
                        else:
                            # No teaser images found, use all images from this folder
                            filtered_images.sort(key=lambda x: os.path.basename(x).lower())
                            all_mobile_teaser_images.extend(filtered_images)
            
            # Check for video folder and add as first slide if exists
            video_folder_processed = self._add_video_frames_slide_if_exists(prs, folder_structure, annotation_option)
            
            # Add slides using Manual tab logic (only approach now)
            self.logger.info("Using Manual tab slide processing logic")
            self._add_slides_in_order_manual(prs, folder_structure, annotation_option, all_desktop_teaser_images, all_mobile_teaser_images, implement_video_frames)
            
            # Add FULL ISI slide with all mainunit-disclaimer.png files (always last)
            if mainunit_disclaimer_files:
                from .slide_creator import SlideCreator
                slide_creator = SlideCreator(self, self.image_processor)
                result = slide_creator._create_full_isi_slide(prs, mainunit_disclaimer_files)
                if isinstance(result, int):
                    self.logger.info(f"Created {result} FULL ISI slides with image splitting")
                else:
                    self.logger.info("Created single FULL ISI slide")
            
            # Post-processing: Remove duplicate slides
            self._remove_duplicate_slides(prs)
            
            # Save the presentation
            # Use original filename if provided, otherwise use date-based naming
            if original_filename:
                filename = f"{original_filename}.pptx"
            else:
                # Fallback to date-based naming
                now = datetime.now()
                date_part = now.strftime("%d-%B-%Y")
                sequence = now.strftime("%S").zfill(3)
                filename = f"{date_part}-{sequence}.pptx"
            
            output_path = os.path.join(output_dir, filename)
            
            prs.save(output_path)
            
            # Get actual slide count (no title slide now)
            actual_slide_count = len(prs.slides)
            
            self.logger.info(f"Presentation saved to {output_path}")
            self.logger.info(f"Total slides created: {len(prs.slides)}")
            
            return output_path, actual_slide_count, video_folder_processed
            
        except Exception as e:
            self.logger.error(f"Error creating presentation: {str(e)}")
            raise
    
    def _remove_duplicate_slides(self, prs):
        """Remove duplicate slides from presentation - identify legitimate vs duplicate slide sequences."""
        try:
            self.logger.info("Starting smart duplicate slide removal for Manual tab")
            
            # Get all slide titles
            all_slides = []
            for i, slide in enumerate(prs.slides):
                title = self._get_slide_title(slide)
                all_slides.append((i, title))
                self.logger.info(f"Slide {i}: {title}")
            
            # Define the expected Manual tab slide sequence
            expected_sequence = [
                'CTV',
                'DESKTOP INSTREAM', 
                'MOBILE INSTREAM',
                'OTT',
                'DESKTOP EXPANDABLE - ALL TEASERS',
                'DESKTOP EXPANDABLE - VPM',
                'DESKTOP EXPANDABLE - ENGAGED',
                'MOBILE EXPANDABLE - ALL TEASERS', 
                'MOBILE EXPANDABLE - ENGAGED',
                'FULL ISI'
            ]
            
            # Find where legitimate sequence ends
            legitimate_end_idx = -1
            found_sequence = []
            
            for i, (slide_idx, title) in enumerate(all_slides):
                if title:
                    # Check if this title matches any expected sequence item
                    base_title = title.split(' VIDEO FRAME')[0].split(' (Contd.)')[0]
                    if any(expected in base_title for expected in expected_sequence):
                        found_sequence.append((slide_idx, title))
                        legitimate_end_idx = slide_idx
            
            self.logger.info(f"Legitimate sequence ends at slide {legitimate_end_idx}")
            self.logger.info(f"Found {len(found_sequence)} legitimate slides")
            
            # Mark slides after legitimate sequence for removal
            slides_to_remove = []
            for slide_idx, title in all_slides:
                if slide_idx > legitimate_end_idx and title:
                    # Only remove if it's a duplicate of something in legitimate sequence
                    base_title = title.split(' VIDEO FRAME')[0].split(' (Contd.)')[0]
                    is_duplicate = any(expected in base_title for expected in expected_sequence)
                    
                    if is_duplicate:
                        slides_to_remove.append(slide_idx)
                        self.logger.info(f"Marking duplicate slide {slide_idx} for removal: {title}")
            
            # Remove slides in reverse order
            slides_to_remove.sort(reverse=True)
            for slide_idx in slides_to_remove:
                slide_part = prs.slides._sldIdLst[slide_idx]
                prs.slides._sldIdLst.remove(slide_part)
                self.logger.info(f"Removed duplicate slide {slide_idx}")
            
            self.logger.info(f"Smart cleanup complete. Removed {len(slides_to_remove)} duplicate slides")
            
        except Exception as e:
            self.logger.error(f"Error removing duplicate slides: {str(e)}")
    
    def _get_slide_title(self, slide):
        """Extract title text from slide."""
        try:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    return shape.text.strip()
        except:
            pass
        return None
    
    def _apply_multi_tab_logic_to_folder(self, prs, folder_name, image_paths, annotation_option):
        """Apply Multi-tab logic to a specific folder to create additional slides for remaining images."""
        # Filter out disclaimer images
        filtered_image_paths = [
            img_path for img_path in image_paths 
            if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                   os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
        ]
        
        if len(filtered_image_paths) <= 2:
            return  # No additional slides needed
            
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        filtered_image_paths.sort(key=self._sort_images_exact_priority)
        
        # Find exact teaser.png and mainunit.png files, plus other images
        teaser_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'teaser.png']
        mainunit_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'mainunit.png']
        other_images = [img for img in filtered_image_paths if os.path.basename(img).lower() not in ['teaser.png', 'mainunit.png']]
        
        # Only create additional slides if there are "other" images
        if not other_images:
            return
            
        self.logger.info(f"Multi-tab: Creating additional slides for {len(other_images)} other images in {folder_name}")
        
        # Determine images per slide based on primary images count
        primary_images = teaser_images + mainunit_images
        images_per_slide = len(primary_images) if primary_images else 2  # Default to 2 if no primary images
        
        # Create slides with grouped images
        for i in range(0, len(other_images), images_per_slide):
            slide_images = other_images[i:i + images_per_slide]
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
            title_bg.line.fill.background()
            title_bg.shadow.inherit = False  # Remove shadow
            
            # Add title text with sequence number
            formatted_name = self._format_folder_name(folder_name)
            slide_number = (i // images_per_slide) + 2  # Start from 2 since first slide is primary
            title_text = f"{formatted_name} ({slide_number})"
            
            # Add title text box with exact specifications
            title_textbox = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),  # Convert cm to inches
                Inches(0.38 / 2.54),  # Convert cm to inches
                Inches(12 / 2.54),    # Convert cm to inches
                Inches(1 / 2.54)      # Convert cm to inches
            )
            title_textbox.text_frame.text = title_text
            
            # Format the title text
            title_paragraph = title_textbox.text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Use same positioning as primary slide for this group of images
            self._arrange_images_on_slide(slide, slide_images, annotation_option, folder_name)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(slide)
    
    
    
    def _add_vpm_slide(self, prs, folder_structure, annotation_option='with_annos'):
        """Add VPM slide only."""
        # Look for vdxdesktopexpandable/engaged folder
        engaged_folder = None
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopexpandable' in folder_name.lower() and 'engaged' in folder_name.lower():
                engaged_folder = (folder_name, image_paths)
                break
        
        if not engaged_folder:
            self.logger.warning("No vdxdesktopexpandable/engaged folder found for VPM slide")
            return
        
        folder_name, image_paths = engaged_folder
        
        # Find VPM.png image
        vpm_image = None
        for img_path in image_paths:
            filename = os.path.basename(img_path).lower()
            if 'vpm.png' in filename:
                vpm_image = img_path
                break
        
        # Add VPM slide
        if vpm_image:
            self._add_special_slide(prs, vpm_image, "DESKTOP EXPANDABLE - VPM (VIDEO PLAYER MODE)", annotation_option)
    
    def _add_desktop_engaged_slide(self, prs, folder_structure, annotation_option='with_annos', implement_video_frames=False):
        """Add desktop engaged slide with all images from engaged folder."""
        # Look for vdxdesktopexpandable/engaged folder
        engaged_folder = None
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopexpandable' in folder_name.lower() and 'engaged' in folder_name.lower():
                engaged_folder = (folder_name, image_paths)
                break
        
        if not engaged_folder:
            self.logger.warning("No vdxdesktopexpandable/engaged folder found for engaged slide")
            return
        
        folder_name, image_paths = engaged_folder
        
        # Filter and sort images using engaged priority
        filtered_images = [
            img_path for img_path in image_paths 
            if os.path.basename(img_path).lower() != 'mainunit-disclaimer.png'
        ]
        
        if not filtered_images:
            self.logger.warning("No valid images found in engaged folder after filtering")
            return
        
        # Sort with engaged priority (vmp.png first, mainunit.png second, then others)
        filtered_images.sort(key=self._sort_images_engaged_priority)
        
        self.logger.info(f"Processing {len(filtered_images)} engaged images in order: {[os.path.basename(img) for img in filtered_images]}")
        
        # Find MainUnit.png for video frames (if needed)
        mainunit_image = None
        for img_path in filtered_images:
            if 'mainunit.png' in os.path.basename(img_path).lower():
                mainunit_image = img_path
                break
        
        if implement_video_frames and mainunit_image:
            # Create video frame slides using mainunit image
            self.logger.info("Video frames enabled - creating video frame slides instead of original Desktop Expandable - Engaged slide")
            engaged_slide_index = len(prs.slides)  # Use current slide count as reference
            self._implement_video_frames_for_desktop_engaged(prs, folder_structure, engaged_slide_index, mainunit_image)
        else:
            # Create regular slides with all engaged images
            self._add_slide_with_images_enhanced(prs, folder_name, filtered_images, annotation_option)

    def _add_desktop_instream_slide(self, prs, folder_structure, annotation_option='with_annos', implement_video_frames=False):
        """Add Desktop Instream slide with images from vdxdesktopinstream folder."""
        # Look for vdxdesktopinstream folder and subfolders
        instream_images = []
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopinstream' in folder_name.lower():
                # Filter out disclaimer files
                filtered_images = [
                    img_path for img_path in image_paths 
                    if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                           os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
                ]
                instream_images.extend(filtered_images)
        
        if not instream_images:
            self.logger.info("No vdxdesktopinstream images found - skipping Desktop Instream slide")
            return
        
        # If video frames are enabled, skip creating the original slide and create video frames directly
        if implement_video_frames:
            self.logger.info("Video frames enabled - creating video frame slides instead of original Desktop Instream slide")
            instream_slide_index = len(prs.slides)  # Use current slide count as reference
            self._implement_video_frames_for_desktop_instream(prs, folder_structure, instream_slide_index)
            return
        
        # Create slide (only when video frames are not enabled)
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        title_text = "DESKTOP INSTREAM"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        instream_images.sort(key=self._sort_images_exact_priority)
        
        # Arrange images on the slide using standard grid layout
        self._arrange_images_on_slide(slide, instream_images, annotation_option, 'desktop_instream')
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)
    
    def _implement_video_frames_for_desktop_instream(self, prs, folder_structure, instream_slide_index):
        """
        Implement video frames with 2x1 grid layout:
        1. Create slides = ceil(video_images_count / 2) because 2 images per slide in 2x1 grid
        2. Each slide has 2x1 grid (2 copies of first Desktop Instream image side by side)
        3. Each video image overlays on corresponding Desktop Instream image copy
        4. Video image dimensions: 7.96cm width × 4.48cm height with 1.28cm Y offset
        """
        # Find video folder images
        video_images = []
        for folder_name, image_paths in folder_structure.items():
            if 'video' in folder_name.lower():
                video_images.extend(image_paths)
        
        if not video_images:
            self.logger.info("Video frames: No video images found, skipping video frames implementation")
            return 0
            
        num_video_images = len(video_images)
        import math
        
        # When video frames are enabled, we don't create the original slide first
        # Instead, use custom positioning for Desktop Instream images from underlying image parameters
        underlying_params = self.video_position_params.get('underlying_image', {}).get('desktop_instream', {})
        original_left = Inches(underlying_params.get('x_pos', 0.82) / 2.54)   # Convert cm to inches
        original_top = Inches(underlying_params.get('y_pos', 3.62) / 2.54)    # Convert cm to inches
        original_width = Inches(underlying_params.get('width', 15.34) / 2.54) # Convert cm to inches  
        original_height = Inches(underlying_params.get('height', 9.26) / 2.54) # Convert cm to inches
        images_per_slide = underlying_params.get('images_per_slide', 2)
        grid_layout = underlying_params.get('grid_layout', 'auto')
        spacing_cm = underlying_params.get('spacing', 0.5)
        
        # Determine actual grid layout to use
        if grid_layout == 'auto':
            if images_per_slide == 1:
                actual_grid = '1x1'
            elif images_per_slide == 2:
                actual_grid = '2x1'
            elif images_per_slide == 3:
                actual_grid = '3x1'
            elif images_per_slide == 4:
                actual_grid = '2x2'
            elif images_per_slide in [5, 6]:
                actual_grid = '3x2'
            else:
                actual_grid = '2x1'  # fallback
        else:
            actual_grid = grid_layout
        
        num_slides_needed = math.ceil(num_video_images / images_per_slide)
        
        self.logger.info(f"Video frames: Found {num_video_images} video images, creating {num_slides_needed} slides using {actual_grid} grid layout")
        
        # Get the first image file path from Desktop Instream folder
        first_image_path = None
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopinstream' in folder_name.lower():
                if image_paths:
                    first_image_path = image_paths[0]  # Get first image path
                    break
        
        if not first_image_path:
            self.logger.warning("Video frames: Could not find first image path for copying")
            return 0
        
        self.logger.info(f"Video frames: Original image positioned at ({original_left}, {original_top}) with size ({original_width}, {original_height})")
        
        # Convert positioning parameters to inches for grid calculation
        x_pos_inches = original_left / Inches(1)  # Convert Inches to float
        y_pos_inches = original_top / Inches(1)   # Convert Inches to float
        width_inches = original_width / Inches(1) # Convert Inches to float
        height_inches = original_height / Inches(1) # Convert Inches to float
        spacing_inches = spacing_cm / 2.54
        
        # Video overlay specs - use custom parameters if provided
        desktop_instream_params = self.video_position_params.get('video_position', {}).get('desktop_instream', {})
        video_width = Inches(desktop_instream_params.get('width', 7.96) / 2.54)
        video_height = Inches(desktop_instream_params.get('height', 4.48) / 2.54)
        video_x_offset = Inches(desktop_instream_params.get('x_offset', 0) / 2.54)
        video_y_offset = Inches(desktop_instream_params.get('y_offset', 1.28) / 2.54)
        
        # Create slides for video frames
        slides_created = 0
        for slide_idx in range(num_slides_needed):
            # Create new slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(new_slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2
            title_bg.shadow.inherit = False
            
            # Add slide title
            title_text_box = new_slide.shapes.add_textbox(
                Inches(0.51 / 2.54), Inches(0.38 / 2.54),  # Position
                Inches(12 / 2.54), Inches(1 / 2.54)        # Size
            )
            title_frame = title_text_box.text_frame
            title_frame.margin_left = Inches(0)
            title_frame.margin_right = Inches(0)
            title_frame.margin_top = Inches(0)
            title_frame.margin_bottom = Inches(0)
            title_frame.word_wrap = False
            
            title_p = title_frame.paragraphs[0]
            if slide_idx == 0:
                title_p.text = "DESKTOP INSTREAM - VIDEO FRAME"
            else:
                title_p.text = "DESKTOP INSTREAM - VIDEO FRAME (Contd.)"
            title_p.alignment = PP_ALIGN.LEFT
            
            title_run = title_p.runs[0]
            title_run.font.name = 'Aptos Display'
            title_run.font.size = Pt(18)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add underlying images based on grid layout
            grid_positions = self._calculate_grid_positions(actual_grid, images_per_slide, 
                                                          x_pos_inches, y_pos_inches, 
                                                          width_inches, height_inches, spacing_inches)
            
            for i in range(images_per_slide):
                if i < len(grid_positions):
                    x_position, y_position = grid_positions[i]
                    
                    # Add copy of Desktop Instream image
                    underlying_picture = new_slide.shapes.add_picture(
                        first_image_path,
                        Inches(x_position),
                        Inches(y_position),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    
                    # Add border to underlying image
                    underlying_picture.line.color.rgb = RGBColor(0, 0, 0)
                    underlying_picture.line.width = Pt(0.5)
                    
                    # Calculate video image index
                    video_index = slide_idx * images_per_slide + i
                    
                    # Add video image overlay if available
                    if video_index < len(video_images):
                        video_x = x_position + (video_x_offset / Inches(1))
                        video_y = y_position + (video_y_offset / Inches(1))
                        
                        video_picture = new_slide.shapes.add_picture(
                            video_images[video_index],
                            Inches(video_x),
                            Inches(video_y),
                            video_width,
                            video_height
                        )
                        
                        # Add black border to video image
                        video_picture.line.color.rgb = RGBColor(0, 0, 0)
                        video_picture.line.width = Pt(0.5)
                        
                        self.logger.info(f"Added video image {os.path.basename(video_images[video_index])} at position ({video_x:.2f}, {video_y:.2f})")
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(new_slide)
            
            slides_created += 1
            slide_name = "DESKTOP INSTREAM - VIDEO FRAME" if slide_idx == 0 else "DESKTOP INSTREAM - VIDEO FRAME (Contd.)"
            self.logger.info(f"Video frames: Created {slide_name} slide with {actual_grid} grid and video overlays")
        
        self.logger.info(f"Video frames: Created {slides_created} Desktop Instream slides with {actual_grid} grid layout for {num_video_images} video images")
        
        # Now add slides for remaining Desktop Instream images (excluding the first image used for copying)
        remaining_instream_images = []
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopinstream' in folder_name.lower():
                if len(image_paths) > 1:  # Skip first image, get the rest
                    remaining_instream_images.extend(image_paths[1:])
        
        if remaining_instream_images:
            self.logger.info(f"Video frames: Found {len(remaining_instream_images)} additional Desktop Instream images to process")
            
            # Create slides for remaining images using same formatting as video frame slides
            additional_slides = self._create_additional_desktop_instream_slides(
                prs, remaining_instream_images, original_left, original_top, 
                original_width, original_height, slides_created
            )
            slides_created += additional_slides
        
        return slides_created

    def _implement_video_frames_for_desktop_inframe_970x250(self, prs, folder_structure, video_images, inframe_970x250_images):
        """Implement video frames for Desktop In-frame 970x250 slides using custom parameters."""
        if not video_images:
            self.logger.info("No video images found for Desktop In-frame 970x250 video frames")
            return 0
        
        if not inframe_970x250_images:
            self.logger.info("No Desktop In-frame 970x250 images found for video frames")
            return 0
        
        # Get the first inframe 970x250 image to copy
        first_inframe_image = inframe_970x250_images[0]
        
        # Get custom parameters for underlying images and video positioning
        self.logger.info(f"DEBUG: video_position_params exists: {hasattr(self, 'video_position_params')}")
        self.logger.info(f"DEBUG: video_position_params content: {getattr(self, 'video_position_params', None)}")
        
        if hasattr(self, 'video_position_params') and self.video_position_params and 'underlying_image' in self.video_position_params and 'desktop_inframe_970x250' in self.video_position_params['underlying_image']:
            underlying_params = self.video_position_params['underlying_image']['desktop_inframe_970x250']
            width_cm = underlying_params['width']
            height_cm = underlying_params['height']
            x_pos_cm = underlying_params['x_pos']
            y_pos_cm = underlying_params['y_pos']
            images_per_slide = underlying_params['images_per_slide']
            spacing_cm = underlying_params['spacing']
            grid_layout = underlying_params.get('grid_layout', 'auto')
            self.logger.info(f"DEBUG: Using CUSTOM parameters - x_pos: {x_pos_cm}, y_pos: {y_pos_cm}, width: {width_cm}, height: {height_cm}")
        else:
            # Default parameters for Desktop In-frame 970x250
            width_cm = 27.59
            height_cm = 7.12
            x_pos_cm = 1
            y_pos_cm = 2.41
            images_per_slide = 2
            spacing_cm = 0.5
            grid_layout = 'auto'
            self.logger.info(f"DEBUG: Using DEFAULT parameters - x_pos: {x_pos_cm}, y_pos: {y_pos_cm}, width: {width_cm}, height: {height_cm}")
        
        # Get custom video parameters
        if hasattr(self, 'video_position_params') and self.video_position_params and 'video_position' in self.video_position_params and 'desktop_inframe_970x250' in self.video_position_params['video_position']:
            video_params = self.video_position_params['video_position']['desktop_inframe_970x250']
            video_width_cm = video_params['width']
            video_height_cm = video_params['height']
            video_x_offset_cm = video_params['x_offset']
            video_y_offset_cm = video_params['y_offset']
        else:
            # Default video parameters
            video_width_cm = 7.96
            video_height_cm = 4.48
            video_x_offset_cm = 0
            video_y_offset_cm = 1.28
        
        # Convert to inches
        width_inches = width_cm / 2.54
        height_inches = height_cm / 2.54
        x_pos_inches = x_pos_cm / 2.54
        y_pos_inches = y_pos_cm / 2.54
        spacing_inches = spacing_cm / 2.54
        
        video_width_inches = video_width_cm / 2.54
        video_height_inches = video_height_cm / 2.54
        video_x_offset_inches = video_x_offset_cm / 2.54
        video_y_offset_inches = video_y_offset_cm / 2.54
        
        # Calculate number of slides needed for video frames
        num_video_slides = math.ceil(len(video_images) / images_per_slide)
        
        # Determine actual grid layout to use
        if grid_layout == 'auto':
            if images_per_slide == 1:
                actual_grid = '1x1'
            elif images_per_slide == 2:
                actual_grid = '2x1'
            elif images_per_slide == 3:
                actual_grid = '3x1'
            elif images_per_slide == 4:
                actual_grid = '2x2'
            elif images_per_slide in [5, 6]:
                actual_grid = '3x2'
            else:
                actual_grid = '2x1'  # fallback
        else:
            actual_grid = grid_layout
        
        self.logger.info(f"Creating {num_video_slides} Desktop In-frame 970x250 video frame slides for {len(video_images)} video images")
        self.logger.info(f"Using underlying image dimensions: {width_cm}x{height_cm}cm at ({x_pos_cm}, {y_pos_cm})cm")
        self.logger.info(f"Using video dimensions: {video_width_cm}x{video_height_cm}cm with offset ({video_x_offset_cm}, {video_y_offset_cm})cm")
        self.logger.info(f"Using grid layout: {actual_grid} (setting: {grid_layout}) for {images_per_slide} images per slide")
        
        slides_created = 0
        
        # Create video frame slides
        for slide_num in range(num_video_slides):
            # Create new slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes
            self._remove_placeholders(new_slide)
            
            # Add gray rectangle background for title
            title_bg = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)
            title_bg.shadow.inherit = False
            
            # Add slide title
            title_text_box = new_slide.shapes.add_textbox(
                Inches(0.51 / 2.54), Inches(0.38 / 2.54),
                Inches(12 / 2.54), Inches(1 / 2.54)
            )
            title_frame = title_text_box.text_frame
            title_frame.margin_left = Inches(0)
            title_frame.margin_right = Inches(0)
            title_frame.margin_top = Inches(0)
            title_frame.margin_bottom = Inches(0)
            title_frame.word_wrap = False
            
            title_p = title_frame.paragraphs[0]
            if num_video_slides > 1:
                title_p.text = f"DESKTOP IN-FRAME - 970x250 VIDEO FRAME {slide_num + 1}"
            else:
                title_p.text = "DESKTOP IN-FRAME - 970x250 VIDEO FRAME"
            title_p.alignment = PP_ALIGN.LEFT
            
            title_run = title_p.runs[0]
            title_run.font.name = 'Aptos Display'
            title_run.font.size = Pt(18)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add underlying images based on grid layout
            grid_positions = self._calculate_grid_positions(actual_grid, images_per_slide, 
                                                          x_pos_inches, y_pos_inches, 
                                                          width_inches, height_inches, spacing_inches)
            
            for i in range(images_per_slide):
                if i < len(grid_positions):
                    x_position, y_position = grid_positions[i]
                    
                    # Add copy of inframe image
                    underlying_picture = new_slide.shapes.add_picture(
                        first_inframe_image,
                        Inches(x_position),
                        Inches(y_position),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    
                    # Add border to underlying image
                    underlying_picture.line.color.rgb = RGBColor(0, 0, 0)
                    underlying_picture.line.width = Pt(0.5)
                    
                    # Calculate video image index
                    video_index = slide_num * images_per_slide + i
                    
                    # Add video image overlay if available
                    if video_index < len(video_images):
                        video_x = x_position + video_x_offset_inches
                        video_y = y_position + video_y_offset_inches
                        
                        video_picture = new_slide.shapes.add_picture(
                            video_images[video_index],
                            Inches(video_x),
                            Inches(video_y),
                            Inches(video_width_inches),
                            Inches(video_height_inches)
                        )
                        
                        # Video images have no borders per user specification
                        self.logger.info(f"Added video image {os.path.basename(video_images[video_index])} at position ({video_x:.2f}, {video_y:.2f})")
            
            # Add VDX TV logo
            self._add_vdx_logo(new_slide)
            
            slides_created += 1
            self.logger.info(f"Created Desktop In-frame 970x250 video frame slide {slide_num + 1}")
        
        self.logger.info(f"Desktop In-frame 970x250 video frames: Created {slides_created} slides")
        return slides_created

    def _calculate_grid_positions(self, grid_layout, images_per_slide, start_x, start_y, image_width, image_height, spacing):
        """Calculate positions for images in various grid layouts."""
        positions = []
        
        # Parse grid layout (e.g., "2x1" means 2 columns, 1 row)
        cols, rows = map(int, grid_layout.split('x'))
        
        # Calculate available space considering margins
        slide_width_inches = 13.33  # Standard slide width minus margins
        slide_height_inches = 7.5   # Standard slide height minus margins
        
        # Center the grid on the slide
        total_grid_width = cols * image_width + (cols - 1) * spacing
        total_grid_height = rows * image_height + (rows - 1) * spacing
        
        # Calculate starting position to center the grid
        grid_start_x = max(start_x, (slide_width_inches - total_grid_width) / 2)
        grid_start_y = max(start_y, (slide_height_inches - total_grid_height) / 2 + 1)  # +1 for title space
        
        # Generate positions
        for i in range(min(images_per_slide, cols * rows)):
            row = i // cols
            col = i % cols
            
            x = grid_start_x + col * (image_width + spacing)
            y = grid_start_y + row * (image_height + spacing)
            
            positions.append((x, y))
            
        self.logger.info(f"Grid {grid_layout}: Generated {len(positions)} positions starting at ({grid_start_x:.2f}, {grid_start_y:.2f})")
        return positions

    def _implement_video_frames_for_desktop_engaged(self, prs, folder_structure, engaged_slide_index, mainunit_image_path):
        """
        Implement video frames with 2x1 grid layout for Desktop Expandable - Engaged slide:
        1. Create slides = ceil(video_images_count / 2) because 2 images per slide in 2x1 grid
        2. Each slide has 2x1 grid (2 copies of Desktop Expandable - Engaged image side by side)
        3. Each video image overlays on corresponding Desktop Expandable - Engaged image copy
        4. Video image dimensions: 7.96cm width × 4.48cm height with 1.28cm Y offset
        5. Video images have NO borders (per user requirement)
        """
        # Look for video folder
        video_images = []
        for folder_name, image_paths in folder_structure.items():
            if 'video' in folder_name.lower():
                video_images.extend(image_paths)
        
        if not video_images:
            self.logger.info("No video folder found - skipping video frames for Desktop Expandable - Engaged")
            return 0
        
        num_video_images = len(video_images)
        self.logger.info(f"Video frames: Found {num_video_images} video images for Desktop Expandable - Engaged")
        
        # When video frames are enabled, we don't create the original slide first
        # Instead, use custom positioning for Desktop Expandable images from underlying image parameters
        underlying_params = self.video_position_params.get('underlying_image', {}).get('desktop_expandable', {})
        original_left = Inches(underlying_params.get('x_pos', 0.82) / 2.54)   # Convert cm to inches
        original_top = Inches(underlying_params.get('y_pos', 3.62) / 2.54)    # Convert cm to inches
        original_width = Inches(underlying_params.get('width', 15.34) / 2.54) # Convert cm to inches  
        original_height = Inches(underlying_params.get('height', 9.26) / 2.54) # Convert cm to inches
        spacing_cm = underlying_params.get('spacing', 0.5)
        images_per_slide = underlying_params.get('images_per_slide', 2)
        grid_layout = underlying_params.get('grid_layout', 'auto')
        
        self.logger.info(f"Video frames: Using Desktop Expandable - Engaged image positioning - left: {original_left}, top: {original_top}, width: {original_width}, height: {original_height}")
        
        # Use the mainunit image path that was passed to the function
        first_image_path = mainunit_image_path
        
        if not first_image_path:
            self.logger.warning("No mainunit image provided for Desktop Expandable - Engaged video frames")
            return 0
        
        self.logger.info(f"Using engaged image for video frames: {os.path.basename(first_image_path)}")
        
        # Determine actual grid layout to use
        if grid_layout == 'auto':
            if images_per_slide == 1:
                actual_grid = '1x1'
            elif images_per_slide == 2:
                actual_grid = '2x1'
            elif images_per_slide == 3:
                actual_grid = '3x1'
            elif images_per_slide == 4:
                actual_grid = '2x2'
            elif images_per_slide in [5, 6]:
                actual_grid = '3x2'
            else:
                actual_grid = '2x1'  # fallback
        else:
            actual_grid = grid_layout
        
        # Calculate number of video frame slides needed based on images per slide
        num_video_slides = math.ceil(num_video_images / images_per_slide)
        self.logger.info(f"Video frames: Creating {num_video_slides} slides for {num_video_images} video images using {actual_grid} grid layout")
        
        # Convert positioning parameters to inches
        x_pos_inches = original_left / Inches(1)  # Convert Inches to float
        y_pos_inches = original_top / Inches(1)   # Convert Inches to float
        width_inches = original_width / Inches(1) # Convert Inches to float
        height_inches = original_height / Inches(1) # Convert Inches to float
        spacing_inches = spacing_cm / 2.54
        
        # Video overlay dimensions and positioning - use custom parameters if provided
        desktop_expandable_params = self.video_position_params.get('video_position', {}).get('desktop_expandable', {})
        video_width = Inches(desktop_expandable_params.get('width', 7.96) / 2.54)
        video_height = Inches(desktop_expandable_params.get('height', 4.48) / 2.54)
        video_x_offset = Inches(desktop_expandable_params.get('x_offset', 0) / 2.54)
        video_y_offset = Inches(desktop_expandable_params.get('y_offset', 1.28) / 2.54)
        
        copy_width = original_width
        copy_height = original_height
        
        slides_created = 0
        
        for slide_idx in range(num_video_slides):
            # Create new slide with same layout
            slide_layout = prs.slide_layouts[5]  # Blank layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(new_slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2
            title_bg.shadow.inherit = False
            
            # Add slide title
            title_text_box = new_slide.shapes.add_textbox(
                Inches(0.51 / 2.54), Inches(0.38 / 2.54),  # Position
                Inches(12 / 2.54), Inches(1 / 2.54)        # Size
            )
            title_frame = title_text_box.text_frame
            title_frame.margin_left = Inches(0)
            title_frame.margin_right = Inches(0)
            title_frame.margin_top = Inches(0)
            title_frame.margin_bottom = Inches(0)
            title_frame.word_wrap = False
            
            title_p = title_frame.paragraphs[0]
            if slide_idx == 0:
                title_p.text = "DESKTOP EXPANDABLE - ENGAGED VIDEO FRAME"
            else:
                title_p.text = "DESKTOP EXPANDABLE - ENGAGED VIDEO FRAME (Contd.)"
            title_p.alignment = PP_ALIGN.LEFT
            
            title_run = title_p.runs[0]
            title_run.font.name = 'Aptos Display'
            title_run.font.size = Pt(18)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add underlying images based on grid layout
            grid_positions = self._calculate_grid_positions(actual_grid, images_per_slide, 
                                                          x_pos_inches, y_pos_inches, 
                                                          width_inches, height_inches, spacing_inches)
            
            for i in range(images_per_slide):
                if i < len(grid_positions):
                    x_position, y_position = grid_positions[i]
                    
                    # Add copy of engaged image
                    underlying_picture = new_slide.shapes.add_picture(
                        first_image_path,
                        Inches(x_position),
                        Inches(y_position),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    
                    # Add border to underlying image
                    underlying_picture.line.color.rgb = RGBColor(0, 0, 0)
                    underlying_picture.line.width = Pt(0.5)
                    
                    # Calculate video image index
                    video_index = slide_idx * images_per_slide + i
                    
                    # Add video image overlay if available
                    if video_index < len(video_images):
                        video_x = x_position + (video_x_offset / Inches(1))
                        video_y = y_position + (video_y_offset / Inches(1))
                        
                        video_picture = new_slide.shapes.add_picture(
                            video_images[video_index],
                            Inches(video_x),
                            Inches(video_y),
                            video_width,
                            video_height
                        )
                        
                        # Video images have no borders per user specification
                        self.logger.info(f"Added video image {os.path.basename(video_images[video_index])} at position ({video_x:.2f}, {video_y:.2f})")
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(new_slide)
            
            slides_created += 1
            slide_name = "DESKTOP EXPANDABLE - ENGAGED VIDEO FRAME" if slide_idx == 0 else "DESKTOP EXPANDABLE - ENGAGED VIDEO FRAME (Contd.)"
            self.logger.info(f"Video frames: Created {slide_name} slide with {actual_grid} grid and video overlays (no borders on video images)")
        
        self.logger.info(f"Video frames: Created {slides_created} Desktop Expandable - Engaged slides with {actual_grid} grid layout for {num_video_images} video images")
        
        # Now add slides for remaining Desktop Expandable - Engaged images (excluding MainUnit.png used for copying)
        remaining_engaged_images = []
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopexpandable' in folder_name.lower() and 'engaged' in folder_name.lower():
                # Skip MainUnit.png, get the rest
                remaining_images = [img for img in image_paths if 'mainunit.png' not in os.path.basename(img).lower()]
                remaining_engaged_images.extend(remaining_images)
        
        if remaining_engaged_images:
            self.logger.info(f"Video frames: Found {len(remaining_engaged_images)} additional Desktop Expandable - Engaged images to process")
            
            # Create slides for remaining images using same formatting as video frame slides
            additional_slides = self._create_additional_desktop_engaged_slides(
                prs, remaining_engaged_images, original_left, original_top, 
                original_width, original_height, slides_created
            )
            slides_created += additional_slides
        
        return slides_created

    def _create_additional_desktop_engaged_slides(self, prs, remaining_images, original_left, original_top, original_width, original_height, start_slide_num):
        """Create additional slides for remaining Desktop Expandable - Engaged images using same formatting as video frame slides."""
        slides_created = 0
        images_per_slide = 2  # Same as video frame slides (2x1 grid)
        
        # Calculate spacing between images (same as video frame slides)
        spacing = Inches(0.5)
        left1_x = original_left
        left1_y = original_top
        left2_x = original_left + original_width + spacing
        left2_y = original_top
        
        # Process images in groups of 2
        for i in range(0, len(remaining_images), images_per_slide):
            slide_images = remaining_images[i:i + images_per_slide]
            slide_number = start_slide_num + slides_created + 1
            
            # Create new slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(new_slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2
            title_bg.shadow.inherit = False
            
            # Add slide title
            title_text_box = new_slide.shapes.add_textbox(
                Inches(0.51 / 2.54), Inches(0.38 / 2.54),  # Position
                Inches(12 / 2.54), Inches(1 / 2.54)        # Size
            )
            title_frame = title_text_box.text_frame
            title_frame.margin_left = Inches(0)
            title_frame.margin_right = Inches(0)
            title_frame.margin_top = Inches(0)
            title_frame.margin_bottom = Inches(0)
            title_frame.word_wrap = False
            
            title_p = title_frame.paragraphs[0]
            title_p.text = "DESKTOP EXPANDABLE - ENGAGED (Contd.)"
            title_p.alignment = PP_ALIGN.LEFT
            
            title_run = title_p.runs[0]
            title_run.font.name = 'Aptos Display'
            title_run.font.size = Pt(18)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add images using same positioning as video frame slides
            # First image (left position)
            if len(slide_images) > 0:
                picture1 = new_slide.shapes.add_picture(
                    slide_images[0],
                    left1_x, left1_y,  # Left position
                    original_width, original_height  # Same size as original
                )
                
                # Add black border
                picture1.line.color.rgb = RGBColor(0, 0, 0)
                picture1.line.width = Pt(0.5)
            
            # Second image (right position)
            if len(slide_images) > 1:
                picture2 = new_slide.shapes.add_picture(
                    slide_images[1],
                    left2_x, left2_y,  # Right position
                    original_width, original_height  # Same size as original
                )
                
                # Add black border
                picture2.line.color.rgb = RGBColor(0, 0, 0)
                picture2.line.width = Pt(0.5)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(new_slide)
            
            slides_created += 1
            self.logger.info(f"Video frames: Created additional Desktop Expandable - Engaged (Contd.) slide with {len(slide_images)} images")
        
        self.logger.info(f"Video frames: Created {slides_created} additional Desktop Expandable - Engaged slides for remaining images")
        return slides_created

    def _create_additional_desktop_instream_slides(self, prs, remaining_images, original_left, original_top, original_width, original_height, start_slide_num):
        """Create additional slides for remaining Desktop Instream images using same formatting as video frame slides."""
        slides_created = 0
        images_per_slide = 2  # Same as video frame slides (2x1 grid)
        
        # Calculate spacing between images (same as video frame slides)
        spacing = Inches(0.5)
        left1_x = original_left
        left1_y = original_top
        left2_x = original_left + original_width + spacing
        left2_y = original_top
        
        # Process images in groups of 2
        for i in range(0, len(remaining_images), images_per_slide):
            slide_images = remaining_images[i:i + images_per_slide]
            slide_number = start_slide_num + slides_created + 1
            
            # Create new slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(new_slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = new_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2
            title_bg.shadow.inherit = False
            
            # Add slide title
            title_text_box = new_slide.shapes.add_textbox(
                Inches(0.51 / 2.54), Inches(0.38 / 2.54),  # Position
                Inches(12 / 2.54), Inches(1 / 2.54)        # Size
            )
            title_frame = title_text_box.text_frame
            title_frame.margin_left = Inches(0)
            title_frame.margin_right = Inches(0)
            title_frame.margin_top = Inches(0)
            title_frame.margin_bottom = Inches(0)
            title_frame.word_wrap = False
            
            title_p = title_frame.paragraphs[0]
            title_p.text = f"DESKTOP INSTREAM - ADDITIONAL {slide_number}"
            title_p.alignment = PP_ALIGN.LEFT
            
            title_run = title_p.runs[0]
            title_run.font.name = 'Aptos Display'
            title_run.font.size = Pt(18)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add images using same positioning as video frame slides
            # First image (left position)
            if len(slide_images) > 0:
                picture1 = new_slide.shapes.add_picture(
                    slide_images[0],
                    left1_x, left1_y,  # Left position
                    original_width, original_height  # Same size as original
                )
                
                # Add black border
                picture1.line.color.rgb = RGBColor(0, 0, 0)
                picture1.line.width = Pt(0.5)
            
            # Second image (right position)
            if len(slide_images) > 1:
                picture2 = new_slide.shapes.add_picture(
                    slide_images[1],
                    left2_x, left2_y,  # Right position
                    original_width, original_height  # Same size as original
                )
                
                # Add black border
                picture2.line.color.rgb = RGBColor(0, 0, 0)
                picture2.line.width = Pt(0.5)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(new_slide)
            
            slides_created += 1
            self.logger.info(f"Video frames: Created additional Desktop Instream slide {slide_number} with {len(slide_images)} images")
        
        self.logger.info(f"Video frames: Created {slides_created} additional Desktop Instream slides for remaining images")
        return slides_created

    def _add_mobile_instream_slide(self, prs, folder_structure, annotation_option='with_annos'):
        """Add Mobile Instream slide with images from vdxmobileinstream folder."""
        # Look for vdxmobileinstream folder and subfolders
        instream_images = []
        for folder_name, image_paths in folder_structure.items():
            if 'vdxmobileinstream' in folder_name.lower():
                # Filter out disclaimer files
                filtered_images = [
                    img_path for img_path in image_paths 
                    if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                           os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
                ]
                instream_images.extend(filtered_images)
        
        if not instream_images:
            self.logger.info("No vdxmobileinstream images found - skipping Mobile Instream slide")
            return
        
        # Create slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        title_text = "MOBILE INSTREAM"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        instream_images.sort(key=self._sort_images_exact_priority)
        
        # Use mobile instream positioning (same as Mobile In-frame 300x600)
        self._arrange_mobile_instream_images(slide, instream_images, annotation_option)
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)

    def _add_ott_slide(self, prs, folder_structure, annotation_option='with_annos'):
        """Add OTT slide with images from ott folder."""
        # Look for ott folder and subfolders - check for various patterns
        ott_images = []
        for folder_name, image_paths in folder_structure.items():
            folder_lower = folder_name.lower()
            # Check for direct ott folders or nested ott folders
            if (folder_lower == 'ott' or 
                folder_lower.startswith('ott/') or 
                folder_lower.endswith('/ott') or 
                '/ott/' in folder_lower or
                folder_lower.split('/')[-1] == 'ott'):
                # Filter out disclaimer files
                filtered_images = [
                    img_path for img_path in image_paths 
                    if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                           os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
                ]
                ott_images.extend(filtered_images)
        
        if not ott_images:
            self.logger.info("No ott images found - creating blank OTT slide")
            # Create blank slide anyway
            pass
        
        # Create slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        title_text = "OTT"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # No specific sorting needed - just arrange images as they are
        # Arrange images on the slide using standard grid layout (no borders for OTT)
        if ott_images:
            self._arrange_images_on_slide(slide, ott_images, annotation_option, 'ott_no_borders')
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)

    def _add_ctv_slide(self, prs, folder_structure, annotation_option='with_annos'):
        """Add CTV slide with images from ctv folder."""
        # Look for ctv folder and subfolders - check for various patterns
        ctv_images = []
        for folder_name, image_paths in folder_structure.items():
            folder_lower = folder_name.lower()
            # Check for direct ctv folders or nested ctv folders
            if (folder_lower == 'ctv' or 
                folder_lower.startswith('ctv/') or 
                folder_lower.endswith('/ctv') or 
                '/ctv/' in folder_lower or
                folder_lower.split('/')[-1] == 'ctv'):
                # Filter out disclaimer files
                filtered_images = [
                    img_path for img_path in image_paths 
                    if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                           os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
                ]
                ctv_images.extend(filtered_images)
        
        if not ctv_images:
            self.logger.info("No ctv images found - creating blank CTV slide")
            # Create blank slide anyway
            pass
        
        # Create slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        title_text = "CTV"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # No specific sorting needed - just arrange images as they are
        # Arrange images on the slide using standard grid layout (no borders for CTV)
        if ctv_images:
            self._arrange_images_on_slide(slide, ctv_images, annotation_option, 'ctv_no_borders')
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)

    def _add_vpm_and_engaged_slides(self, prs, folder_structure, annotation_option='with_annos', implement_video_frames=False):
        """Add VPM and engaged slides after Desktop In-frame 970x250 slide."""
        # Look for vdxdesktopexpandable/engaged folder
        engaged_folder = None
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopexpandable' in folder_name.lower() and 'engaged' in folder_name.lower():
                engaged_folder = (folder_name, image_paths)
                break
        
        if not engaged_folder:
            self.logger.warning("No vdxdesktopexpandable/engaged folder found for VPM and engaged slides")
            return
        
        folder_name, image_paths = engaged_folder
        
        # Filter out mainunit-disclaimer.png files
        filtered_images = [
            img for img in image_paths 
            if os.path.basename(img).lower() != 'mainunit-disclaimer.png'
        ]
        
        # Sort images to ensure consistent ordering
        sorted_images = sorted(filtered_images)
        
        # Find vmp.png and mainunit.png images using exact matching
        vmp_image = None
        mainunit_image = None
        other_images = []
        
        for img_path in sorted_images:
            filename = os.path.basename(img_path).lower()
            if filename == 'vpm.png':
                vmp_image = img_path
            elif filename == 'mainunit.png':
                mainunit_image = img_path
            else:
                other_images.append(img_path)
        
        self.logger.info(f"Found engaged folder with {len(sorted_images)} images (excluding disclaimer)")
        self.logger.info(f"VMP image: {vmp_image}")
        self.logger.info(f"MainUnit image: {mainunit_image}")
        
        # Create combined slide with both vmp.png and mainunit.png
        if vmp_image and mainunit_image:
            # If video frames are enabled, handle video frames logic
            if implement_video_frames:
                self.logger.info("Video frames enabled - creating video frame slides instead of original Desktop Expandable - Engaged slide")
                engaged_slide_index = len(prs.slides)  # Use current slide count as reference
                self._implement_video_frames_for_desktop_engaged(prs, folder_structure, engaged_slide_index, mainunit_image)
            else:
                # Create combined slide with both images
                self._add_combined_desktop_expandable_slide(prs, vmp_image, mainunit_image, annotation_option)
        elif vmp_image or mainunit_image:
            # If only one image exists, create single image slide
            image_to_use = vmp_image if vmp_image else mainunit_image
            if implement_video_frames and mainunit_image:
                self.logger.info("Video frames enabled - creating video frame slides instead of original Desktop Expandable - Engaged slide")
                engaged_slide_index = len(prs.slides)
                self._implement_video_frames_for_desktop_engaged(prs, folder_structure, engaged_slide_index, mainunit_image)
            else:
                self._add_special_slide(prs, image_to_use, "DESKTOP EXPANDABLE - ENGAGED", annotation_option)
        
        # Handle additional images if any exist
        if other_images and not implement_video_frames:
            self._add_additional_desktop_expandable_slides(prs, other_images, annotation_option)
    
    def _add_combined_desktop_expandable_slide(self, prs, vmp_image_path, mainunit_image_path, annotation_option='with_annos'):
        """Add combined slide with vmp.png and mainunit.png side by side with 0.5cm spacing."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = "DESKTOP EXPANDABLE - ENGAGED"
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Image dimensions (same as Desktop Expandable: width 15.34cm, height 8.64cm)
        image_width = 15.34 / 2.54  # Convert cm to inches
        image_height = 8.64 / 2.54  # Convert cm to inches
        image_y = 3.62 / 2.54  # Convert cm to inches (Y position)
        
        # Calculate X positions for two images with 0.5cm spacing
        spacing_cm = 0.5
        spacing_inches = spacing_cm / 2.54
        
        # First image (vmp.png) - left position
        first_image_x = 0.82 / 2.54  # Convert cm to inches (same X as single image)
        
        # Second image (mainunit.png) - right position with 0.5cm gap
        second_image_x = first_image_x + image_width + spacing_inches
        
        try:
            # Add first image (vmp.png)
            picture1_shape = slide.shapes.add_picture(
                vmp_image_path,
                Inches(first_image_x),
                Inches(image_y),
                Inches(image_width),
                Inches(image_height)
            )
            
            # Add black border (0.5pt)
            picture1_shape.line.color.rgb = RGBColor(0, 0, 0)
            picture1_shape.line.width = Pt(0.5)
            
            
            
            self.logger.info(f"Added VMP image {os.path.basename(vmp_image_path)} at position ({first_image_x:.2f}, {image_y:.2f})")
            
        except Exception as e:
            self.logger.error(f"Error adding VMP image {vmp_image_path}: {str(e)}")
        
        try:
            # Add second image (mainunit.png)
            picture2_shape = slide.shapes.add_picture(
                mainunit_image_path,
                Inches(second_image_x),
                Inches(image_y),
                Inches(image_width),
                Inches(image_height)
            )
            
            # Add black border (0.5pt)
            picture2_shape.line.color.rgb = RGBColor(0, 0, 0)
            picture2_shape.line.width = Pt(0.5)
            
            
            
            self.logger.info(f"Added MainUnit image {os.path.basename(mainunit_image_path)} at position ({second_image_x:.2f}, {image_y:.2f})")
            
        except Exception as e:
            self.logger.error(f"Error adding MainUnit image {mainunit_image_path}: {str(e)}")
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)
        
        self.logger.info("Created combined DESKTOP EXPANDABLE - ENGAGED slide with vmp.png and mainunit.png")
    
    def _add_additional_desktop_expandable_slides(self, prs, additional_images, annotation_option='with_annos'):
        """Add additional slides for remaining Desktop Expandable images using same layout as combined slide."""
        if not additional_images:
            return
            
        # Process images in chunks of 2 (same as combined slide layout)
        images_per_slide = 2
        slide_num = 2  # Start with (Contd.) since first slide is the combined one
        
        for i in range(0, len(additional_images), images_per_slide):
            slide_images = additional_images[i:i + images_per_slide]
            
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
            title_bg.line.fill.background()
            title_bg.shadow.inherit = False  # Remove shadow
            
            # Add title text box with continuation marker
            title_textbox = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),  # Convert cm to inches
                Inches(0.38 / 2.54),  # Convert cm to inches
                Inches(12 / 2.54),    # Convert cm to inches
                Inches(1 / 2.54)      # Convert cm to inches
            )
            title_textbox.text_frame.text = "DESKTOP EXPANDABLE - ENGAGED (CONTD.)"
            
            # Format the title text
            title_paragraph = title_textbox.text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Image dimensions (same as combined slide)
            image_width = 15.34 / 2.54  # Convert cm to inches
            image_height = 9.26 / 2.54  # Convert cm to inches
            image_y = 3.62 / 2.54  # Convert cm to inches (Y position)
            spacing_inches = 0.5 / 2.54  # 0.5cm spacing
            
            # Position images
            for img_idx, img_path in enumerate(slide_images):
                if img_idx == 0:
                    image_x = 0.82 / 2.54  # First image position
                else:
                    image_x = (0.82 / 2.54) + image_width + spacing_inches  # Second image position
                
                try:
                    picture_shape = slide.shapes.add_picture(
                        img_path,
                        Inches(image_x),
                        Inches(image_y),
                        Inches(image_width),
                        Inches(image_height)
                    )
                    
                    # Add black border (0.5pt)
                    picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                    picture_shape.line.width = Pt(0.5)
                    
                    
                    
                    self.logger.info(f"Added additional Desktop Expandable image {os.path.basename(img_path)} at position ({image_x:.2f}, {image_y:.2f})")
                    
                except Exception as e:
                    self.logger.error(f"Error adding additional Desktop Expandable image {img_path}: {str(e)}")
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(slide)
            
            self.logger.info(f"Created additional DESKTOP EXPANDABLE - ENGAGED slide {slide_num} with {len(slide_images)} images")
            slide_num += 1
    
    def _add_special_slide(self, prs, image_path, title_text, annotation_option='with_annos'):
        """Add a special slide with single image at specified position and dimensions."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add image with same dimensions as Desktop Instream
        # Width - 15.34cm Height - 8.64cm , X = 0.82cm Y = 3.62cm (same X/Y positions)
        image_width = 15.34 / 2.54  # Convert cm to inches
        image_height = 8.64 / 2.54  # Convert cm to inches
        image_x = 0.82 / 2.54  # Convert cm to inches
        image_y = 3.62 / 2.54  # Convert cm to inches
        
        try:
            picture_shape = slide.shapes.add_picture(
                image_path,
                Inches(image_x),
                Inches(image_y),
                Inches(image_width),
                Inches(image_height)
            )
            
            # Add black border (0.5pt)
            picture_shape.line.color.rgb = RGBColor(0, 0, 0)
            picture_shape.line.width = Pt(0.5)
            
            self.logger.info(f"Added special slide image {os.path.basename(image_path)} at position ({image_x:.2f}, {image_y:.2f})")
            

                
        except Exception as e:
            self.logger.error(f"Error adding special slide image {image_path}: {str(e)}")
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)

    def _add_video_frames_slide_if_exists(self, prs, folder_structure, annotation_option='with_annos'):
        """Add video frames slides as first slides if video folder exists. Creates multiple slides for more than 6 images."""
        video_folder = None
        video_images = []
        
        # Look for folder named "video"
        for folder_name, image_paths in folder_structure.items():
            if 'video' in folder_name.lower():
                video_folder = folder_name
                video_images = image_paths
                break
        
        if not video_folder or not video_images:
            self.logger.info("No video folder detected in uploaded files")
            return False
            
        self.logger.info(f"Found video folder: {video_folder} with {len(video_images)} images")
        
        # Create multiple slides if needed (6 images per slide for 3x2 grid)
        images_per_slide = 6
        total_slides = (len(video_images) + images_per_slide - 1) // images_per_slide  # Round up division
        
        for slide_num in range(total_slides):
            # Get images for this slide
            start_idx = slide_num * images_per_slide
            end_idx = min(start_idx + images_per_slide, len(video_images))
            slide_images = video_images[start_idx:end_idx]
            
            # Create slide with "Video Frames" title
            slide_layout = prs.slide_layouts[5]  # Use blank slide layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove any placeholder shapes
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title
            rectangle = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                prs.slide_width,
                Inches(1.79 / 2.54)  # 1.79cm converted to inches
            )
            rectangle.fill.solid()
            rectangle.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2
            rectangle.line.color.rgb = RGBColor(242, 242, 242)
            rectangle.shadow.inherit = False  # Remove shadow
            
            # Add title text box
            title_left = Inches(0.51 / 2.54)  # 0.51cm converted to inches
            title_top = Inches(0.38 / 2.54)   # 0.38cm converted to inches
            title_width = Inches(12 / 2.54)   # 12cm converted to inches
            title_height = Inches(1 / 2.54)   # 1cm converted to inches
            
            title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
            title_frame = title_box.text_frame
            
            # Set title with continuation marker for additional slides
            if slide_num == 0:
                title_frame.text = "Video Frames"
            else:
                title_frame.text = "Video Frames (Contd.)"
            
            # Style the title
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Arrange images in 3x2 grid below gray rectangle
            start_frame_number = slide_num * 6 + 1  # Continuous numbering across slides
            self._arrange_video_images_3x2_grid(slide, slide_images, annotation_option, start_frame_number)
            
            # Add VDX TV logo
            self._add_vdx_logo(slide)
            
            self.logger.info(f"Added Video Frames slide {slide_num + 1} with {len(slide_images)} images in 3x2 grid")
        
        self.logger.info(f"Created {total_slides} Video Frames slides for {len(video_images)} total images")
        return True

    def _arrange_video_images_3x2_grid(self, slide, image_paths, annotation_option='with_annos', start_frame_number=1):
        """Arrange up to 6 images in a 3x2 grid layout below the gray rectangle."""
        # Take up to 6 images for 3x2 grid
        images_to_use = image_paths[:6]
        
        if not images_to_use:
            return
            
        # Grid specifications: 3 columns, 2 rows
        cols = 3
        rows = 2
        
        # Image dimensions (convert cm to inches)
        image_width_cm = 10.33
        image_height_cm = 5.81
        img_width = image_width_cm / 2.54   # inches
        img_height = image_height_cm / 2.54  # inches
        
        # Position at specific Y coordinates as requested
        first_row_y = 3.5 / 2.54  # 3.5cm converted to inches
        second_row_y = 11.30 / 2.54  # 11.30cm converted to inches
        
        # Calculate horizontal spacing to center the grid
        total_image_width = cols * img_width
        available_width = 13.33  # Total slide width
        total_spacing = available_width - total_image_width
        spacing_x = total_spacing / (cols + 1)  # Equal spacing on sides and between images
        start_x = spacing_x
        
        self.logger.info(f"Video grid: {len(images_to_use)} images, {img_width:.2f}x{img_height:.2f} inches each")
        
        # Place images in grid
        for i, img_path in enumerate(images_to_use):
            row = i // cols
            col = i % cols
            
            # Calculate position with specific Y coordinates
            x = start_x + col * (img_width + spacing_x)
            if row == 0:
                y = first_row_y  # First row at 3.5cm
            else:
                y = second_row_y  # Second row at 11.30cm
            
            try:
                # Add image to slide
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x),
                    Inches(y),
                    Inches(img_width),
                    Inches(img_height)
                )
                
                # Add black border (0.5pt)
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                # Add frame label below image (Frame-01, Frame-02, etc.) with continuous numbering
                frame_number = start_frame_number + i
                frame_label = f"Frame-{frame_number:02d}"
                
                # Position label centered below image
                label_y = y + img_height + 0.1  # 0.1 inch spacing below image
                label_width = img_width
                label_height = 0.3  # Small height for label
                
                label_textbox = slide.shapes.add_textbox(
                    Inches(x),
                    Inches(label_y),
                    Inches(label_width),
                    Inches(label_height)
                )
                
                label_frame = label_textbox.text_frame
                label_frame.text = frame_label
                label_frame.margin_left = Inches(0)
                label_frame.margin_right = Inches(0)
                label_frame.margin_top = Inches(0)
                label_frame.margin_bottom = Inches(0)
                
                # Format label text - Aptos Display, 9px, bold
                label_paragraph = label_frame.paragraphs[0]
                label_paragraph.font.name = "Aptos Display"
                label_paragraph.font.size = Pt(9)
                label_paragraph.font.bold = True
                label_paragraph.alignment = PP_ALIGN.CENTER
                label_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                
                self.logger.info(f"Added video image {os.path.basename(img_path)} at position ({x:.2f}, {y:.2f}) with label '{frame_label}'")
                
            except Exception as e:
                self.logger.error(f"Error adding video image {img_path}: {str(e)}")

    def _add_mobile_expandable_engaged_slide(self, prs, folder_structure, annotation_option='with_annos'):
        """Add Mobile Expandable engaged slide(s) with same 3-position layout as consolidated teaser slide."""
        # Look for vdxmobileexpandable/engaged folder
        engaged_folder = None
        for folder_name, image_paths in folder_structure.items():
            if 'vdxmobileexpandable' in folder_name.lower() and 'engaged' in folder_name.lower():
                engaged_folder = (folder_name, image_paths)
                break
        
        if not engaged_folder:
            self.logger.info("No vdxmobileexpandable/engaged folder found - skipping Mobile Expandable engaged slide")
            return
        
        folder_name, image_paths = engaged_folder
        
        # Filter out disclaimer files
        filtered_image_paths = [
            img_path for img_path in image_paths 
            if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                   os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
        ]
        
        if not filtered_image_paths:
            self.logger.info("No valid images found in vdxmobileexpandable/engaged folder")
            return
        
        # Process images in chunks of 3 (same as teaser slide layout)
        images_per_slide = 3
        slide_num = 1
        
        for i in range(0, len(filtered_image_paths), images_per_slide):
            slide_images = filtered_image_paths[i:i + images_per_slide]
            
            # Create slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
            title_bg.line.fill.background()
            title_bg.shadow.inherit = False  # Remove shadow
            
            # Add title text with continuation marker if needed
            if slide_num == 1:
                title_text = "MOBILE EXPANDABLE - ENGAGED"
            else:
                title_text = "MOBILE EXPANDABLE - ENGAGED (CONTD.)"
            
            # Add title text box with exact specifications
            title_textbox = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),  # Convert cm to inches
                Inches(0.38 / 2.54),  # Convert cm to inches
                Inches(12 / 2.54),    # Convert cm to inches
                Inches(1 / 2.54)      # Convert cm to inches
            )
            title_textbox.text_frame.text = title_text
            
            # Format the title text
            title_paragraph = title_textbox.text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Use same 3-position layout as Mobile Expandable - All Teasers slide
            self._arrange_mobile_engaged_images_with_custom_positions(slide, slide_images, annotation_option)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(slide)
            
            self.logger.info(f"Created Mobile Expandable - Engaged slide {slide_num} with {len(slide_images)} images")
            slide_num += 1
    
    def _arrange_mobile_engaged_images_with_custom_positions(self, slide, engaged_images, annotation_option='with_annos'):
        """Arrange mobile engaged images using same 3-position layout as teaser slide."""
        # Define positioning specifications for mobile engaged images
        # Same positions as teaser slide: 300x250, 300x600, 320x50 positions
        positioning_specs = [
            {  # Position 1 (300x250 position)
                'width_cm': 8.22,
                'height_cm': 16,
                'horizontal_cm': 2.6,
                'vertical_cm': 2.17
            },
            {  # Position 2 (300x600 position)
                'width_cm': 8.22,
                'height_cm': 16,
                'horizontal_cm': 12.11,
                'vertical_cm': 2.17
            },
            {  # Position 3 (320x50 position)
                'width_cm': 8.22,
                'height_cm': 16,
                'horizontal_cm': 21.56,
                'vertical_cm': 2.17
            }
        ]
        
        # Place images sequentially in the 3 positions
        for i, img_path in enumerate(engaged_images):
            if i >= len(positioning_specs):
                break  # Maximum 3 images per slide
                
            spec = positioning_specs[i]
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['horizontal_cm'] / 2.54
            y_inches = spec['vertical_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Mobile images don't get borders
                
                # Add filename annotation if requested
                if annotation_option == 'with_annos':
                    # Add filename text below the image
                    filename = os.path.basename(img_path)
                    text_y = y_inches + height_inches + 0.1  # 0.1 inches below image
                    filename_box = slide.shapes.add_textbox(
                        Inches(x_inches), Inches(text_y), Inches(width_inches), Inches(0.3)
                    )
                    filename_box.text_frame.text = filename
                    p = filename_box.text_frame.paragraphs[0]
                    p.font.size = Pt(8)
                    p.font.color.rgb = RGBColor(0, 0, 0)
                    p.alignment = PP_ALIGN.CENTER
                
                self.logger.info(f"Added mobile engaged image {os.path.basename(img_path)} at position {i+1} ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding mobile engaged image {img_path}: {str(e)}")

    def _add_desktop_inframe_970x250_with_additional_slides(self, prs, folder_name, image_paths, annotation_option='with_annos'):
        """Add Desktop In-frame 970x250 slides with additional slides for remaining images using same layout."""
        # Filter out disclaimer files
        filtered_image_paths = [
            img_path for img_path in image_paths 
            if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                   os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
        ]
        
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        filtered_image_paths.sort(key=self._sort_images_exact_priority)
        
        # Skip if no images left after filtering
        if not filtered_image_paths:
            return
        
        # Find exact teaser.png and mainunit.png files, plus other images
        teaser_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'teaser.png']
        mainunit_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'mainunit.png']
        other_images = [img for img in filtered_image_paths if os.path.basename(img).lower() not in ['teaser.png', 'mainunit.png']]
        
        self.logger.info(f"Auto Tab Desktop In-frame 970x250: Processing folder {folder_name}")
        self.logger.info(f"Total filtered images: {len(filtered_image_paths)}")
        self.logger.info(f"Teaser images: {len(teaser_images)} - {[os.path.basename(img) for img in teaser_images]}")
        self.logger.info(f"Mainunit images: {len(mainunit_images)} - {[os.path.basename(img) for img in mainunit_images]}")
        self.logger.info(f"Other images: {len(other_images)} - {[os.path.basename(img) for img in other_images]}")
        
        # Always create first slide with primary images (teaser + mainunit)
        primary_images = teaser_images + mainunit_images
        
        # If we have primary images, create the first slide
        if primary_images:
            self._create_desktop_inframe_970x250_slide(prs, folder_name, primary_images, annotation_option, slide_number=1)
        
        # Always check for additional images (both when we have primary images and when we don't)
        all_additional_images = other_images
        
        # If no primary images were found, treat all filtered images as additional images
        if not primary_images and filtered_image_paths:
            all_additional_images = filtered_image_paths
            self.logger.info(f"Auto Tab: No teaser/mainunit images found, treating all {len(all_additional_images)} images as additional images")
        
        # Create additional slides for remaining images
        if all_additional_images:
            images_per_slide = 2  # Always use 2 images per slide for Desktop In-frame 970x250
            
            # Determine starting slide number
            start_slide_number = 2 if primary_images else 1
            
            self.logger.info(f"Auto Tab: Creating additional Desktop In-frame 970x250 slides for {len(all_additional_images)} additional images, {images_per_slide} images per slide")
            
            # Create additional slides with remaining images using same layout
            for i in range(0, len(all_additional_images), images_per_slide):
                slide_images = all_additional_images[i:i + images_per_slide]
                slide_number = (i // images_per_slide) + start_slide_number
                self.logger.info(f"Auto Tab: Creating slide {slide_number} with images: {[os.path.basename(img) for img in slide_images]}")
                self._create_desktop_inframe_970x250_slide(prs, folder_name, slide_images, annotation_option, slide_number)

    def _create_desktop_inframe_970x250_slide(self, prs, folder_name, image_list, annotation_option, slide_number=1):
        """Create a single Desktop In-frame 970x250 slide with specific layout and positioning."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        formatted_name = self._format_folder_name(folder_name)
        if slide_number == 1:
            title_text = formatted_name
        else:
            title_text = f"{formatted_name} ({slide_number})"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Use Desktop In-frame 970x250 specific positioning
        self._arrange_desktop_inframe_970x250_images(slide, image_list, annotation_option, folder_name)
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)
        
        self.logger.info(f"Auto Tab: Created Desktop In-frame 970x250 slide {slide_number} with {len(image_list)} images")

    def _add_vdx_logo(self, slide, folder_name=None):
        """Add VDX TV logo to the slide - position varies by slide type."""
        try:
            # Logo specifications: height=0.51cm, width=1.85cm
            logo_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "static", "vdx-tv-logo.png")
            
            # Check if logo file exists
            if not os.path.exists(logo_path):
                self.logger.warning(f"VDX TV logo not found at {logo_path}")
                return
            
            # Convert cm to inches
            logo_width = 1.85 / 2.54  # 1.85cm to inches
            logo_height = 0.51 / 2.54  # 0.51cm to inches
            
            # Position VDX TV logo at exact coordinates for all slides
            logo_x = 31.42 / 2.54  # 31.42cm to inches
            logo_y = 0.63 / 2.54   # 0.63cm to inches
            
            # Add logo to slide
            logo_shape = slide.shapes.add_picture(
                logo_path,
                Inches(logo_x),
                Inches(logo_y),
                Inches(logo_width),
                Inches(logo_height)
            )
            
            self.logger.info(f"Added VDX TV logo to slide at position ({logo_x:.2f}, {logo_y:.2f})")
            
        except Exception as e:
            self.logger.error(f"Error adding VDX TV logo: {str(e)}")
    
    def _add_title_slide(self, prs, folder_count):
        """Add a title slide to the presentation."""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = "Image Collection"
        
        # Set subtitle
        subtitle = slide.placeholders[1]
        subtitle.text = f"Generated from {folder_count} folder{'s' if folder_count != 1 else ''}\n{datetime.now().strftime('%B %d, %Y')}"
        
        # Style the title
        title_paragraph = title.text_frame.paragraphs[0]
        title_paragraph.font.size = Pt(44)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.CENTER
        
        # Style the subtitle
        subtitle_paragraph = subtitle.text_frame.paragraphs[0]
        subtitle_paragraph.font.size = Pt(18)
        subtitle_paragraph.alignment = PP_ALIGN.CENTER
        
        # Add VDX TV logo to title slide
        self._add_vdx_logo(slide)
    
    def _add_slide_with_images(self, prs, folder_name, image_paths, annotation_option='with_annos'):
        """Add slides with folder name as title and images arranged on the slides."""
        # Filter out teaser-disclaimer.png and mainunit-disclaimer.png files
        filtered_image_paths = [
            img_path for img_path in image_paths 
            if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                   os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
        ]
        
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        filtered_image_paths.sort(key=self._sort_images_exact_priority)
        
        # Skip if no images left after filtering
        if not filtered_image_paths:
            return
        
        # Split images into chunks that fit on slides (max 9 per slide)
        max_images_per_slide = 9
        image_chunks = [filtered_image_paths[i:i + max_images_per_slide] for i in range(0, len(filtered_image_paths), max_images_per_slide)]
        
        for chunk_index, image_chunk in enumerate(image_chunks):
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
            title_bg.line.fill.background()
            title_bg.shadow.inherit = False  # Remove shadow
            
            # Add title text (with page number if multiple slides)
            formatted_name = self._format_folder_name(folder_name)
            title_text = formatted_name
            if len(image_chunks) > 1:
                title_text = f"{formatted_name} ({chunk_index + 1}/{len(image_chunks)})"
            
            # Add title text box with exact specifications
            title_textbox = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),  # Convert cm to inches
                Inches(0.38 / 2.54),  # Convert cm to inches
                Inches(12 / 2.54),    # Convert cm to inches
                Inches(1 / 2.54)      # Convert cm to inches
            )
            title_textbox.text_frame.text = title_text
            
            # Format the title text
            title_paragraph = title_textbox.text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Arrange images on the slide
            self._arrange_images_on_slide(slide, image_chunk, annotation_option, folder_name)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(slide)
    
    def _add_slide_with_images_enhanced(self, prs, folder_name, image_paths, annotation_option='with_annos'):
        """Add slides with folder name as title and images arranged on the slides - Enhanced version with additional slides for remaining images."""
        # Filter out teaser-disclaimer.png and mainunit-disclaimer.png files
        filtered_image_paths = [
            img_path for img_path in image_paths 
            if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                   os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
        ]
        
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        filtered_image_paths.sort(key=self._sort_images_exact_priority)
        
        # Skip if no images left after filtering
        if not filtered_image_paths:
            return
        
        # Find exact teaser.png and mainunit.png files, plus other images
        teaser_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'teaser.png']
        mainunit_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'mainunit.png']
        other_images = [img for img in filtered_image_paths if os.path.basename(img).lower() not in ['teaser.png', 'mainunit.png']]
        
        # Debug logging
        self.logger.info(f"Enhanced processing folder: {folder_name}")
        self.logger.info(f"Total filtered images: {len(filtered_image_paths)}")
        self.logger.info(f"All image names: {[os.path.basename(img).lower() for img in filtered_image_paths]}")
        self.logger.info(f"Teaser images: {len(teaser_images)} - {[os.path.basename(img) for img in teaser_images]}")
        self.logger.info(f"Mainunit images: {len(mainunit_images)} - {[os.path.basename(img) for img in mainunit_images]}")
        self.logger.info(f"Other images: {len(other_images)} - {[os.path.basename(img) for img in other_images]}")
        
        # Create first slide with teaser and mainunit images
        primary_images = teaser_images + mainunit_images
        if primary_images:
            self._create_slide_with_images(prs, folder_name, primary_images, annotation_option, slide_number=1)
            
            # If there are other images, create additional slides with same count as primary slide
            if other_images:
                images_per_slide = len(primary_images) if primary_images else 2  # Default to 2 if no primary images
                self.logger.info(f"Creating additional slides for {len(other_images)} other images, {images_per_slide} images per slide")
                
                # Create additional slides with remaining images
                for i in range(0, len(other_images), images_per_slide):
                    slide_images = other_images[i:i + images_per_slide]
                    slide_number = (i // images_per_slide) + 2  # Start from 2 since first slide is primary
                    self._create_slide_with_images(prs, folder_name, slide_images, annotation_option, slide_number)
        else:
            # No teaser/mainunit images, just create slides with other images
            if other_images:
                self.logger.info(f"No teaser/mainunit images found, creating slides with {len(other_images)} other images")
                # Split images into chunks that fit on slides (max 9 per slide)
                max_images_per_slide = 9
                image_chunks = [other_images[i:i + max_images_per_slide] for i in range(0, len(other_images), max_images_per_slide)]
                
                for chunk_index, image_chunk in enumerate(image_chunks):
                    slide_number = chunk_index + 1
                    self._create_slide_with_images(prs, folder_name, image_chunk, annotation_option, slide_number, len(image_chunks))
    
    def _create_slide_with_images(self, prs, folder_name, image_list, annotation_option, slide_number=1, total_slides=None):
        """Create a single slide with the given images."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        formatted_name = self._format_folder_name(folder_name)
        if slide_number == 1 and total_slides is None:
            title_text = formatted_name
        elif total_slides is not None:
            title_text = f"{formatted_name} ({slide_number}/{total_slides})"
        else:
            title_text = f"{formatted_name} ({slide_number})"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Use same positioning as primary slide for this group of images
        self._arrange_images_on_slide(slide, image_list, annotation_option, folder_name)
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)

    def _add_slide_with_images_multi_tab(self, prs, folder_name, image_paths, annotation_option='with_annos'):
        """Add slides with folder name as title and images arranged on the slides - Multi-tab version with additional slides for remaining images."""
        # Filter out teaser-disclaimer.png and mainunit-disclaimer.png files
        filtered_image_paths = [
            img_path for img_path in image_paths 
            if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                   os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
        ]
        
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        filtered_image_paths.sort(key=self._sort_images_exact_priority)
        
        # Skip if no images left after filtering
        if not filtered_image_paths:
            return
        
        # Find exact teaser.png and mainunit.png files, plus other images
        teaser_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'teaser.png']
        mainunit_images = [img for img in filtered_image_paths if os.path.basename(img).lower() == 'mainunit.png']
        other_images = [img for img in filtered_image_paths if os.path.basename(img).lower() not in ['teaser.png', 'mainunit.png']]
        
        # Debug logging
        self.logger.info(f"Multi-tab processing folder: {folder_name}")
        self.logger.info(f"Total filtered images: {len(filtered_image_paths)}")
        self.logger.info(f"Teaser images: {len(teaser_images)}")
        self.logger.info(f"Mainunit images: {len(mainunit_images)}")
        self.logger.info(f"Other images: {len(other_images)}")
        for i, img in enumerate(other_images):
            self.logger.info(f"Other image {i+1}: {os.path.basename(img)}")
        
        # Create first slide with teaser and mainunit images
        primary_images = teaser_images + mainunit_images
        if primary_images:
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
            title_bg.line.fill.background()
            title_bg.shadow.inherit = False  # Remove shadow
            
            # Add title text
            formatted_name = self._format_folder_name(folder_name)
            title_text = formatted_name
            
            # Add title text box with exact specifications
            title_textbox = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),  # Convert cm to inches
                Inches(0.38 / 2.54),  # Convert cm to inches
                Inches(12 / 2.54),    # Convert cm to inches
                Inches(1 / 2.54)      # Convert cm to inches
            )
            title_textbox.text_frame.text = title_text
            
            # Format the title text
            title_paragraph = title_textbox.text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Arrange primary images on the slide
            self._arrange_images_on_slide(slide, primary_images, annotation_option, folder_name)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(slide)
        
        # Create additional slides for remaining images using same positioning
        if other_images:
            # Group other images by the number of primary images (teaser + mainunit)
            images_per_slide = len(primary_images) if primary_images else 2  # Default to 2 if no primary images
            
            # Create slides with grouped images
            for i in range(0, len(other_images), images_per_slide):
                slide_images = other_images[i:i + images_per_slide]
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Remove all placeholder shapes to prevent "Title 1" text
                self._remove_placeholders(slide)
                
                # Add gray rectangle background for title (1.79cm height)
                title_bg = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0), Inches(0), 
                    prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
                )
                title_bg.fill.solid()
                title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
                title_bg.line.fill.background()
                title_bg.shadow.inherit = False  # Remove shadow
                
                # Add title text with sequence number
                formatted_name = self._format_folder_name(folder_name)
                slide_number = (i // images_per_slide) + 2  # Start from 2 since first slide is primary
                title_text = f"{formatted_name} ({slide_number})"
                
                # Add title text box with exact specifications
                title_textbox = slide.shapes.add_textbox(
                    Inches(0.51 / 2.54),  # Convert cm to inches
                    Inches(0.38 / 2.54),  # Convert cm to inches
                    Inches(12 / 2.54),    # Convert cm to inches
                    Inches(1 / 2.54)      # Convert cm to inches
                )
                title_textbox.text_frame.text = title_text
                
                # Format the title text
                title_paragraph = title_textbox.text_frame.paragraphs[0]
                title_paragraph.font.name = "Aptos Display"
                title_paragraph.font.size = Pt(18)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.LEFT
                title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                
                # Use same positioning as primary slide for this group of images
                self._arrange_images_on_slide(slide, slide_images, annotation_option, folder_name)
                
                # Add VDX TV logo to slide
                self._add_vdx_logo(slide)
    
    def _add_consolidated_teaser_slide(self, prs, teaser_images, annotation_option='with_annos'):
        """Add a consolidated slide with all teaser images."""
        if not teaser_images:
            return
            
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text
        title_text = "DESKTOP EXPANDABLE - ALL TEASERS"
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12 / 2.54),    # Convert cm to inches
            Inches(1 / 2.54)      # Convert cm to inches
        )
        title_textbox.text_frame.text = title_text
        
        # Format the title text
        title_paragraph = title_textbox.text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Arrange teaser images with custom positioning
        self._arrange_teaser_images_with_custom_positions(slide, teaser_images, annotation_option)
        
        # Add VDX TV logo to slide
        self._add_vdx_logo(slide)
    
    def _add_consolidated_mobile_teaser_slide(self, prs, mobile_teaser_images, annotation_option='with_annos'):
        """Add consolidated slides with all teaser images from Mobile Expandable folders using 3-position layout."""
        if not mobile_teaser_images:
            return
            
        # Process images in chunks of 3 (same as engaged slide layout)
        images_per_slide = 3
        slide_num = 1
        
        for i in range(0, len(mobile_teaser_images), images_per_slide):
            slide_images = mobile_teaser_images[i:i + images_per_slide]
            
            # Create slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove all placeholder shapes to prevent "Title 1" text
            self._remove_placeholders(slide)
            
            # Add gray rectangle background for title (1.79cm height)
            title_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), 
                prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
            )
            title_bg.fill.solid()
            title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
            title_bg.line.fill.background()
            title_bg.shadow.inherit = False  # Remove shadow
            
            # Add title text with continuation marker if needed
            if slide_num == 1:
                title_text = "MOBILE EXPANDABLE - ENGAGED"
            else:
                title_text = "MOBILE EXPANDABLE - ENGAGED (CONTD.)"
            
            # Add title text box with exact specifications
            title_textbox = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),  # Convert cm to inches
                Inches(0.38 / 2.54),  # Convert cm to inches
                Inches(12 / 2.54),    # Convert cm to inches
                Inches(1 / 2.54)      # Convert cm to inches
            )
            title_textbox.text_frame.text = title_text
            
            # Format the title text
            title_paragraph = title_textbox.text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Use same 3-position layout as Mobile Expandable - Engaged slide
            self._arrange_mobile_engaged_images_with_custom_positions(slide, slide_images, annotation_option)
            
            # Add VDX TV logo to slide
            self._add_vdx_logo(slide)
            
            slide_num += 1
            self.logger.info(f"Created Mobile Expandable - Engaged slide {slide_num - 1} with {len(slide_images)} images")
    
    def _arrange_teaser_images_with_custom_positions(self, slide, teaser_images, annotation_option='with_annos'):
        """Arrange teaser images with specific dimensions and positions."""
        # Define custom positioning for different ad sizes
        positioning_specs = {
            '970x250': {
                'width_cm': 17.84,
                'height_cm': 4.6,
                'horizontal_cm': 0.7,
                'vertical_cm': 2.13
            },
            '728x90': {
                'width_cm': 16.55,
                'height_cm': 2.05,
                'horizontal_cm': 0.7,
                'vertical_cm': 7.95
            },
            '300x250': {
                'width_cm': 7.35,
                'height_cm': 6.13,
                'horizontal_cm': 0.7,
                'vertical_cm': 11.27
            },
            '300x600': {
                'width_cm': 7.57,
                'height_cm': 15.13,
                'horizontal_cm': 20.64,
                'vertical_cm': 2.17
            },
            '160x600': {
                'width_cm': 4.04,
                'height_cm': 15.13,
                'horizontal_cm': 28.99,
                'vertical_cm': 2.14
            }
        }
        
        for img_path in teaser_images:
            # Extract size from file path or folder name
            size_key = self._extract_size_from_path(img_path)
            
            if size_key in positioning_specs:
                spec = positioning_specs[size_key]
                
                # Convert cm to inches
                width_inches = spec['width_cm'] / 2.54
                height_inches = spec['height_cm'] / 2.54
                x_inches = spec['horizontal_cm'] / 2.54
                y_inches = spec['vertical_cm'] / 2.54
                
                # Add image with exact positioning
                try:
                    picture_shape = slide.shapes.add_picture(
                        img_path,
                        Inches(x_inches),
                        Inches(y_inches),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    
                    # Add border (desktop images get borders)
                    picture_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black color
                    picture_shape.line.width = Pt(0.5)  # 0.5pt width
                    
                    # Add dimension label underneath the image
                    label_y = y_inches + height_inches + 0.1  # 0.1 inch spacing below image
                    label_x = x_inches + (width_inches / 2) - 0.5  # Center horizontally (approximate)
                    
                    # Create text box for dimension label
                    label_textbox = slide.shapes.add_textbox(
                        Inches(label_x),
                        Inches(label_y),
                        Inches(1),  # Width of text box
                        Inches(0.3)  # Height of text box
                    )
                    
                    # Set the label text
                    label_textbox.text_frame.text = size_key
                    
                    # Format the label text
                    label_paragraph = label_textbox.text_frame.paragraphs[0]
                    label_paragraph.font.name = "Aptos Display"
                    label_paragraph.font.size = Pt(10)
                    label_paragraph.font.bold = True
                    label_paragraph.alignment = PP_ALIGN.CENTER
                    label_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                    
                    self.logger.info(f"Added teaser image {os.path.basename(img_path)} with size {size_key} at position ({x_inches:.2f}, {y_inches:.2f}) and label at ({label_x:.2f}, {label_y:.2f})")
                        
                except Exception as e:
                    self.logger.error(f"Error adding teaser image {img_path}: {str(e)}")
            else:
                self.logger.warning(f"No positioning spec found for size {size_key} in image {img_path}")
    
    def _extract_size_from_path(self, img_path):
        """Extract size key from image path or folder name."""
        # Check folder path for size patterns
        path_parts = img_path.split('/')
        for part in path_parts:
            if '970x250' in part:
                return '970x250'
            elif '728x90' in part:
                return '728x90'
            elif '300x250' in part:
                return '300x250'
            elif '300x600' in part:
                return '300x600'
            elif '160x600' in part:
                return '160x600'
            elif '320x50' in part:
                return '320x50'
        
        # If no size found in path, return empty string
        return ''
    
    def _arrange_mobile_teaser_images_with_custom_positions(self, slide, mobile_teaser_images, annotation_option='with_annos'):
        """Arrange mobile teaser images with custom positioning for different sizes."""
        # Define positioning specifications for mobile teaser images
        # All images have same height (16cm) and width (8.22cm), only x positions differ
        positioning_specs = {
            '300x250': {
                'width_cm': 8.22,
                'height_cm': 16,
                'horizontal_cm': 2.6,
                'vertical_cm': 2.17
            },
            '300x600': {
                'width_cm': 8.22,
                'height_cm': 16,
                'horizontal_cm': 12.11,
                'vertical_cm': 2.17
            },
            '320x50': {
                'width_cm': 8.22,
                'height_cm': 16,
                'horizontal_cm': 21.56,
                'vertical_cm': 2.17
            }
        }
        
        for img_path in mobile_teaser_images:
            # Extract size from file path or folder name
            size_key = self._extract_size_from_path(img_path)
            
            # Check for 320x50 size which might not be in folder name
            if '320x50' in img_path or '320x50' in os.path.basename(img_path):
                size_key = '320x50'
            
            if size_key in positioning_specs:
                spec = positioning_specs[size_key]
                
                # Convert cm to inches
                width_inches = spec['width_cm'] / 2.54
                height_inches = spec['height_cm'] / 2.54
                x_inches = spec['horizontal_cm'] / 2.54
                y_inches = spec['vertical_cm'] / 2.54
                
                # Add image with exact positioning
                try:
                    picture_shape = slide.shapes.add_picture(
                        img_path,
                        Inches(x_inches),
                        Inches(y_inches),
                        Inches(width_inches),
                        Inches(height_inches)
                    )
                    
                    # Mobile images don't get borders
                    # No border for mobile images
                    
                    self.logger.info(f"Added mobile teaser image {os.path.basename(img_path)} with size {size_key} at position ({x_inches:.2f}, {y_inches:.2f})")
                        
                except Exception as e:
                    self.logger.error(f"Error adding mobile teaser image {img_path}: {str(e)}")
            else:
                self.logger.warning(f"No positioning spec found for size {size_key} in mobile teaser image {img_path}")
    
    def _arrange_images_on_slide(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange images on a slide based on the number of images."""
        if not image_paths:
            return
        
        # Special handling for Desktop In-frame images
        if 'vdxdesktopinframe' in folder_name.lower():
            if '160x600' in folder_name:
                self._arrange_desktop_inframe_160x600_images(slide, image_paths, annotation_option, folder_name)
                return
            elif '300x250' in folder_name:
                self._arrange_desktop_inframe_300x250_images(slide, image_paths, annotation_option, folder_name)
                return
            elif '300x600' in folder_name:
                self._arrange_desktop_inframe_300x600_images(slide, image_paths, annotation_option, folder_name)
                return
            elif '970x250' in folder_name:
                self._arrange_desktop_inframe_970x250_images(slide, image_paths, annotation_option, folder_name)
                return
            elif '728x90' in folder_name:
                self._arrange_desktop_inframe_728x90_images(slide, image_paths, annotation_option, folder_name)
                return
        
        # Special handling for Mobile In-frame images
        if 'vdxmobileinframe' in folder_name.lower():
            if '300x250' in folder_name:
                self._arrange_mobile_inframe_300x250_images(slide, image_paths, annotation_option, folder_name)
                return
            elif '300x600' in folder_name:
                self._arrange_mobile_inframe_300x600_images(slide, image_paths, annotation_option, folder_name)
                return
        
        # Special handling for Mobile Instream images (same positioning as Mobile In-frame 300x600)
        if 'vdxmobileinstream' in folder_name.lower():
            self._arrange_mobile_instream_images(slide, image_paths, annotation_option, folder_name)
            return
        
        # Calculate available space (excluding title area)
        available_width = Inches(12.33)
        available_height = Inches(5.5)
        start_x = Inches(0.5)
        start_y = Inches(1.79 / 2.54 + 0.2)  # Below title with margin
        
        num_images = len(image_paths)
        
        # Calculate grid dimensions
        cols = 3 if num_images > 6 else 2
        rows = (num_images + cols - 1) // cols
        
        # Calculate image dimensions with proper spacing
        spacing = 0.25  # 0.25 inch spacing between images
        img_width = (available_width.inches - spacing * (cols - 1)) / cols
        img_height = (available_height.inches - spacing * (rows - 1)) / rows
        

        
        # Place images in grid
        for i, img_path in enumerate(image_paths):
            row = i // cols
            col = i % cols
            
            # Calculate position using simple arithmetic
            x = start_x.inches + col * (img_width + spacing)
            y = start_y.inches + row * (img_height + spacing + (0.4 if annotation_option == 'no_annos' else 0))
            
            # Pass raw numbers to avoid Inches arithmetic issues
            self._add_image_to_slide(slide, img_path, x, y, img_width, img_height, annotation_option, folder_name)
    
    def _arrange_desktop_inframe_160x600_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Desktop In-frame 160x600 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Desktop In-frame 160x600
        positioning_specs = {
            'teaser': {
                'height_cm': 15.98,
                'width_cm': 4.26,
                'x_cm': 9.93,
                'y_cm': 2.3
            },
            'mainunit': {
                'height_cm': 15.98,
                'width_cm': 4.26,
                'x_cm': 16.23,
                'y_cm': 2.3
            }
        }
        
        teaser_position = positioning_specs['teaser']
        mainunit_position = positioning_specs['mainunit']
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position
            if os.path.basename(img_path).lower() == 'teaser.png':
                spec = teaser_position
                image_type = 'teaser'
            elif os.path.basename(img_path).lower() == 'mainunit.png':
                spec = mainunit_position
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = teaser_position
                    image_type = 'additional_1'
                else:
                    spec = mainunit_position
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Add border (desktop images get borders)
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black color
                picture_shape.line.width = Pt(0.5)  # 0.5pt width
                
                self.logger.info(f"Added Desktop In-frame 160x600 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Desktop In-frame 160x600 image {img_path}: {str(e)}")
    
    def _arrange_desktop_inframe_300x250_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Desktop In-frame 300x250 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Desktop In-frame 300x250
        positioning_specs = {
            'teaser': {
                'height_cm': 8.53,
                'width_cm': 10.23,
                'x_cm': 5.35,
                'y_cm': 4.59
            },
            'mainunit': {
                'height_cm': 8.53,
                'width_cm': 10.23,
                'x_cm': 17.86,
                'y_cm': 4.59
            }
        }
        
        teaser_position = positioning_specs['teaser']
        mainunit_position = positioning_specs['mainunit']
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position
            if os.path.basename(img_path).lower() == 'teaser.png':
                spec = teaser_position
                image_type = 'teaser'
            elif os.path.basename(img_path).lower() == 'mainunit.png':
                spec = mainunit_position
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = teaser_position
                    image_type = 'additional_1'
                else:
                    spec = mainunit_position
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Add border (desktop images get borders)
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black color
                picture_shape.line.width = Pt(0.5)  # 0.5pt width
                
                self.logger.info(f"Added Desktop In-frame 300x250 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Desktop In-frame 300x250 image {img_path}: {str(e)}")
    
    def _arrange_desktop_inframe_300x600_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Desktop In-frame 300x600 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Desktop In-frame 300x600
        positioning_specs = {
            'teaser': {
                'height_cm': 15.33,
                'width_cm': 7.66,
                'x_cm': 7.63,
                'y_cm': 2.3
            },
            'mainunit': {
                'height_cm': 15.33,
                'width_cm': 7.66,
                'x_cm': 19.69,
                'y_cm': 2.3
            }
        }
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position
            if os.path.basename(img_path).lower() == 'teaser.png':
                spec = positioning_specs['teaser']
                image_type = 'teaser'
            elif os.path.basename(img_path).lower() == 'mainunit.png':
                spec = positioning_specs['mainunit']
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = positioning_specs['teaser']
                    image_type = 'additional_1'
                else:
                    spec = positioning_specs['mainunit']
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Add border (desktop images get borders)
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Added Desktop In-frame 300x600 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Desktop In-frame 300x600 image {img_path}: {str(e)}")
    
    def _arrange_desktop_inframe_970x250_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Desktop In-frame 970x250 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Desktop In-frame 970x250
        positioning_specs = {
            'teaser': {
                'height_cm': 5.7,
                'width_cm': 22.1,
                'x_cm': 1,
                'y_cm': 2.41
            },
            'mainunit': {
                'height_cm': 5.7,
                'width_cm': 22.1,
                'x_cm': 1,
                'y_cm': 10.94
            }
        }
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position using exact filename matching
            filename = os.path.basename(img_path).lower()
            if filename == 'teaser.png':
                spec = positioning_specs['teaser']
                image_type = 'teaser'
            elif filename == 'mainunit.png':
                spec = positioning_specs['mainunit']
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = positioning_specs['teaser']
                    image_type = 'additional_1'
                else:
                    spec = positioning_specs['mainunit']
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Add border (desktop images get borders)
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Added Desktop In-frame 970x250 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Desktop In-frame 970x250 image {img_path}: {str(e)}")
        
        # Annotations disabled for Desktop In-frame 970x250 per user request
    
    def _add_desktop_inframe_970x250_annotations(self, slide):
        """Add specific text boxes for Desktop In-frame 970x250 slide with annotations."""
        # Define annotation text boxes based on user requirements (excluding title as it's added separately)
        annotations = [
            {
                'text': 'ROLLOVER TO EXPLORE" and "VDX-TV" animate in peel back upper right; user click expands...',
                'x_cm': 0.51,
                'y_cm': 1.79,
                'width_cm': 25.40,
                'height_cm': 2.00
            },
            {
                'text': 'Clicking the "X" minimizes the teaser to 970x90.',
                'x_cm': 0.51,
                'y_cm': 3.79,
                'width_cm': 25.40,
                'height_cm': 2.00
            },
            {
                'text': 'Global: The ISI auto scrolls, but the user also has the ability to manually scroll the ISI.',
                'x_cm': 0.51,
                'y_cm': 5.79,
                'width_cm': 25.40,
                'height_cm': 2.00
            },
            {
                'text': 'Teaser State/Pre-engagement State In this state, users see the teaser, ISI, logo, CTA, "ROLLOVER TO EXPLORE"...',
                'x_cm': 0.51,
                'y_cm': 7.79,
                'width_cm': 25.40,
                'height_cm': 2.00
            },
            {
                'text': 'Engaged State Once the users click or hover, the expandable overlays the 970x250 area...',
                'x_cm': 0.51,
                'y_cm': 9.79,
                'width_cm': 25.40,
                'height_cm': 2.00
            }
        ]
        
        for annotation in annotations:
            try:
                # Convert cm to inches
                x_inches = annotation['x_cm'] / 2.54
                y_inches = annotation['y_cm'] / 2.54
                width_inches = annotation['width_cm'] / 2.54
                height_inches = annotation['height_cm'] / 2.54
                
                # Add text box
                text_box = slide.shapes.add_textbox(
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Set text content
                text_frame = text_box.text_frame
                text_frame.text = annotation['text']
                
                # Format the text
                paragraph = text_frame.paragraphs[0]
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Arial"
                paragraph.alignment = PP_ALIGN.LEFT
                paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                
                # Set text frame properties
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.05)
                text_frame.margin_bottom = Inches(0.05)
                text_frame.word_wrap = True
                
                self.logger.info(f"Added annotation text box at position ({x_inches:.2f}, {y_inches:.2f})")
                
            except Exception as e:
                self.logger.error(f"Error adding annotation text box: {str(e)}")
    
    def _arrange_desktop_inframe_728x90_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Desktop In-frame 728x90 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Desktop In-frame 728x90
        positioning_specs = {
            'teaser': {
                'height_cm': 2.68,
                'width_cm': 21.67,
                'x_cm': 1,
                'y_cm': 3.28
            },
            'mainunit': {
                'height_cm': 2.68,
                'width_cm': 21.67,
                'x_cm': 1,
                'y_cm': 8.66
            }
        }
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position using exact filename matching
            filename = os.path.basename(img_path).lower()
            if filename == 'teaser.png':
                spec = positioning_specs['teaser']
                image_type = 'teaser'
            elif filename == 'mainunit.png':
                spec = positioning_specs['mainunit']
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = positioning_specs['teaser']
                    image_type = 'additional_1'
                else:
                    spec = positioning_specs['mainunit']
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Add border (desktop images get borders)
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Added Desktop In-frame 728x90 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Desktop In-frame 728x90 image {img_path}: {str(e)}")
    
    def _arrange_mobile_inframe_300x250_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Mobile In-frame 300x250 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Mobile In-frame 300x250
        positioning_specs = {
            'teaser': {
                'height_cm': 16.02,
                'width_cm': 8.23,
                'x_cm': 5.02,
                'y_cm': 2
            },
            'mainunit': {
                'height_cm': 16.02,
                'width_cm': 8.23,
                'x_cm': 18.7,
                'y_cm': 2
            }
        }
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position
            if os.path.basename(img_path).lower() == 'teaser.png':
                spec = positioning_specs['teaser']
                image_type = 'teaser'
            elif os.path.basename(img_path).lower() == 'mainunit.png':
                spec = positioning_specs['mainunit']
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = positioning_specs['teaser']
                    image_type = 'additional_1'
                else:
                    spec = positioning_specs['mainunit']
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Mobile images don't get borders
                # No border for mobile images
                
                self.logger.info(f"Added Mobile In-frame 300x250 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Mobile In-frame 300x250 image {img_path}: {str(e)}")
    
    def _arrange_mobile_inframe_300x600_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Mobile In-frame 300x600 images with specific positioning."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Define specific positioning for Mobile In-frame 300x600
        positioning_specs = {
            'teaser': {
                'height_cm': 16.02,
                'width_cm': 8.23,
                'x_cm': 5.02,
                'y_cm': 2
            },
            'mainunit': {
                'height_cm': 16.02,
                'width_cm': 8.23,
                'x_cm': 18.7,
                'y_cm': 2
            }
        }
        
        for i, img_path in enumerate(sorted_images):
            # Determine image type and position
            if os.path.basename(img_path).lower() == 'teaser.png':
                spec = positioning_specs['teaser']
                image_type = 'teaser'
            elif os.path.basename(img_path).lower() == 'mainunit.png':
                spec = positioning_specs['mainunit']
                image_type = 'mainunit'
            else:
                # For additional images, use positioning based on index
                if i == 0:
                    spec = positioning_specs['teaser']
                    image_type = 'additional_1'
                else:
                    spec = positioning_specs['mainunit']
                    image_type = 'additional_2'
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            # Add image with exact positioning
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path,
                    Inches(x_inches),
                    Inches(y_inches),
                    Inches(width_inches),
                    Inches(height_inches)
                )
                
                # Mobile images don't get borders
                # No border for mobile images
                
                self.logger.info(f"Added Mobile In-frame 300x600 {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding Mobile In-frame 300x600 image {img_path}: {str(e)}")
    
    def _arrange_mobile_instream_images(self, slide, image_paths, annotation_option='with_annos', folder_name=''):
        """Arrange Mobile Instream images with cropping and specific dimensions for Auto tab."""
        # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
        sorted_images = sorted(image_paths, key=self._sort_images_exact_priority)
        
        # Calculate slide dimensions for centering
        slide_width_cm = 33.867  # Standard slide width in cm
        slide_height_cm = 19.05  # Standard slide height in cm
        
        # Image specifications: height 11.28cm, width 12.76cm
        img_height_cm = 11.28
        img_width_cm = 12.76
        
        # Calculate spacing for centering two images
        total_images_width = 2 * img_width_cm
        horizontal_spacing = 2.0  # cm between images
        total_width_with_spacing = total_images_width + horizontal_spacing
        start_x_cm = (slide_width_cm - total_width_with_spacing) / 2
        
        # Vertical centering with proper spacing from gray rectangle (4cm from top)
        y_position_cm = 4.0  # Below gray rectangle with proper spacing
        
        positioning_specs = {
            'first': {'height_cm': img_height_cm, 'width_cm': img_width_cm, 'x_cm': start_x_cm, 'y_cm': y_position_cm},
            'second': {'height_cm': img_height_cm, 'width_cm': img_width_cm, 'x_cm': start_x_cm + img_width_cm + horizontal_spacing, 'y_cm': y_position_cm}
        }
        
        for i, img_path in enumerate(sorted_images[:2]):  # Only take first 2 images
            position_key = 'first' if i == 0 else 'second'
            spec = positioning_specs[position_key]
            
            # Crop image from bottom to 774px height
            cropped_img_path = self.image_processor._crop_image_from_bottom(img_path, 774)
            
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    cropped_img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Mobile images don't get borders
                
                # Determine image type for logging
                if os.path.basename(img_path).lower() == 'teaser.png':
                    image_type = 'teaser'
                elif os.path.basename(img_path).lower() == 'mainunit.png':
                    image_type = 'mainunit'
                else:
                    image_type = f'image_{i+1}'
                
                self.logger.info(f"Added Mobile Instream {image_type} image {os.path.basename(img_path)} at position ({x_inches:.2f}, {y_inches:.2f}) with cropped height 774px")
                    
            except Exception as e:
                self.logger.error(f"Error adding Mobile Instream image {img_path}: {str(e)}")
    
    def _add_image_to_slide(self, slide, image_path, x, y, max_width, max_height, annotation_option='with_annos', folder_name=''):
        """Add an image to the slide with proper sizing."""
        # Convert all input values to inches (numbers) first
        x_inches = x.inches if hasattr(x, 'inches') else float(x)
        y_inches = y.inches if hasattr(y, 'inches') else float(y)
        max_w_inches = max_width.inches if hasattr(max_width, 'inches') else float(max_width)
        max_h_inches = max_height.inches if hasattr(max_height, 'inches') else float(max_height)
        
        try:
            
            # Get image dimensions
            with Image.open(image_path) as img:
                img_width, img_height = img.size
            
            # Calculate aspect ratio
            aspect_ratio = img_width / img_height
            
            # Special handling for mobile images: fixed height of 16cm (6.299 inches)
            if "mobile" in folder_name.lower():
                fixed_height_cm = 16
                fixed_height_inches = fixed_height_cm / 2.54  # Convert cm to inches
                height_inches = fixed_height_inches
                width_inches = fixed_height_inches * aspect_ratio
                
                # Ensure width doesn't exceed max available width
                if width_inches > max_w_inches:
                    width_inches = max_w_inches
                    height_inches = max_w_inches / aspect_ratio
            else:
                # Calculate the size to fit within the max dimensions for non-mobile images
                if aspect_ratio > (max_w_inches / max_h_inches):
                    # Width is the limiting factor
                    width_inches = max_w_inches
                    height_inches = max_w_inches / aspect_ratio
                else:
                    # Height is the limiting factor
                    height_inches = max_h_inches
                    width_inches = max_h_inches * aspect_ratio
            
            # Center the image within the available space
            center_x_inches = x_inches + (max_w_inches - width_inches) / 2
            center_y_inches = y_inches + (max_h_inches - height_inches) / 2
            
            # Convert back to Inches objects for pptx
            center_x = Inches(center_x_inches)
            center_y = Inches(center_y_inches)
            width = Inches(width_inches)
            height = Inches(height_inches)
            
            # Add the image to the slide
            picture_shape = slide.shapes.add_picture(image_path, center_x, center_y, width, height)
            
            # Add border if folder name doesn't contain "mobile" and isn't OTT/CTV
            if ("mobile" not in folder_name.lower() and 
                "ott" not in folder_name.lower() and 
                "ctv" not in folder_name.lower()):
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black color
                picture_shape.line.width = Pt(0.5)  # 0.5pt width
            
            self.logger.info(f"Successfully added image {os.path.basename(image_path)} at position ({center_x_inches:.2f}, {center_y_inches:.2f}) with size ({width_inches:.2f}, {height_inches:.2f})")
            
            # Add filename annotation if "with_annos" is selected

            
        except Exception as e:
            self.logger.error(f"Error adding image {image_path}: {str(e)}")
            # Add a placeholder text if image fails to load
            placeholder = slide.shapes.add_textbox(Inches(x_inches), Inches(y_inches), Inches(max_w_inches), Inches(max_h_inches))
            placeholder.text_frame.text = f"Error loading image:\n{os.path.basename(image_path)}"
    
    def _add_full_isi_slide(self, prs, mainunit_disclaimer_files, annotation_option='with_annos'):
        """Add FULL ISI slide(s) using the SlideCreator."""
        from .slide_creator import SlideCreator
        slide_creator = SlideCreator(self, self.image_processor)
        result = slide_creator._create_full_isi_slide(prs, mainunit_disclaimer_files)
        
        # Add FULL ISI text boxes with "With Annos" option if it's a single slide
        if annotation_option == 'with_annos' and not isinstance(result, int):
            if prs.slides:
                last_slide = prs.slides[-1]
                self._add_full_isi_textboxes(last_slide)
        
        return result
    
    def _add_disclaimer_images(self, slide, disclaimer_files):
        """Add disclaimer images with fixed height of 16.5cm."""
        if not disclaimer_files:
            return
        
        # Fixed height in cm converted to inches
        fixed_height_cm = 16.5
        fixed_height_inches = fixed_height_cm / 2.54
        
        # Fixed positioning as requested: x = 0.5cm, y = 2.42cm
        start_x_cm = 0.5
        start_y_cm = 2.42
        start_x = Inches(start_x_cm / 2.54)
        start_y = Inches(start_y_cm / 2.54)
        available_width = Inches(12.33)
        
        # Calculate how many images can fit horizontally
        num_images = len(disclaimer_files)
        spacing = 0.25  # 0.25 inch spacing between images
        
        for i, img_path in enumerate(disclaimer_files):
            try:
                # Get image dimensions to calculate width based on aspect ratio
                with Image.open(img_path) as img:
                    img_width, img_height = img.size
                
                # Calculate aspect ratio
                aspect_ratio = img_width / img_height
                
                # Calculate width based on fixed height and aspect ratio
                width_inches = fixed_height_inches * aspect_ratio
                
                # Use fixed positioning: x = 0.5cm, y = 2.42cm for first image
                # Subsequent images positioned horizontally with spacing
                x_pos = start_x.inches + (i * (width_inches + spacing))
                y_pos = start_y.inches
                
                # Add the image to the slide
                picture_shape = slide.shapes.add_picture(
                    img_path, 
                    Inches(x_pos), 
                    Inches(y_pos), 
                    Inches(width_inches), 
                    Inches(fixed_height_inches)
                )
                
                # Log the positioning for confirmation
                self.logger.info(f"Added disclaimer image {os.path.basename(img_path)} at position ({x_pos*2.54:.2f}cm, {y_pos*2.54:.2f}cm) with dimensions {width_inches*2.54:.2f}cm x {fixed_height_inches*2.54:.2f}cm")
                
                # No border for FULL ISI disclaimer images
                
                self.logger.info(f"Added disclaimer image {os.path.basename(img_path)} with height {fixed_height_cm}cm")
                
            except Exception as e:
                self.logger.error(f"Error adding disclaimer image {img_path}: {str(e)}")
                # Add placeholder text if image fails to load
                placeholder = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2), Inches(4), Inches(1)
                )
                placeholder.text_frame.text = f"Error loading disclaimer image:\n{os.path.basename(img_path)}"
    
    def _get_aspect_ratio(self, img_path):
        """Get aspect ratio of an image."""
        try:
            with Image.open(img_path) as img:
                img_width, img_height = img.size
                return img_width / img_height
        except:
            return 1.0  # Default aspect ratio if image can't be opened
    
    def _remove_placeholders(self, slide):
        """Remove all placeholder shapes from a slide to prevent unwanted text like 'Title 1'."""
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                shapes_to_remove.append(shape)
        
        # Remove placeholder shapes
        for shape in shapes_to_remove:
            slide.shapes._spTree.remove(shape._element)
    
    def _add_slides_in_order_manual(self, prs, folder_structure, annotation_option, all_desktop_teaser_images, all_mobile_teaser_images, implement_video_frames=False):
        """Add slides in order for Manual tab - sequential processing regardless of naming."""
        self.logger.info("MANUAL TAB ENTRY: _add_slides_in_order_manual function STARTING!!!")
        self.logger.info(f"MANUAL TAB: annotation_option={annotation_option}, implement_video_frames={implement_video_frames}")
        self.logger.info(f"MANUAL TAB: Number of folders to process: {len(folder_structure)}")
        
        # Helper function to create manual slides for a folder
        def create_manual_slides(folder_name, image_paths):
            """Create slides with sequential image processing, maintaining format-specific positioning."""
            if not image_paths:
                return
                
            # Filter out disclaimer images
            filtered_images = [
                img_path for img_path in image_paths 
                if not (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                       os.path.basename(img_path).lower() == 'mainunit-disclaimer.png')
            ]
            
            if not filtered_images:
                return
                
            # Use exact filename priority for regular slides: teaser.png first, then mainunit.png, then sequential order
            filtered_images.sort(key=self._sort_images_exact_priority)
            
            self.logger.info(f"Manual tab: Processing {len(filtered_images)} images for {folder_name}")
            self.logger.info(f"Manual tab: Images in order: {[os.path.basename(img) for img in filtered_images]}")
            
            # Determine images per slide based on folder type
            images_per_slide = 2  # Default for most formats
            
            # Video folder special handling - 6 images per slide in 3x2 grid
            if 'video' in folder_name.lower():
                images_per_slide = 6
            
            # Desktop In-frame specific image capacities
            elif 'vdxdesktopinframe' in folder_name.lower():
                if '160x600' in folder_name:
                    images_per_slide = 7
                elif '300x250' in folder_name:
                    images_per_slide = 6
                elif '300x600' in folder_name:
                    images_per_slide = 4
                elif '728x90' in folder_name:
                    images_per_slide = 5
            
            # Mobile slides (any folder with "mobile" in name) - different capacities
            elif 'mobile' in folder_name.lower():
                if 'vdxmobileinstream' in folder_name.lower():
                    images_per_slide = 2  # Mobile Instream: 2 images per slide
                else:
                    images_per_slide = 4  # Other mobile: 4 images per slide
            
            # Process images in chunks
            for i in range(0, len(filtered_images), images_per_slide):
                slide_images = filtered_images[i:i + images_per_slide]
                slide_number = (i // images_per_slide) + 1
                
                # Create slide
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Remove placeholder shapes
                self._remove_placeholders(slide)
                
                # Add title background and text
                title_text = self._format_folder_name(folder_name)
                if slide_number > 1:
                    title_text += " (Contd.)"
                
                # Add gray rectangle background
                rectangle = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0),
                    Inches(0),
                    Inches(13.33),
                    Inches(1.79 / 2.54)
                )
                rectangle.fill.solid()
                rectangle.fill.fore_color.rgb = RGBColor(242, 242, 242)
                rectangle.line.fill.background()
                rectangle.shadow.inherit = False  # Remove shadow
                
                # Add title text
                title_text_box = slide.shapes.add_textbox(
                    Inches(0.51 / 2.54),
                    Inches(0.38 / 2.54),
                    Inches(12 / 2.54),
                    Inches(1 / 2.54)
                )
                title_text_frame = title_text_box.text_frame
                title_text_frame.text = title_text.upper()
                title_paragraph = title_text_frame.paragraphs[0]
                title_paragraph.font.name = "Aptos Display"
                title_paragraph.font.size = Pt(18)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.LEFT
                title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
                
                # Add images with format-specific positioning
                is_first_slide = (slide_number == 1)
                self._add_manual_images_to_slide(slide, slide_images, annotation_option, folder_name, is_first_slide)
                
                # Add VDX TV logo
                self._add_vdx_logo(slide, folder_name)
                
                self.logger.info(f"Manual tab: Created slide {slide_number} for {folder_name} with {len(slide_images)} images")
        
        # Process slides in specific order: Desktop In-frame first, then Desktop Instream, Mobile Instream
        # CTV and OTT are handled separately before FULL ISI slide
        folder_order = ['vdxdesktopinframe', 'vdxdesktopinstream', 'vdxmobileinstream']
        
        desktop_inframe_processed = False
        for folder_type in folder_order:
            self.logger.info(f"Manual tab: Processing folder type {folder_type}")
            
            # For Desktop In-frame folders, apply custom sorting
            if folder_type == 'vdxdesktopinframe':
                # Define the desired order for Desktop In-frame formats
                desktop_inframe_order = ['970x250', '300x250', '300x600', '160x600', '728x90']
                
                # Get all vdxdesktopinframe folders
                desktop_inframe_folders = [(folder_name, image_paths) for folder_name, image_paths in folder_structure.items()
                                         if 'vdxdesktopinframe' in folder_name.lower()]
                
                # Sort Desktop In-frame folders according to the desired order
                def get_desktop_inframe_priority(folder_item):
                    folder_name = folder_item[0]
                    for idx, format_name in enumerate(desktop_inframe_order):
                        if format_name in folder_name:
                            return idx
                    return 999  # Unknown formats go last
                
                desktop_inframe_folders.sort(key=get_desktop_inframe_priority)
                self.logger.info(f"Manual tab: Desktop In-frame folders will be processed in order: {[folder_name for folder_name, _ in desktop_inframe_folders]}")
                
                # Process Desktop In-frame folders in the sorted order
                for folder_name, image_paths in desktop_inframe_folders:
                    # Skip special folders that are handled separately
                    if any(skip_folder in folder_name.lower() for skip_folder in ['engaged']):
                        self.logger.info(f"Manual tab: Skipping engaged folder {folder_name}")
                        continue
                        
                    # Skip desktop expandable folders (they're handled by consolidated slides)
                    if 'vdxdesktopexpandable' in folder_name.lower():
                        self.logger.info(f"Manual tab: Skipping desktop expandable folder {folder_name}")
                        continue
                    
                    # Skip mobile expandable folders (they're handled by consolidated slides)
                    if 'vdxmobileexpandable' in folder_name.lower():
                        self.logger.info(f"Manual tab: Skipping mobile expandable folder {folder_name}")
                        continue
                    
                    # Skip video folder since it's already processed as first slide
                    if 'video' in folder_name.lower():
                        self.logger.info(f"Manual tab: Skipping video folder {folder_name}")
                        continue
                    
                    self.logger.info(f"Manual tab: PROCESSING Desktop In-frame folder {folder_name}")
                    # Store slide count before creating slides
                    slides_before = len(prs.slides)
                    create_manual_slides(folder_name, image_paths)
                    desktop_inframe_processed = True
            else:
                # For non-Desktop In-frame folders, use the original logic
                for folder_name, image_paths in folder_structure.items():
                    # Skip special folders that are handled separately
                    if any(skip_folder in folder_name.lower() for skip_folder in ['engaged']):
                        self.logger.info(f"Manual tab: Skipping engaged folder {folder_name}")
                        continue
                        
                    # Skip desktop expandable folders (they're handled by consolidated slides)
                    if 'vdxdesktopexpandable' in folder_name.lower():
                        self.logger.info(f"Manual tab: Skipping desktop expandable folder {folder_name}")
                        continue
                    
                    # Skip mobile expandable folders (they're handled by consolidated slides)
                    if 'vdxmobileexpandable' in folder_name.lower():
                        self.logger.info(f"Manual tab: Skipping mobile expandable folder {folder_name}")
                        continue
                    
                    # Skip video folder since it's already processed as first slide
                    if 'video' in folder_name.lower():
                        self.logger.info(f"Manual tab: Skipping video folder {folder_name}")
                        continue
                    
                    # Process only the current folder type
                    if folder_type in folder_name.lower():
                        self.logger.info(f"Manual tab: PROCESSING folder {folder_name} for type {folder_type}")
                        # Store slide count before creating slides
                        slides_before = len(prs.slides)
                        create_manual_slides(folder_name, image_paths)
                        
                        # Add video frames functionality for Desktop Instream slides
                        self.logger.info(f"Manual tab: Checking video frames for {folder_type} - implement_video_frames={implement_video_frames}")
                        if folder_type == 'vdxdesktopinstream' and implement_video_frames:
                            self.logger.info("Manual tab: Video frames enabled for Desktop Instream, creating NEW slides")
                            # Create NEW slides with video frames (not overlay on existing slides)
                            slides_after = len(prs.slides)
                            for slide_idx in range(slides_before, slides_after):
                                # For each Desktop Instream slide created, add video frame slides
                                additional_slides = self._implement_video_frames_for_desktop_instream(prs, folder_structure, slide_idx)
                                self.logger.info(f"Manual tab: Created {additional_slides} additional video frame slides for Desktop Instream slide {slide_idx}")
                        elif folder_type == 'vdxdesktopinstream':
                            self.logger.info("Manual tab: Desktop Instream found but video frames disabled")
                        
                        # Track if Desktop In-frame slides were processed
                        if folder_type == 'vdxdesktopinframe':
                            desktop_inframe_processed = True
                    else:
                        self.logger.info(f"Manual tab: Folder {folder_name} does NOT match type {folder_type}")
            
            # Add Desktop Expandable consolidated teaser slide and engaged slides after Desktop In-frame slides are processed
            if folder_type == 'vdxdesktopinframe' and desktop_inframe_processed:
                # First add Desktop Expandable - All Teasers slide if available
                if all_desktop_teaser_images:
                    self._add_consolidated_teaser_slide(prs, all_desktop_teaser_images, annotation_option)
                # Then add Desktop Expandable - Engaged slides (VPM and engaged)
                self._add_vpm_and_engaged_slides(prs, folder_structure, annotation_option, implement_video_frames)
                # Finally add regular Desktop Expandable slide
                self._add_desktop_expandable_slide_manual(prs, folder_structure, annotation_option)
        
        # Process any remaining folders not in the specific order
        for folder_name, image_paths in folder_structure.items():
            # Skip special folders that are handled separately
            if any(skip_folder in folder_name.lower() for skip_folder in ['engaged']):
                self.logger.info(f"Manual tab: Skipping engaged folder {folder_name}")
                continue
                
            # Skip desktop expandable folders (they're handled by consolidated slides)
            if 'vdxdesktopexpandable' in folder_name.lower():
                self.logger.info(f"Manual tab: Skipping desktop expandable folder {folder_name}")
                continue
            
            # Skip mobile expandable folders (they're handled by consolidated slides)
            if 'vdxmobileexpandable' in folder_name.lower():
                self.logger.info(f"Manual tab: Skipping mobile expandable folder {folder_name}")
                continue
            
            # Skip video folder since it's already processed as first slide
            if 'video' in folder_name.lower():
                self.logger.info(f"Manual tab: Skipping video folder {folder_name}")
                continue
                
            # Skip folders already processed in specific order or CTV/OTT (which will be processed before FULL ISI)
            skip_folder = False
            all_processed_folders = folder_order + ['ctv', 'ott']  # Include CTV and OTT in skip list
            for folder_type in all_processed_folders:
                if folder_type in folder_name.lower():
                    self.logger.info(f"Manual tab: Skipping already processed folder {folder_name} (matches {folder_type})")
                    skip_folder = True
                    break
            
            if skip_folder:
                continue
                
            # For Manual tab, create individual slides for each folder with sequential processing
            self.logger.info(f"Manual tab: Processing remaining folder {folder_name}")
            create_manual_slides(folder_name, image_paths)
        
        # Desktop Expandable - All Teasers slide is now handled earlier (before regular Desktop Expandable slide)
        
        
        # Add mobile consolidated slides if needed
        if all_mobile_teaser_images:
            self._add_consolidated_mobile_teaser_slide(prs, all_mobile_teaser_images, annotation_option)
        
        # Add Mobile Expandable engaged slide (special slide)
        self._add_mobile_expandable_engaged_slide(prs, folder_structure, annotation_option)
        
        # Add CTV and OTT slides before FULL ISI slide
        ctv_ott_order = ['ctv', 'ott']
        for folder_type in ctv_ott_order:
            self.logger.info(f"Manual tab: Processing {folder_type} slide before FULL ISI")
            for folder_name, image_paths in folder_structure.items():
                # Skip special folders that are handled separately
                if any(skip_folder in folder_name.lower() for skip_folder in ['engaged']):
                    continue
                    
                # Skip desktop expandable folders (they're handled by consolidated slides)
                if 'vdxdesktopexpandable' in folder_name.lower():
                    continue
                
                # Skip mobile expandable folders (they're handled by consolidated slides)
                if 'vdxmobileexpandable' in folder_name.lower():
                    continue
                
                # Skip video folder since it's already processed as first slide
                if 'video' in folder_name.lower():
                    continue
                
                # Process only the current folder type
                if folder_type in folder_name.lower():
                    self.logger.info(f"Manual tab: PROCESSING folder {folder_name} for type {folder_type} before FULL ISI")
                    create_manual_slides(folder_name, image_paths)
        
        # Add FULL ISI slide as the last slide only if no disclaimer images exist
        # Check if there are any mainunit-disclaimer.png files
        has_disclaimer_images = False
        for folder_name, image_paths in folder_structure.items():
            for img_path in image_paths:
                if os.path.basename(img_path).lower() == 'mainunit-disclaimer.png':
                    has_disclaimer_images = True
                    break
            if has_disclaimer_images:
                break
        
        # Only add blank FULL ISI slide if no disclaimer images found
        if not has_disclaimer_images:
            self._add_full_isi_slide_manual(prs, annotation_option)
        
        # Note: OTT, CTV, Desktop Instream, Mobile Instream are already handled by individual folder processing above
        # No need to add them again with consolidated slide methods
    
    def _add_manual_images_to_slide(self, slide, image_paths, annotation_option, folder_name, is_first_slide=True):
        """Add images to slide using sequential positioning for Manual tab."""
        
        self.logger.info(f"_add_manual_images_to_slide called: folder_name={folder_name}, annotation_option={annotation_option}, is_first_slide={is_first_slide}")
        
        # Video folder special handling - use 3x2 grid layout
        if 'video' in folder_name.lower():
            self.logger.info("Taking video path")
            self._arrange_video_images_3x2_grid(slide, image_paths, annotation_option)
        # Use the same positioning logic as Auto tab but with sequential image assignment
        elif 'vdxdesktopinframe' in folder_name.lower():
            self.logger.info("Taking vdxdesktopinframe path")
            if '970x250' in folder_name:
                self.logger.info("Calling _arrange_desktop_inframe_970x250_images_manual")
                self._arrange_desktop_inframe_970x250_images_manual(slide, image_paths, annotation_option, folder_name, is_first_slide)
            elif '300x250' in folder_name:
                self._arrange_desktop_inframe_300x250_images_manual(slide, image_paths, annotation_option, folder_name)
            elif '300x600' in folder_name:
                self._arrange_desktop_inframe_300x600_images_manual(slide, image_paths, annotation_option, folder_name)
            elif '160x600' in folder_name:
                self._arrange_desktop_inframe_160x600_images_manual(slide, image_paths, annotation_option, folder_name)
            elif '728x90' in folder_name:
                self._arrange_desktop_inframe_728x90_images_manual(slide, image_paths, annotation_option, folder_name)
            else:
                self._arrange_images_on_slide(slide, image_paths, annotation_option, folder_name)
        elif 'mobile' in folder_name.lower():
            if 'vdxmobileinstream' in folder_name.lower():
                self._arrange_mobile_instream_images_manual(slide, image_paths, annotation_option, folder_name, is_first_slide)
            else:
                self._arrange_mobile_images_manual(slide, image_paths, annotation_option, folder_name, is_first_slide)
        elif 'ott' in folder_name.lower():
            self._arrange_ott_images_manual(slide, image_paths, annotation_option, folder_name, is_first_slide)
        else:
            self._arrange_images_on_slide(slide, image_paths, annotation_option, folder_name)
    
    def _arrange_desktop_inframe_970x250_images_manual(self, slide, image_paths, annotation_option, folder_name, is_first_slide=True):
        """Arrange Desktop In-frame 970x250 images sequentially for Manual tab."""
        # Different dimensions for "With Annos" first slide
        if annotation_option == 'with_annos' and is_first_slide:
            positioning_specs = {
                'first': {'height_cm': 6.12, 'width_cm': 23.7, 'x_cm': 1, 'y_cm': 2.41},
                'second': {'height_cm': 6.12, 'width_cm': 23.7, 'x_cm': 1, 'y_cm': 10.94}
            }
        else:
            # Standard dimensions for all other cases
            positioning_specs = {
                'first': {'height_cm': 7.12, 'width_cm': 27.59, 'x_cm': 1, 'y_cm': 2.41},
                'second': {'height_cm': 7.12, 'width_cm': 27.59, 'x_cm': 1, 'y_cm': 10.94}
            }
        
        for i, img_path in enumerate(image_paths[:2]):  # Max 2 images per slide
            position_key = 'first' if i == 0 else 'second'
            spec = positioning_specs[position_key]
            
            # Convert cm to inches
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Add border
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Manual: Added Desktop In-frame 970x250 image {i+1} at position ({x_inches:.2f}, {y_inches:.2f}) with dims {spec['width_cm']}x{spec['height_cm']}cm")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Desktop In-frame 970x250 image {img_path}: {str(e)}")
        
        # Add teaser state text box only for "With Annos" first slide
        self.logger.info(f"Desktop In-frame 970x250 annotation check: annotation_option={annotation_option}, is_first_slide={is_first_slide}")
        if annotation_option == 'with_annos' and is_first_slide:
            self.logger.info("Calling _add_teaser_state_textbox function")
            self._add_teaser_state_textbox(slide)
            # Also add the new "Engaged State" text box
            self._add_engaged_state_textbox(slide)
            # Also add the new "Global" text box
            self._add_global_textbox(slide)
        else:
            self.logger.info(f"Skipping annotations - condition not met (annotation_option={annotation_option}, is_first_slide={is_first_slide})")
    
    def _add_teaser_state_textbox(self, slide):
        """Add teaser state description text box and all annotations for Desktop In-frame 970x250 'With Annos' slides."""
        
        # Text box 1: Teaser State/Pre-engagement State
        # Size: Height 2.82cm, width 8.11cm, x = 25.16cm, y = 3.3cm
        textbox1 = slide.shapes.add_textbox(
            Inches(25.16 / 2.54),  # x position
            Inches(3.3 / 2.54),    # y position  
            Inches(8.11 / 2.54),   # width
            Inches(2.82 / 2.54)    # height
        )
        
        # Set background color to #F8FFD1 and black border
        textbox1.fill.solid()
        textbox1.fill.fore_color.rgb = RGBColor(248, 255, 209)  # #F8FFD1 in RGB
        textbox1.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        textbox1.line.width = Pt(0.5)  # 0.5pt border width
        
        # Add text content
        text_content1 = "Teaser State/Pre-engagement State\n\nIn this state, the user sees the ad unit for the first time on the publisher's site and can either roll over or click the ad, both of which count as an engagement. The ISI auto scrolls."
        textbox1.text_frame.text = text_content1
        
        # Format textbox1
        text_frame1 = textbox1.text_frame
        text_frame1.margin_left = Inches(0.1)
        text_frame1.margin_right = Inches(0.1)
        text_frame1.margin_top = Inches(0.1)
        text_frame1.margin_bottom = Inches(0.1)
        text_frame1.word_wrap = True
        
        for i, paragraph in enumerate(text_frame1.paragraphs):
            if i == 0:  # First paragraph - bold
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.alignment = PP_ALIGN.LEFT
            else:  # Remaining paragraphs - regular
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(9)
                paragraph.font.bold = False
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.alignment = PP_ALIGN.LEFT
        
        # Text box 2: Global ISI auto scrolls
        # Size: Height 1.12cm, width 8.2cm, x = 0.98cm, y = 8.71cm
        textbox2 = slide.shapes.add_textbox(
            Inches(0.98 / 2.54),   # x position
            Inches(8.71 / 2.54),   # y position  
            Inches(8.2 / 2.54),    # width
            Inches(1.12 / 2.54)    # height
        )
        
        # Set red border (no background fill)
        textbox2.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        textbox2.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content with bold formatting
        text_content2 = "Global: The ISI auto scrolls, but the user also has the option to manually scroll through"
        textbox2.text_frame.text = text_content2
        
        # Format textbox2 with bold "Global:"
        text_frame2 = textbox2.text_frame
        text_frame2.margin_left = Inches(0.1)
        text_frame2.margin_right = Inches(0.1)
        text_frame2.margin_top = Inches(0.05)
        text_frame2.margin_bottom = Inches(0.05)
        text_frame2.word_wrap = True
        
        paragraph2 = text_frame2.paragraphs[0]
        paragraph2.font.name = "Aptos"
        paragraph2.font.size = Pt(10)
        paragraph2.font.color.rgb = RGBColor(0, 0, 0)
        paragraph2.alignment = PP_ALIGN.CENTER
        
        # Make "Global:" bold by splitting the text
        run1 = paragraph2.runs[0]
        run1.text = "Global: "
        run1.font.bold = True
        
        run2 = paragraph2.add_run()
        run2.text = "The ISI auto scrolls, but the user also has the option to manually scroll through"
        run2.font.bold = False
        
        # Text box 3: ROLLOVER TO EXPLORE and VDX.TV animate
        # Size: Height 1.09cm, width 8.48cm, x = 16.23cm, y = 8.98cm
        textbox3 = slide.shapes.add_textbox(
            Inches(16.23 / 2.54),  # x position
            Inches(8.98 / 2.54),   # y position  
            Inches(8.48 / 2.54),   # width
            Inches(1.09 / 2.54)    # height
        )
        
        # Set red border (no background fill)
        textbox3.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        textbox3.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        text_content3 = '"ROLLOVER TO EXPLORE" and "VDX.TV" animate in peel back every 3 sec.'
        textbox3.text_frame.text = text_content3
        
        # Format textbox3
        text_frame3 = textbox3.text_frame
        text_frame3.margin_left = Inches(0.1)
        text_frame3.margin_right = Inches(0.1)
        text_frame3.margin_top = Inches(0.05)
        text_frame3.margin_bottom = Inches(0.05)
        text_frame3.word_wrap = True
        
        paragraph3 = text_frame3.paragraphs[0]
        paragraph3.font.name = "Aptos"
        paragraph3.font.size = Pt(10)
        paragraph3.font.color.rgb = RGBColor(0, 0, 0)
        paragraph3.alignment = PP_ALIGN.CENTER
        
        # Make quoted text bold
        run3_1 = paragraph3.runs[0]
        run3_1.text = '"ROLLOVER TO EXPLORE" '
        run3_1.font.bold = True
        
        run3_2 = paragraph3.add_run()
        run3_2.text = 'and '
        run3_2.font.bold = False
        
        run3_3 = paragraph3.add_run()
        run3_3.text = '"VDX.TV" '
        run3_3.font.bold = True
        
        run3_4 = paragraph3.add_run()
        run3_4.text = 'animate in peel back every 3 sec.'
        run3_4.font.bold = False
        
        # Text box 4: Clicking the "X" minimizes
        # Size: Height 1.15cm, width 5.81cm, x = 25.15cm, y = 10.94cm
        textbox4 = slide.shapes.add_textbox(
            Inches(25.15 / 2.54),  # x position
            Inches(10.94 / 2.54),  # y position  
            Inches(5.81 / 2.54),   # width
            Inches(1.15 / 2.54)    # height
        )
        
        # Set red border (no background fill)
        textbox4.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        textbox4.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        text_content4 = 'Clicking the "X" minimizes the teaser to 970x90.'
        textbox4.text_frame.text = text_content4
        
        # Format textbox4
        text_frame4 = textbox4.text_frame
        text_frame4.margin_left = Inches(0.1)
        text_frame4.margin_right = Inches(0.1)
        text_frame4.margin_top = Inches(0.05)
        text_frame4.margin_bottom = Inches(0.05)
        text_frame4.word_wrap = True
        
        paragraph4 = text_frame4.paragraphs[0]
        paragraph4.font.name = "Aptos"
        paragraph4.font.size = Pt(10)
        paragraph4.font.color.rgb = RGBColor(0, 0, 0)
        paragraph4.alignment = PP_ALIGN.CENTER
        
        # Make "X" and "970x90" bold
        run4_1 = paragraph4.runs[0]
        run4_1.text = 'Clicking the '
        run4_1.font.bold = False
        
        run4_2 = paragraph4.add_run()
        run4_2.text = '"X" '
        run4_2.font.bold = True
        
        run4_3 = paragraph4.add_run()
        run4_3.text = 'minimizes the teaser to '
        run4_3.font.bold = False
        
        run4_4 = paragraph4.add_run()
        run4_4.text = '970x90'
        run4_4.font.bold = True
        
        run4_5 = paragraph4.add_run()
        run4_5.text = '.'
        run4_5.font.bold = False
        
        # Text box 5: Global volume icon functionality
        # Size: Height 1.17cm, width 6.67cm, x = 10.01cm, y = 17.5cm
        textbox5 = slide.shapes.add_textbox(
            Inches(10.01 / 2.54),  # x position
            Inches(17.5 / 2.54),   # y position
            Inches(6.67 / 2.54),   # width
            Inches(1.17 / 2.54)    # height
        )
        
        # Set red border (no background fill)
        textbox5.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        textbox5.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        text_frame5 = textbox5.text_frame
        text_frame5.margin_left = Inches(0.1)
        text_frame5.margin_right = Inches(0.1)
        text_frame5.margin_top = Inches(0.05)
        text_frame5.margin_bottom = Inches(0.05)
        text_frame5.word_wrap = True
        
        paragraph5 = text_frame5.paragraphs[0]
        paragraph5.font.name = "Aptos"
        paragraph5.font.size = Pt(10)
        paragraph5.font.color.rgb = RGBColor(0, 0, 0)
        paragraph5.alignment = PP_ALIGN.CENTER
        
        # Add "Global" as bold text
        run5_1 = paragraph5.runs[0] if paragraph5.runs else paragraph5.add_run()
        run5_1.text = 'Global'
        run5_1.font.bold = True
        
        # Add ": The volume icon can be used to mute or unmute the video sound." as regular text
        run5_2 = paragraph5.add_run()
        run5_2.text = ': The volume icon can be used to mute or unmute the video sound.'
        run5_2.font.bold = False
        
        # Text box 6: Global tab switching functionality  
        # Size: Height 1.11cm, width 6.67cm, x = 18.03cm, y = 17.5cm
        textbox6 = slide.shapes.add_textbox(
            Inches(18.03 / 2.54),  # x position
            Inches(17.5 / 2.54),   # y position
            Inches(6.67 / 2.54),   # width
            Inches(1.11 / 2.54)    # height
        )
        
        # Set red border (no background fill)
        textbox6.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        textbox6.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        text_frame6 = textbox6.text_frame
        text_frame6.margin_left = Inches(0.1)
        text_frame6.margin_right = Inches(0.1)
        text_frame6.margin_top = Inches(0.05)
        text_frame6.margin_bottom = Inches(0.05)
        text_frame6.word_wrap = True
        
        paragraph6 = text_frame6.paragraphs[0]
        paragraph6.font.name = "Aptos"
        paragraph6.font.size = Pt(10)
        paragraph6.font.color.rgb = RGBColor(0, 0, 0)
        paragraph6.alignment = PP_ALIGN.CENTER
        
        # Add "Global: " as bold text
        run6_1 = paragraph6.runs[0] if paragraph6.runs else paragraph6.add_run()
        run6_1.text = 'Global: '
        run6_1.font.bold = True
        
        # Add "The user can switch between tabs using these buttons." as regular text
        run6_2 = paragraph6.add_run()
        run6_2.text = 'The user can switch between tabs using these buttons.'
        run6_2.font.bold = False
        
        self.logger.info("Added all 6 annotation text boxes for Desktop In-frame 970x250 slide")
    
    def _add_engaged_state_textbox(self, slide):
        """Add Engaged State text box for Desktop In-frame 970x250 'With Annos' slides."""
        
        # Engaged State text box
        # Size: Height 3.21cm, width 8.11cm, x = 25.16cm, y = 12.45cm
        engaged_textbox = slide.shapes.add_textbox(
            Inches(25.16 / 2.54),  # x position
            Inches(12.45 / 2.54),  # y position  
            Inches(8.11 / 2.54),   # width
            Inches(3.21 / 2.54)    # height
        )
        
        # Set background color to #F8FFD1 and black border
        engaged_textbox.fill.solid()
        engaged_textbox.fill.fore_color.rgb = RGBColor(248, 255, 209)  # #F8FFD1 in RGB
        engaged_textbox.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        engaged_textbox.line.width = Pt(0.5)  # 0.5pt border width
        
        # Add text content with bold "Engaged State"
        text_content = "Engaged State\n\nOnce the users click or hover, the full video starts playing with sound. They will then see the complete unit, which includes the logo, text, CTA, and tabs. The ISI will automatically start scrolling from the beginning, but the user can also scroll it manually."
        engaged_textbox.text_frame.text = text_content
        
        # Format the text box
        text_frame = engaged_textbox.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True
        
        for i, paragraph in enumerate(text_frame.paragraphs):
            if i == 0:  # First paragraph - bold "Engaged State"
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(9)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.alignment = PP_ALIGN.LEFT
            else:  # Remaining paragraphs - regular
                paragraph.font.name = "Aptos"
                paragraph.font.size = Pt(9)
                paragraph.font.bold = False
                paragraph.font.color.rgb = RGBColor(0, 0, 0)
                paragraph.alignment = PP_ALIGN.LEFT
        
        self.logger.info("Added Engaged State text box for Desktop In-frame 970x250 slide")
    
    def _add_global_textbox(self, slide):
        """Add Global text box for Desktop In-frame 970x250 'With Annos' slides."""
        
        # Global text box
        # Size: Height 1.11cm, width 6.67cm, x = 18.03cm, y = 17.5cm
        global_textbox = slide.shapes.add_textbox(
            Inches(18.03 / 2.54),  # x position
            Inches(17.5 / 2.54),   # y position
            Inches(6.67 / 2.54),   # width
            Inches(1.11 / 2.54)    # height
        )
        
        # Set red border (no background fill)
        global_textbox.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        global_textbox.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        text_frame = global_textbox.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.word_wrap = True
        
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.alignment = PP_ALIGN.CENTER
        
        # Add "Global: " as bold text
        run1 = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        run1.text = 'Global: '
        run1.font.bold = True
        
        # Add "The user can switch between tabs using these buttons." as regular text
        run2 = paragraph.add_run()
        run2.text = 'The user can switch between tabs using these buttons.'
        run2.font.bold = False
        
        self.logger.info("Added Global text box for Desktop In-frame 970x250 slide")
    
    def _arrange_desktop_inframe_300x250_images_manual(self, slide, image_paths, annotation_option, folder_name):
        """Arrange Desktop In-frame 300x250 images sequentially for Manual tab with up to 6 images per slide."""
        # Calculate positioning for 6 images in 2 rows of 3 columns
        # Gray rectangle height is 1.79cm, so start Y position is below it with breathing space
        start_y_cm = 1.79 + 0.5  # 0.5cm breathing space below gray rectangle
        
        # Available slide dimensions (standard 16:9 slide is 33.87cm x 19.05cm)
        available_width_cm = 33.87 - 1.0  # 0.5cm margin on each side
        available_height_cm = 19.05 - start_y_cm - 0.5  # 0.5cm bottom margin
        
        # Calculate dimensions for 6 images in 3x2 grid
        columns = 3
        rows = 2
        spacing_cm = 0.5  # Space between images
        
        # Fixed dimensions per requirements
        image_width_cm = 9.48  # Fixed width 9.48cm
        image_height_cm = 7.9  # Fixed height 7.9cm
        
        # Calculate starting position to center the grid
        total_width_needed = columns * image_width_cm + (columns - 1) * spacing_cm
        slide_width_cm = 33.87  # Standard 16:9 slide width
        start_x_cm = (slide_width_cm - total_width_needed) / 2  # Center horizontally
        
        # Process up to 6 images per slide
        for i, img_path in enumerate(image_paths[:6]):
            row = i // columns
            col = i % columns
            
            x_cm = start_x_cm + col * (image_width_cm + spacing_cm)
            y_cm = start_y_cm + row * (image_height_cm + spacing_cm)
            
            width_inches = image_width_cm / 2.54
            height_inches = image_height_cm / 2.54
            x_inches = x_cm / 2.54
            y_inches = y_cm / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Add black border
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Manual Desktop In-frame 300x250: Added image {i+1} at position ({x_cm:.2f}, {y_cm:.2f})cm, size ({image_width_cm:.2f}x{image_height_cm:.2f})cm")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Desktop In-frame 300x250 image {img_path}: {str(e)}")
    
    def _arrange_desktop_inframe_300x600_images_manual(self, slide, image_paths, annotation_option, folder_name):
        """Arrange Desktop In-frame 300x600 images sequentially for Manual tab with up to 4 images per slide."""
        # Calculate positioning for 4 images in 1 row of 4 columns
        # Gray rectangle height is 1.79cm, so start Y position is below it with breathing space
        start_y_cm = 1.79 + 0.5  # 0.5cm breathing space below gray rectangle
        
        # Available slide dimensions (standard 16:9 slide is 33.87cm x 19.05cm)
        available_width_cm = 33.87 - 1.0  # 0.5cm margin on each side
        available_height_cm = 19.05 - start_y_cm - 0.5  # 0.5cm bottom margin
        
        # Calculate dimensions for 4 images in 1x4 grid
        columns = 4
        rows = 1
        spacing_cm = 0.5  # Space between images
        
        # Calculate image dimensions
        image_width_cm = (available_width_cm - (columns - 1) * spacing_cm) / columns
        image_height_cm = max(15.33, available_height_cm)  # Minimum 15.33cm height
        
        # Starting positions
        start_x_cm = 0.5  # Left margin
        
        # Process up to 4 images per slide
        for i, img_path in enumerate(image_paths[:4]):
            col = i % columns
            
            x_cm = start_x_cm + col * (image_width_cm + spacing_cm)
            y_cm = start_y_cm
            
            width_inches = image_width_cm / 2.54
            height_inches = image_height_cm / 2.54
            x_inches = x_cm / 2.54
            y_inches = y_cm / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Add black border
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Manual Desktop In-frame 300x600: Added image {i+1} at position ({x_cm:.2f}, {y_cm:.2f})cm, size ({image_width_cm:.2f}x{image_height_cm:.2f})cm")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Desktop In-frame 300x600 image {img_path}: {str(e)}")
    
    def _arrange_desktop_inframe_160x600_images_manual(self, slide, image_paths, annotation_option, folder_name):
        """Arrange Desktop In-frame 160x600 images sequentially for Manual tab with up to 7 images per slide."""
        # X positions for up to 7 images as specified by user
        x_positions_cm = [0.83, 5.58, 10.33, 15.01, 19.68, 24.36, 29.04]
        
        # Fixed dimensions and Y position for all images
        height_cm = 15.98
        width_cm = 4.26
        y_cm = 2.3
        
        # Process up to 7 images per slide
        for i, img_path in enumerate(image_paths[:7]):
            x_cm = x_positions_cm[i]
            
            width_inches = width_cm / 2.54
            height_inches = height_cm / 2.54
            x_inches = x_cm / 2.54
            y_inches = y_cm / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Add black border
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Manual Desktop In-frame 160x600: Added image {i+1} at position ({x_cm}, {y_cm})cm")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Desktop In-frame 160x600 image {img_path}: {str(e)}")
    
    def _arrange_desktop_inframe_728x90_images_manual(self, slide, image_paths, annotation_option, folder_name):
        """Arrange Desktop In-frame 728x90 images sequentially for Manual tab with up to 5 images stacked vertically."""
        # Calculate positioning for 5 images stacked vertically
        # Gray rectangle height is 1.79cm, so start Y position is below it with breathing space
        start_y_cm = 1.79 + 0.5  # 0.5cm breathing space below gray rectangle
        
        # Available slide dimensions (standard 16:9 slide is 33.87cm x 19.05cm)
        available_width_cm = 33.87 - 1.0  # 0.5cm margin on each side
        available_height_cm = 19.05 - start_y_cm - 0.5  # 0.5cm bottom margin
        
        # Calculate dimensions for 5 images stacked vertically
        rows = 5
        columns = 1
        spacing_cm = 0.3  # Space between images (smaller for vertical stack)
        
        # Calculate image dimensions with constraints
        calculated_height = (available_height_cm - (rows - 1) * spacing_cm) / rows
        image_height_cm = max(2.68, min(3.0, calculated_height))  # Minimum 2.68cm, maximum 3cm height
        image_width_cm = min(24.21, available_width_cm)  # Maximum 24.21cm width
        
        # Starting positions
        start_x_cm = 0.5  # Left margin
        
        # Process up to 5 images per slide
        for i, img_path in enumerate(image_paths[:5]):
            row = i
            
            x_cm = start_x_cm
            y_cm = start_y_cm + row * (image_height_cm + spacing_cm)
            
            width_inches = image_width_cm / 2.54
            height_inches = image_height_cm / 2.54
            x_inches = x_cm / 2.54
            y_inches = y_cm / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Add black border
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
                
                self.logger.info(f"Manual Desktop In-frame 728x90: Added image {i+1} at position ({x_cm:.2f}, {y_cm:.2f})cm, size ({image_width_cm:.2f}x{image_height_cm:.2f})cm")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Desktop In-frame 728x90 image {img_path}: {str(e)}")
    
    def _arrange_mobile_images_manual(self, slide, image_paths, annotation_option, folder_name, is_first_slide=True):
        """Arrange mobile images sequentially for Manual tab with up to 4 images - 1x1 grid layout."""
        # Calculate positioning for mobile images
        # Gray rectangle height is 1.79cm, so start Y position is below it with breathing space
        start_y_cm = 1.79 + 0.5  # 0.5cm breathing space below gray rectangle
        
        # Fixed dimensions per requirements
        image_width_cm = 8.06  # Maximum width 8.06cm
        image_height_cm = 15.67  # Fixed height 15.67cm
        spacing_cm = 0.5  # Space between images
        
        # Calculate starting X position based on number of images
        num_images = min(len(image_paths), 4)  # Process up to 4 images
        slide_width_cm = 33.87  # Standard 16:9 slide width
        
        if num_images == 1:
            # Single image - center horizontally
            start_x_cm = (slide_width_cm - image_width_cm) / 2
        else:
            # Multiple images - center the group horizontally
            total_width_needed = num_images * image_width_cm + (num_images - 1) * spacing_cm
            start_x_cm = (slide_width_cm - total_width_needed) / 2
        
        # Process up to 4 images per slide
        for i, img_path in enumerate(image_paths[:4]):
            if num_images == 1:
                x_cm = start_x_cm  # Single image centered
            else:
                x_cm = start_x_cm + i * (image_width_cm + spacing_cm)  # Multiple images spaced
            
            y_cm = start_y_cm
            
            width_inches = image_width_cm / 2.54
            height_inches = image_height_cm / 2.54
            x_inches = x_cm / 2.54
            y_inches = y_cm / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Mobile images have no borders (following existing mobile image rules)
                # No border for mobile images
                
                self.logger.info(f"Manual mobile: Added image {i+1} at position ({x_cm:.2f}, {y_cm:.2f})cm, size ({image_width_cm:.2f}x{image_height_cm:.2f})cm")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual mobile image {img_path}: {str(e)}")
        
        # Add descriptive text boxes for Mobile In-frame 300x250 first slide with "With Annos" option
        if is_first_slide and annotation_option == 'with_annos' and 'vdxmobileinframe' in folder_name.lower() and '300x250' in folder_name:
            self._add_mobile_inframe_300x250_textbox(slide)
            self._add_mobile_inframe_300x250_animation_textbox(slide)
    
    def _arrange_mobile_inframe_300x250_images_manual(self, slide, image_paths, annotation_option, folder_name):
        """Arrange Mobile In-frame 300x250 images sequentially for Manual tab."""
        positioning_specs = {
            'first': {'height_cm': 16.02, 'width_cm': 8.23, 'x_cm': 5.02, 'y_cm': 2},
            'second': {'height_cm': 16.02, 'width_cm': 8.23, 'x_cm': 18.7, 'y_cm': 2}
        }
        
        for i, img_path in enumerate(image_paths[:2]):
            position_key = 'first' if i == 0 else 'second'
            spec = positioning_specs[position_key]
            
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Mobile images don't get borders
                
                self.logger.info(f"Manual Mobile In-frame 300x250: Added image {i+1} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Mobile In-frame 300x250 image {img_path}: {str(e)}")
    
    def _arrange_mobile_inframe_300x600_images_manual(self, slide, image_paths, annotation_option, folder_name):
        """Arrange Mobile In-frame 300x600 images sequentially for Manual tab."""
        positioning_specs = {
            'first': {'height_cm': 16.02, 'width_cm': 8.23, 'x_cm': 5.02, 'y_cm': 2},
            'second': {'height_cm': 16.02, 'width_cm': 8.23, 'x_cm': 18.7, 'y_cm': 2}
        }
        
        for i, img_path in enumerate(image_paths[:2]):
            position_key = 'first' if i == 0 else 'second'
            spec = positioning_specs[position_key]
            
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Mobile images don't get borders
                
                self.logger.info(f"Manual Mobile In-frame 300x600: Added image {i+1} at position ({x_inches:.2f}, {y_inches:.2f})")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Mobile In-frame 300x600 image {img_path}: {str(e)}")
    
    def _arrange_mobile_instream_images_manual(self, slide, image_paths, annotation_option, folder_name, is_first_slide=True):
        """Arrange Mobile Instream images sequentially for Manual tab with cropping and specific dimensions."""
        # Fixed dimensions: height 11.28cm, width 12.76cm
        # Two images with specific positioning
        
        # Image specifications
        img_height_cm = 11.28
        img_width_cm = 12.76
        
        if is_first_slide:
            # Special positioning for first slide: x=6.13cm, y=4cm with 0.5cm gap
            start_x_cm = 6.13
            y_position_cm = 4.0
            horizontal_spacing = 0.5  # 0.5cm gap between images
        else:
            # Calculate spacing for centering two images for continuation slides
            slide_width_cm = 33.867  # Standard slide width in cm
            total_images_width = 2 * img_width_cm
            horizontal_spacing = 2.0  # cm between images
            total_width_with_spacing = total_images_width + horizontal_spacing
            start_x_cm = (slide_width_cm - total_width_with_spacing) / 2
            y_position_cm = 4.0  # Below gray rectangle with proper spacing
        
        positioning_specs = {
            'first': {'height_cm': img_height_cm, 'width_cm': img_width_cm, 'x_cm': start_x_cm, 'y_cm': y_position_cm},
            'second': {'height_cm': img_height_cm, 'width_cm': img_width_cm, 'x_cm': start_x_cm + img_width_cm + horizontal_spacing, 'y_cm': y_position_cm}
        }
        
        for i, img_path in enumerate(image_paths[:2]):
            position_key = 'first' if i == 0 else 'second'
            spec = positioning_specs[position_key]
            
            # Crop image from bottom to 774px height
            cropped_img_path = self.image_processor._crop_image_from_bottom(img_path, 774)
            
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    cropped_img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Mobile images don't get borders
                
                self.logger.info(f"Manual Mobile Instream: Added image {i+1} at position ({x_inches:.2f}, {y_inches:.2f}) with cropped height 774px")
                    
                # Clean up temporary cropped file if it was created
                if cropped_img_path != img_path:
                    try:
                        os.unlink(cropped_img_path)
                    except:
                        pass  # Ignore cleanup errors
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Mobile Instream image {img_path}: {str(e)}")
                # Clean up temporary file on error
                if cropped_img_path != img_path:
                    try:
                        os.unlink(cropped_img_path)
                    except:
                        pass
        
        # Add descriptive text box only for first slide and "With Annos" option
        if is_first_slide and annotation_option == 'with_annos':
            self._add_mobile_instream_textbox(slide)
    
    def _add_mobile_instream_textbox(self, slide):
        """Add Mobile Instream descriptive text box for first Mobile Instream slide with 'With Annos'."""
        
        # Mobile Instream text box
        # Size: Height 10.53cm, width 5.31cm, x = 0.8cm, y = 4.59cm
        mobile_textbox = slide.shapes.add_textbox(
            Inches(0.8 / 2.54),    # x position
            Inches(4.59 / 2.54),   # y position  
            Inches(5.31 / 2.54),   # width
            Inches(10.53 / 2.54)   # height
        )
        
        # Set background color to #F8FFD1 and black border
        mobile_textbox.fill.solid()
        mobile_textbox.fill.fore_color.rgb = RGBColor(248, 255, 209)  # #F8FFD1 in RGB
        mobile_textbox.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        mobile_textbox.line.width = Pt(0.5)  # 0.5pt border width
        
        # Add text content with bold title and bullet points
        text_content = "Mobile Instream/Inread units:\n\n• Mobile Instream units run on video playing sites like Dailymotion or Vimeo.\n• Instreams differ from other desktop In-frame/expandable units by not showing a pre-engagement state. They are displayed immediately before the main video begins to play for the user. The main video here refers to the actual video that user wished to play from the website. During this time, the user will get a message like \"Video Play Soon\" or option to \"Skip\" the ad.\n• To engage, user will tap the video. Upon engagement, the ISI will auto-scroll from the beginning and the video will continue playing.\n• If the user does not engage, the video will play in full."
        mobile_textbox.text_frame.text = text_content
        
        # Format the text box
        text_frame = mobile_textbox.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True
        
        # Format all paragraphs
        for i, paragraph in enumerate(text_frame.paragraphs):
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.alignment = PP_ALIGN.LEFT
            
            if i == 0:  # First paragraph - bold title
                paragraph.font.bold = True
            else:  # Bullet points - regular text
                paragraph.font.bold = False
        
        self.logger.info("Added Mobile Instream descriptive text box for first slide")
    
    def _arrange_ott_images_manual(self, slide, image_paths, annotation_option, folder_name, is_first_slide=True):
        """Arrange OTT images sequentially for Manual tab with standard grid layout and descriptive text box."""
        # Use standard image arrangement (no special positioning needed for OTT)
        self._arrange_images_on_slide(slide, image_paths, annotation_option, folder_name)
        
        # Add descriptive text box only for first slide and "With Annos" option
        if is_first_slide and annotation_option == 'with_annos':
            self._add_ott_textbox(slide)
    
    def _add_ott_textbox(self, slide):
        """Add OTT descriptive text box for first OTT slide with 'With Annos'."""
        
        # OTT text box
        # Size: Height 1.91cm, width 12.98cm, x = 2.87cm, y = 14.75cm
        ott_textbox = slide.shapes.add_textbox(
            Inches(2.87 / 2.54),   # x position
            Inches(14.75 / 2.54),  # y position  
            Inches(12.98 / 2.54),  # width
            Inches(1.91 / 2.54)    # height
        )
        
        # Set background color to #F8FFD1 and black border
        ott_textbox.fill.solid()
        ott_textbox.fill.fore_color.rgb = RGBColor(248, 255, 209)  # #F8FFD1 in RGB
        ott_textbox.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        ott_textbox.line.width = Pt(0.5)  # 0.5pt border width
        
        # Add text content with bold title and description
        text_content = "OTT Units:\nOTT ads run on streaming platforms like Hulu, and Roku, viewed on Smart TVs, Fire TV, or Apple TV. These are full-screen, non-interactive, non-skippable video ads shown before or during content, ensuring 100% viewability in a lean-back environment."
        ott_textbox.text_frame.text = text_content
        
        # Format the text box
        text_frame = ott_textbox.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True
        
        # Format all paragraphs
        for i, paragraph in enumerate(text_frame.paragraphs):
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.alignment = PP_ALIGN.LEFT
            
            if i == 0:  # First paragraph - bold title
                paragraph.font.bold = True
            else:  # Description - regular text
                paragraph.font.bold = False
        
        self.logger.info("Added OTT descriptive text box for first slide")
    
    def _add_mobile_inframe_300x250_textbox(self, slide):
        """Add Mobile In-frame 300x250 descriptive text box for first slide with 'With Annos'."""
        
        # Mobile In-frame text box
        # Size: Height 2.18cm, width 5.62cm, x = 9.85cm, y = 12.06cm
        mobile_inframe_textbox = slide.shapes.add_textbox(
            Inches(9.85 / 2.54),   # x position
            Inches(12.06 / 2.54),  # y position  
            Inches(5.62 / 2.54),   # width
            Inches(2.18 / 2.54)    # height
        )
        
        # Set background color to #F8FFD1 and black border
        mobile_inframe_textbox.fill.solid()
        mobile_inframe_textbox.fill.fore_color.rgb = RGBColor(248, 255, 209)  # #F8FFD1 in RGB
        mobile_inframe_textbox.line.color.rgb = RGBColor(0, 0, 0)  # Black border
        mobile_inframe_textbox.line.width = Pt(0.5)  # 0.5pt border width
        
        # Add text content with bold title and description
        text_content = "Mobile In-frame Teaser Sizes:\nThe mobile In-frame units include two teaser sizes 300x250 and 300x600 which are similar to Desktop In-frame units 300x250 and 300x600."
        mobile_inframe_textbox.text_frame.text = text_content
        
        # Format the text box
        text_frame = mobile_inframe_textbox.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True
        
        # Format all paragraphs
        for i, paragraph in enumerate(text_frame.paragraphs):
            paragraph.font.name = "Aptos"
            paragraph.font.size = Pt(9)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)
            paragraph.alignment = PP_ALIGN.LEFT
            
            if i == 0:  # First paragraph - bold title
                paragraph.font.bold = True
            else:  # Description - regular text
                paragraph.font.bold = False
        
        self.logger.info("Added Mobile In-frame 300x250 descriptive text box for first slide")
    
    def _add_mobile_inframe_300x250_animation_textbox(self, slide):
        """Add Mobile In-frame 300x250 animation text box for first slide with 'With Annos'."""
        
        # Mobile In-frame animation text box
        # Size: Height 1.88cm, width 5.55cm, x = 1.35cm, y = 13.83cm
        animation_textbox = slide.shapes.add_textbox(
            Inches(1.35 / 2.54),   # x position
            Inches(13.83 / 2.54),  # y position  
            Inches(5.55 / 2.54),   # width
            Inches(1.88 / 2.54)    # height
        )
        
        # Set background color to white and red border
        animation_textbox.fill.solid()
        animation_textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        animation_textbox.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        animation_textbox.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content with bold formatting for specific phrases
        text_content = 'The message "TAP TO EXPLORE" and "VDX.TV" animate one after another in this area every 3 sec.'
        animation_textbox.text_frame.text = text_content
        
        # Format the text box
        text_frame = animation_textbox.text_frame
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.1)
        text_frame.margin_bottom = Inches(0.1)
        text_frame.word_wrap = True
        
        # Format the paragraph
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.bold = False  # Default to non-bold
        
        # Apply bold formatting to specific phrases
        # Note: PowerPoint's text formatting is limited in python-pptx
        # The bold formatting for "TAP TO EXPLORE" and "VDX.TV" will be applied to the entire text
        # For selective bold formatting, we would need to use runs, but this is complex
        # For now, the entire text will be regular weight with the bold phrases indicated in quotes
        
        self.logger.info("Added Mobile In-frame 300x250 animation text box for first slide")
    
    def _add_full_isi_textboxes(self, slide):
        """Add two text boxes to FULL ISI slide with 'With Annos'."""
        
        # First text box - Full Prescribing Information
        # Size: width 10.27cm, height 1.8cm, x = 2.41cm, y = 15.52cm
        prescribing_textbox = slide.shapes.add_textbox(
            Inches(2.41 / 2.54),   # x position
            Inches(15.52 / 2.54),  # y position  
            Inches(10.27 / 2.54),  # width
            Inches(1.8 / 2.54)     # height
        )
        
        # Set background color to white and red border
        prescribing_textbox.fill.solid()
        prescribing_textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        prescribing_textbox.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        prescribing_textbox.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        prescribing_text = 'The Full Prescribing Information clicks to:'
        prescribing_textbox.text_frame.text = prescribing_text
        
        # Format the text box
        prescribing_frame = prescribing_textbox.text_frame
        prescribing_frame.margin_left = Inches(0.1)
        prescribing_frame.margin_right = Inches(0.1)
        prescribing_frame.margin_top = Inches(0.1)
        prescribing_frame.margin_bottom = Inches(0.1)
        prescribing_frame.word_wrap = True
        
        # Format the paragraph
        prescribing_paragraph = prescribing_frame.paragraphs[0]
        prescribing_paragraph.font.name = "Aptos"
        prescribing_paragraph.font.size = Pt(10)
        prescribing_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        prescribing_paragraph.alignment = PP_ALIGN.LEFT
        prescribing_paragraph.font.bold = False
        
        # Second text box - Medication Guide
        # Size: width 10.27cm, height 1.8cm, x = 22.04cm, y = 15.53cm
        medication_textbox = slide.shapes.add_textbox(
            Inches(22.04 / 2.54),  # x position
            Inches(15.53 / 2.54),  # y position  
            Inches(10.27 / 2.54),  # width
            Inches(1.8 / 2.54)     # height
        )
        
        # Set background color to white and red border
        medication_textbox.fill.solid()
        medication_textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
        medication_textbox.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        medication_textbox.line.width = Pt(0.75)  # 0.75pt border width
        
        # Add text content
        medication_text = 'The Medication Guide clicks to:'
        medication_textbox.text_frame.text = medication_text
        
        # Format the text box
        medication_frame = medication_textbox.text_frame
        medication_frame.margin_left = Inches(0.1)
        medication_frame.margin_right = Inches(0.1)
        medication_frame.margin_top = Inches(0.1)
        medication_frame.margin_bottom = Inches(0.1)
        medication_frame.word_wrap = True
        
        # Format the paragraph
        medication_paragraph = medication_frame.paragraphs[0]
        medication_paragraph.font.name = "Aptos"
        medication_paragraph.font.size = Pt(10)
        medication_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        medication_paragraph.alignment = PP_ALIGN.LEFT
        medication_paragraph.font.bold = False
        
        self.logger.info("Added FULL ISI text boxes for 'With Annos' option")
    
    def _add_desktop_expandable_slide_manual(self, prs, folder_structure, annotation_option):
        """Add Desktop Expandable slide for Manual tab with images from vdxdesktopexpandable folder."""
        
        # Collect all images from vdxdesktopexpandable folder and subfolders
        desktop_expandable_images = []
        
        for folder_name, image_paths in folder_structure.items():
            if 'vdxdesktopexpandable' in folder_name.lower():
                # Filter out disclaimer images and check dimensions
                filtered_images = []
                for img_path in image_paths:
                    if (os.path.basename(img_path).lower() == 'teaser-disclaimer.png' or 
                        os.path.basename(img_path).lower() == 'mainunit-disclaimer.png' or
                        'engaged' in img_path.lower()):
                        continue
                    
                    # Check image dimensions - skip if smaller than 1900x1092
                    try:
                        from PIL import Image
                        with Image.open(img_path) as img:
                            width, height = img.size
                            if width >= 1900 and height >= 1092:
                                filtered_images.append(img_path)
                                self.logger.info(f"Desktop Expandable: Including image {os.path.basename(img_path)} ({width}x{height})")
                            else:
                                self.logger.info(f"Desktop Expandable: Skipping image {os.path.basename(img_path)} ({width}x{height}) - too small (requires ≥1900x1092)")
                    except Exception as e:
                        self.logger.warning(f"Desktop Expandable: Could not read dimensions for {os.path.basename(img_path)}: {str(e)}")
                        # Skip images that can't be read
                        continue
                
                desktop_expandable_images.extend(filtered_images)
        
        if not desktop_expandable_images:
            self.logger.info("No vdxdesktopexpandable images found - skipping Desktop Expandable slide")
            return
        
        # Use Desktop Expandable specific priority: special handling for 3 images (teaser, mainunit, vmp)
        desktop_expandable_images.sort(key=lambda img: self._sort_images_desktop_expandable_priority(img, desktop_expandable_images))
        
        self.logger.info(f"Manual tab: Processing {len(desktop_expandable_images)} images for Desktop Expandable slide")
        
        # Process images in chunks of 2 per slide
        images_per_slide = 2
        for i in range(0, len(desktop_expandable_images), images_per_slide):
            slide_images = desktop_expandable_images[i:i + images_per_slide]
            slide_number = (i // images_per_slide) + 1
            
            # Create slide
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Remove placeholder shapes
            self._remove_placeholders(slide)
            
            # Add title background and text
            title_text = "DESKTOP EXPANDABLE"
            if slide_number > 1:
                title_text += f" ({slide_number})"
            
            # Add gray rectangle background
            rectangle = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(0),
                Inches(13.33),
                Inches(1.79 / 2.54)
            )
            rectangle.fill.solid()
            rectangle.fill.fore_color.rgb = RGBColor(242, 242, 242)
            rectangle.line.fill.background()
            rectangle.shadow.inherit = False  # Remove shadow
            
            # Add title text
            title_text_box = slide.shapes.add_textbox(
                Inches(0.51 / 2.54),
                Inches(0.38 / 2.54),
                Inches(12 / 2.54),
                Inches(1 / 2.54)
            )
            title_text_frame = title_text_box.text_frame
            title_text_frame.text = title_text.upper()
            title_paragraph = title_text_frame.paragraphs[0]
            title_paragraph.font.name = "Aptos Display"
            title_paragraph.font.size = Pt(18)
            title_paragraph.font.bold = True
            title_paragraph.alignment = PP_ALIGN.LEFT
            title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
            
            # Add images with specific positioning
            self._arrange_desktop_expandable_images_manual(slide, slide_images, annotation_option)
            
            # Add VDX TV logo
            self._add_vdx_logo(slide, 'desktop_expandable')
            
            self.logger.info(f"Manual tab: Created Desktop Expandable slide {slide_number} with {len(slide_images)} images")
    
    def _arrange_desktop_expandable_images_manual(self, slide, image_paths, annotation_option):
        """Arrange Desktop Expandable images with specific positioning for Manual tab."""
        
        # User requirements: height 8.64cm, width 16.02cm, y = 4.69cm
        # Two images per slide with proper positioning
        positioning_specs = {
            'first': {'height_cm': 8.64, 'width_cm': 16.02, 'x_cm': 1, 'y_cm': 4.69},
            'second': {'height_cm': 8.64, 'width_cm': 16.02, 'x_cm': 17.45, 'y_cm': 4.69}
        }
        
        for i, img_path in enumerate(image_paths[:2]):
            position_key = 'first' if i == 0 else 'second'
            spec = positioning_specs[position_key]
            
            width_inches = spec['width_cm'] / 2.54
            height_inches = spec['height_cm'] / 2.54
            x_inches = spec['x_cm'] / 2.54
            y_inches = spec['y_cm'] / 2.54
            
            try:
                picture_shape = slide.shapes.add_picture(
                    img_path, Inches(x_inches), Inches(y_inches), Inches(width_inches), Inches(height_inches)
                )
                
                # Desktop Expandable images get black borders
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)  # Black color
                picture_shape.line.width = Pt(0.5)  # 0.5pt width
                
                self.logger.info(f"Manual Desktop Expandable: Added image {i+1} at position ({x_inches:.2f}, {y_inches:.2f})in, size ({width_inches:.2f}x{height_inches:.2f})in")
                    
            except Exception as e:
                self.logger.error(f"Error adding manual Desktop Expandable image {img_path}: {str(e)}")
    
    def _add_full_isi_slide_manual(self, prs, annotation_option):
        """Add blank FULL ISI slide as the last slide for Manual tab."""
        
        # Create slide
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove placeholder shapes
        self._remove_placeholders(slide)
        
        # Add gray rectangle background
        rectangle = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.33),
            Inches(1.79 / 2.54)
        )
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = RGBColor(242, 242, 242)
        rectangle.line.fill.background()
        rectangle.shadow.inherit = False  # Remove shadow
        
        # Add title text
        title_text_box = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),
            Inches(0.38 / 2.54),
            Inches(12 / 2.54),
            Inches(1 / 2.54)
        )
        title_text_frame = title_text_box.text_frame
        title_text_frame.text = "FULL ISI"
        title_paragraph = title_text_frame.paragraphs[0]
        title_paragraph.font.name = "Aptos Display"
        title_paragraph.font.size = Pt(18)
        title_paragraph.font.bold = True
        title_paragraph.alignment = PP_ALIGN.LEFT
        title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add VDX TV logo
        self._add_vdx_logo(slide, 'full_isi')
        
        # Add FULL ISI text boxes with "With Annos" option
        if annotation_option == 'with_annos':
            self._add_full_isi_textboxes(slide)
        
        self.logger.info("Manual tab: Created blank FULL ISI slide as last slide")