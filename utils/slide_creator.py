"""
Slide creation and layout management for PowerPoint presentations.
Handles different slide types and their specific arrangements.
"""

import os
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class SlideCreator:
    def __init__(self, base_generator, image_processor):
        self.base_generator = base_generator
        self.image_processor = image_processor
        self.logger = base_generator.logger
    
    def _create_slide_with_title(self, prs, title_text):
        """Create a new slide with gray title background."""
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Remove all placeholder shapes to prevent "Title 1" text
        self.base_generator._remove_placeholders(slide)
        
        # Add gray rectangle background for title (1.79cm height)
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), 
            prs.slide_width, Inches(1.79 / 2.54)  # Convert cm to inches
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = RGBColor(242, 242, 242)  # #F2F2F2 color
        title_bg.line.fill.background()
        title_bg.shadow.inherit = False  # Remove shadow
        
        # Add title text box with exact specifications
        title_textbox = slide.shapes.add_textbox(
            Inches(0.51 / 2.54),  # Convert cm to inches
            Inches(0.38 / 2.54),  # Convert cm to inches
            Inches(12),           # 12cm width
            Inches(1)             # 1cm height
        )
        text_frame = title_textbox.text_frame
        text_frame.text = title_text
        text_frame.margin_left = Inches(0)
        text_frame.margin_right = Inches(0)
        text_frame.margin_top = Inches(0)
        text_frame.margin_bottom = Inches(0)
        text_frame.word_wrap = True
        
        # Format the title text
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = 'Aptos Display'
        paragraph.font.size = Pt(18)
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
        
        return slide
    
    def _add_disclaimer_images(self, slide, disclaimer_files):
        """Add disclaimer images to a slide with fixed height of 14cm."""
        if not disclaimer_files:
            return
        
        # Fixed height for disclaimer images
        fixed_height = Inches(14 / 2.54)  # Convert 14cm to inches
        y_position = Inches(2 / 2.54)  # Start at y=2cm as requested
        
        # Calculate total width needed and starting position
        total_width = 0
        image_widths = []
        
        for img_path in disclaimer_files:
            try:
                width, height, aspect_ratio = self.image_processor._calculate_image_dimensions(
                    img_path, target_height=fixed_height
                )
                image_widths.append(width)
                total_width += width
            except Exception as e:
                self.logger.error(f"Error calculating width for disclaimer image {img_path}: {str(e)}")
                image_widths.append(Inches(4))  # Default width
                total_width += Inches(4)
        
        # Add spacing between images
        spacing = Inches(0.5)
        total_width += spacing * (len(disclaimer_files) - 1)
        
        # Start at x=0.51cm as requested
        start_x = Inches(0.51 / 2.54)
        
        current_x = start_x
        for i, img_path in enumerate(disclaimer_files):
            width = image_widths[i]
            
            # Add image without border (no borders as requested)
            self.image_processor._add_image_to_slide(
                slide, img_path, current_x, y_position, width, fixed_height, "disclaimer_no_border"
            )
            
            current_x += width + spacing
        
        self.logger.info(f"Added {len(disclaimer_files)} disclaimer images to FULL ISI slide")
    
    def _check_if_images_need_splitting(self, disclaimer_files):
        """Check if any disclaimer images have height > 1000px and need splitting."""
        if not disclaimer_files:
            return False
            
        from PIL import Image
        
        for img_path in disclaimer_files:
            try:
                with Image.open(img_path) as img:
                    width, height = img.size
                    if height > 1000:
                        self.logger.info(f"Image {img_path} has height {height}px > 1000px, will need splitting")
                        return True
            except Exception as e:
                self.logger.error(f"Error checking image dimensions for {img_path}: {str(e)}")
                continue
        
        return False
    
    def _create_split_full_isi_slides(self, prs, disclaimer_files):
        """Create multiple FULL ISI slides for tall disclaimer images."""
        from PIL import Image
        import tempfile
        import os
        
        slides_created = 0
        
        for img_path in disclaimer_files:
            try:
                with Image.open(img_path) as img:
                    width, height = img.size
                    
                    if height <= 1000:
                        # Regular height image - add to single slide
                        slide = self._create_slide_with_title(prs, "FULL ISI")
                        self._add_disclaimer_images(slide, [img_path])
                        self.image_processor._add_vdx_logo(slide)
                        slides_created += 1
                        self.logger.info(f"Added regular height image {img_path} to FULL ISI slide")
                    else:
                        # Tall image - split into parts based on 1000px height
                        max_height_per_part = 1000
                        parts_needed = (height + max_height_per_part - 1) // max_height_per_part  # Ceiling division
                        
                        self.logger.info(f"Splitting image {img_path} ({width}x{height}) into {parts_needed} parts of max {max_height_per_part}px each")
                        
                        temp_files = []
                        
                        for part_num in range(parts_needed):
                            # Calculate crop coordinates
                            top = part_num * max_height_per_part
                            bottom = min(top + max_height_per_part, height)
                            actual_part_height = bottom - top
                            
                            self.logger.info(f"Part {part_num + 1}: cropping from y={top} to y={bottom} (height={actual_part_height}px)")
                            
                            # Crop the image part
                            cropped_img = img.crop((0, top, width, bottom))
                            
                            # Save to temporary file
                            temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                            cropped_img.save(temp_file.name, format='PNG')
                            temp_file.close()
                            temp_files.append(temp_file.name)
                            
                            # Create slide for this part
                            slide_title = "FULL ISI" if part_num == 0 else f"FULL ISI (CONTD.)"
                            slide = self._create_slide_with_title(prs, slide_title)
                            self._add_disclaimer_images(slide, [temp_file.name])
                            self.image_processor._add_vdx_logo(slide)
                            slides_created += 1
                            
                            self.logger.info(f"Created FULL ISI slide part {part_num + 1}/{parts_needed} for {os.path.basename(img_path)} with {actual_part_height}px height")
                        
                        # Clean up temporary files
                        for temp_file in temp_files:
                            try:
                                os.unlink(temp_file)
                            except:
                                pass
                                
            except Exception as e:
                self.logger.error(f"Error processing disclaimer image {img_path}: {str(e)}")
                # Fallback - create regular slide
                slide = self._create_slide_with_title(prs, "FULL ISI")
                self._add_disclaimer_images(slide, [img_path])
                self.image_processor._add_vdx_logo(slide)
                slides_created += 1
        
        self.logger.info(f"Created {slides_created} FULL ISI slides (including split parts)")
        return slides_created
    
    def _create_full_isi_slide(self, prs, disclaimer_files):
        """Create FULL ISI slide with disclaimer images, splitting tall images if needed."""
        # Check if any images need to be split due to height > 1000px
        split_needed = self._check_if_images_need_splitting(disclaimer_files)
        
        if split_needed:
            slides_created = self._create_split_full_isi_slides(prs, disclaimer_files)
            self.logger.info(f"Created {slides_created} FULL ISI slides due to image splitting")
            # Return the count of slides created for proper tracking
            return slides_created
        else:
            slide = self._create_slide_with_title(prs, "FULL ISI")
            
            # Add disclaimer images if any
            if disclaimer_files:
                self._add_disclaimer_images(slide, disclaimer_files)
            
            # Add VDX TV logo
            self.image_processor._add_vdx_logo(slide)
            
            return slide
    
    def _create_blank_full_isi_slide(self, prs):
        """Create blank FULL ISI slide for Manual tab."""
        slide = self._create_slide_with_title(prs, "FULL ISI")
        
        # Add VDX TV logo
        self.image_processor._add_vdx_logo(slide)
        
        self.logger.info("Manual tab: Created blank FULL ISI slide as last slide")
        return slide
    
    def _arrange_images_in_grid(self, slide, image_paths, folder_name, annotation_option, max_per_slide=9):
        """Arrange images in a standard grid layout."""
        if not image_paths:
            return
        
        # Limit to max_per_slide images
        images_to_show = image_paths[:max_per_slide]
        
        # Calculate grid dimensions
        num_images = len(images_to_show)
        if num_images <= 3:
            cols = num_images
            rows = 1
        elif num_images <= 6:
            cols = 3
            rows = 2
        else:
            cols = 3
            rows = 3
        
        # Available space (leave space for title and margins)
        start_y = Inches(2.5)  # Below title
        available_width = Inches(12)
        available_height = Inches(5)
        
        # Calculate image size and spacing
        image_width = available_width / cols * 0.8  # 80% of available space
        image_height = available_height / rows * 0.8
        
        spacing_x = (available_width - (image_width * cols)) / (cols + 1)
        spacing_y = (available_height - (image_height * rows)) / (rows + 1)
        
        # Place images
        for i, img_path in enumerate(images_to_show):
            row = i // cols
            col = i % cols
            
            x = spacing_x + col * (image_width + spacing_x)
            y = start_y + spacing_y + row * (image_height + spacing_y)
            
            self.image_processor._add_image_to_slide(
                slide, img_path, x, y, image_width, image_height, folder_name
            )
            
            # Add annotation if needed
            if annotation_option == 'with_annos':
                self.image_processor._add_image_annotation(slide, img_path, x, y, image_width, image_height)
    
    def _create_consolidated_slide(self, prs, title, teaser_images, arrangement_method):
        """Create a consolidated slide with custom arrangement."""
        slide = self._create_slide_with_title(prs, title)
        
        # Use the specified arrangement method
        if arrangement_method and teaser_images:
            arrangement_method(slide, teaser_images)
        
        # Add VDX TV logo
        self.image_processor._add_vdx_logo(slide)
        
        return slide