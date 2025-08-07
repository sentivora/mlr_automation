"""
Image processing and positioning logic for PowerPoint presentations.
Handles image placement, sizing, and arrangement on slides.
"""

import os
import tempfile
from PIL import Image
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt


class ImageProcessor:
    def __init__(self, logger):
        self.logger = logger

    def _add_image_to_slide(self, slide, img_path, x, y, width, height, folder_name):
        """Add an image to a slide with proper formatting."""
        try:
            picture_shape = slide.shapes.add_picture(img_path, x, y, width, height)

            # Add borders based on folder type
            if folder_name == "disclaimer_no_border":
                # No border for disclaimer images in FULL ISI slide
                pass
            elif "disclaimer" in folder_name.lower():
                # Other disclaimer images get black borders
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
            elif any(pattern in folder_name.lower() for pattern in ['desktop', 'expandable', 'inframe']):
                # Desktop images get black borders
                picture_shape.line.color.rgb = RGBColor(0, 0, 0)
                picture_shape.line.width = Pt(0.5)
            # Mobile images don't get borders by default

            self.logger.info(f"Added image from {folder_name} at ({x}, {y}) with size ({width}, {height})")

        except Exception as e:
            self.logger.error(f"Error adding image {img_path}: {str(e)}")

    def _add_vdx_logo(self, slide, folder_name=""):
        """Add VDX TV logo to slide at specified position."""
        try:
            logo_path = os.path.join('static', 'vdx-tv-logo.png')
            if os.path.exists(logo_path):
                # Exact positioning: X = 31.42cm, Y = 0.63cm
                left = Inches(31.42 / 2.54)  # Convert cm to inches
                top = Inches(0.63 / 2.54)    # Convert cm to inches
                width = Inches(1.85 / 2.54)  # 1.85cm width
                height = Inches(0.51 / 2.54) # 0.51cm height

                slide.shapes.add_picture(logo_path, left, top, width, height)
                self.logger.info(f"Added VDX TV logo to slide at position ({left:.2f}, {top:.2f})")
            else:
                self.logger.warning(f"VDX TV logo not found at {logo_path}")
        except Exception as e:
            self.logger.error(f"Error adding VDX TV logo: {str(e)}")

    def _crop_image_from_bottom(self, image_path, target_height_px=774):
        """Crop image from bottom to specified height and save to temporary file."""
        try:
            with Image.open(image_path) as img:
                original_width, original_height = img.size

                if original_height <= target_height_px:
                    # No cropping needed, return original path
                    return image_path

                # Crop from bottom: keep top portion
                cropped_img = img.crop((0, 0, original_width, target_height_px))

                # Save to temporary file
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
                cropped_img.save(temp_file.name, format='PNG')
                temp_file.close()

                self.logger.info(f"Cropped image {image_path} from {original_height}px to {target_height_px}px height")
                return temp_file.name

        except Exception as e:
            self.logger.error(f"Error cropping image {image_path}: {str(e)}")
            return image_path  # Return original if cropping fails

    def _calculate_image_dimensions(self, image_path, target_width=None, target_height=None, max_width=None, max_height=None):
        """Calculate appropriate dimensions for an image while maintaining aspect ratio."""
        try:
            with Image.open(image_path) as img:
                original_width, original_height = img.size
                aspect_ratio = original_width / original_height

                if target_height and not target_width:
                    # Fixed height, calculate width
                    height = target_height
                    width = height * aspect_ratio
                    if max_width and width > max_width:
                        width = max_width
                        height = width / aspect_ratio
                elif target_width and not target_height:
                    # Fixed width, calculate height
                    width = target_width
                    height = width / aspect_ratio
                    if max_height and height > max_height:
                        height = max_height
                        width = height * aspect_ratio
                else:
                    # Use provided dimensions or defaults
                    width = target_width or original_width
                    height = target_height or original_height

                return width, height, aspect_ratio
        except Exception as e:
            self.logger.error(f"Error calculating dimensions for {image_path}: {str(e)}")
            return target_width or Inches(2), target_height or Inches(2), 1.0

    def _resize_image_if_needed(self, image_path, max_width_px=1920, max_height_px=1080):
        """Resize image if it's too large, return path to processed image."""
        try:
            with Image.open(image_path) as img:
                width, height = img.size

                # Check if resize is needed
                if width <= max_width_px and height <= max_height_px:
                    return image_path

                # Calculate new size maintaining aspect ratio
                ratio = min(max_width_px / width, max_height_px / height)
                new_width = int(width * ratio)
                new_height = int(height * ratio)

                # Create temporary resized image
                resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)

                # Save to temporary file
                temp_fd, temp_path = tempfile.mkstemp(suffix='.png')
                os.close(temp_fd)
                resized_img.save(temp_path, 'PNG')

                self.logger.info(f"Resized image from {width}x{height} to {new_width}x{new_height}")
                return temp_path

        except Exception as e:
            self.logger.error(f"Error resizing image {image_path}: {str(e)}")
            return image_path

    def _add_image_annotation(self, slide, image_path, left, top, width, height):
        """Add filename annotation below an image."""
        try:
            filename = os.path.basename(image_path)

            # Add text box below the image
            text_left = left
            text_top = top + height + Inches(0.1)  # Small gap below image
            text_width = width
            text_height = Inches(0.3)

            textbox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            text_frame = textbox.text_frame
            text_frame.text = filename

            # Format the text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(8)
            paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

        except Exception as e:
            self.logger.error(f"Error adding annotation for {image_path}: {str(e)}")