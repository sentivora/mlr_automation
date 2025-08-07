import os
import tempfile
import logging
import shutil
from datetime import datetime
from flask import Flask, render_template, request, redirect, url_for, send_file, jsonify, flash
from werkzeug.utils import secure_filename
from werkzeug.middleware.proxy_fix import ProxyFix
import zipfile
from zipfile import ZipFile
from utils.presentation_generator import PresentationGenerator
from utils.image_processor import ImageProcessor
from utils.presentation_assembler import PresentationAssembler
from utils.slide_creator import SlideCreator
from utils.base_generator import BaseGenerator
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.utils import ImageReader
from PIL import Image
# Removed Windows-specific imports for serverless compatibility
# import pythoncom
# import win32com.client
import subprocess


def is_image_file(filename):
    """Check if file is an image."""
    image_extensions = {'png', 'jpg', 'jpeg', 'gif', 'bmp'}
    return '.' in filename and filename.rsplit(
        '.', 1)[1].lower() in image_extensions


def extract_folder_structure(upload_path):
    folder_structure = {}
    for root, dirs, files in os.walk(upload_path):
        rel_path = os.path.relpath(root, upload_path)
        if rel_path == '.':
            folder_name = 'Images'
        else:
            folder_name = rel_path.replace("\\", "/")  # Use full relative path
        image_files = []
        for file in files:
            if is_image_file(file):
                full_path = os.path.join(root, file)
                try:
                    with Image.open(full_path) as img:
                        img.verify()
                    image_files.append(full_path)
                except Exception as e:
                    print(
                        f"Skipping invalid/corrupt image: {full_path} | Error: {e}"
                    )
        if image_files:
            folder_structure[folder_name] = image_files
    total = sum(len(v) for v in folder_structure.values())
    print(f"Total images found after extraction and verification: {total}")
    return folder_structure


# Configure logging for serverless environment
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Create Flask app
app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET",
                                "dev-secret-key-change-in-production")
app.wsgi_app = ProxyFix(app.wsgi_app, x_proto=1, x_host=1)

# Configuration
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
MAX_CONTENT_LENGTH = 500 * 1024 * 1024  # 500MB max file size
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'zip'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


def allowed_file(filename):
    """Check if file has allowed extension."""
    return '.' in filename and filename.rsplit(
        '.', 1)[1].lower() in ALLOWED_EXTENSIONS


@app.route('/')
def index():
    """Main page with upload form."""
    try:
        logger.info("Accessing main page")
        return render_template('index.html')
    except Exception as e:
        logger.error(f"Error in index route: {str(e)}")
        return f"Error: {str(e)}", 500


def process_uploaded_file(file, annotation_option='with_annos', is_multi_tab=False, implement_video_frames=False):
    """Process uploaded file and return results. Common logic for both tabs."""
    if not file or file.filename == '':
        return {'error': 'No file selected'}

    if not allowed_file(file.filename):
        return {'error': 'Invalid file type. Please upload a ZIP file or image.'}

    try:
        # Create temporary directory for this upload
        temp_dir = tempfile.mkdtemp()
        filename = secure_filename(file.filename)
        file_path = os.path.join(temp_dir, filename)
        file.save(file_path)
        
        # Extract original filename without extension for PowerPoint naming
        original_filename = None
        if filename.lower().endswith('.zip'):
            original_filename = filename[:-4]  # Remove .zip extension

        # Extract folder structure
        if filename.lower().endswith('.zip'):
            # Extract zip file
            extract_dir = os.path.join(temp_dir, 'extracted')
            os.makedirs(extract_dir, exist_ok=True)

            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                zip_ref.extractall(extract_dir)

            folder_structure = extract_folder_structure(extract_dir)
        else:
            # Single file upload
            if is_image_file(filename):
                folder_structure = {'Images': [file_path]}
            else:
                return {'error': 'Unsupported file type'}

        if not folder_structure:
            return {'error': 'No images found in uploaded content'}

        # Generate presentation
        app.logger.info(f"Creating presentation with parameters: annotation_option={annotation_option}, is_multi_tab={is_multi_tab}, implement_video_frames={implement_video_frames}")
        app.logger.info(f"Folder structure keys: {list(folder_structure.keys())}")
        
        assembler = PresentationAssembler()
        ppt_path, slide_count, video_folder_found = assembler.create_presentation(folder_structure,
                                                              OUTPUT_FOLDER,
                                                              annotation_option,
                                                              is_multi_tab,
                                                              implement_video_frames,
                                                              None,  # video_position_params
                                                              original_filename)

        # Clean up temporary directory
        shutil.rmtree(temp_dir)

        # Get relative paths for download
        ppt_filename = os.path.basename(ppt_path)

        return {
            'success': True,
            'ppt_file': ppt_filename,
            'folder_count': len(folder_structure),
            'slide_count': slide_count,
            'video_folder_found': video_folder_found
        }

    except Exception as e:
        app.logger.error(f"Error processing upload: {str(e)}")
        return {'error': f'Error processing file: {str(e)}'}


@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file upload and process images for Auto tab."""
    if 'file' not in request.files:
        flash('No file selected', 'error')
        return redirect(request.url)

    file = request.files['file']
    annotation_option = request.form.get('annotation_option', 'with_annos')
    implement_video_frames = 'implementVideoFrames' in request.form
    
    app.logger.info(f"Auto upload: annotation_option={annotation_option}, implement_video_frames={implement_video_frames}")

    result = process_uploaded_file(file, annotation_option, is_multi_tab=True, implement_video_frames=implement_video_frames)
    
    if 'error' in result:
        flash(result['error'], 'error')
        return redirect(request.url)

    return render_template('result.html',
                           ppt_file=result['ppt_file'],
                           folder_count=result['folder_count'],
                           slide_count=result['slide_count'],
                           video_folder_found=result['video_folder_found'],
                           implement_video_frames=implement_video_frames,
                           annotation_option=annotation_option)


@app.route('/manual-upload', methods=['POST'])
def manual_upload():
    """Handle file upload and process images for Manual tab."""
    if 'file' not in request.files:
        flash('No file selected', 'error')
        return redirect(request.url)

    file = request.files['file']
    annotation_option = request.form.get('annotation_option', 'with_annos')
    implement_video_frames = 'implementVideoFrames' in request.form
    
    app.logger.info(f"Manual upload: annotation_option={annotation_option}, implement_video_frames={implement_video_frames}")

    result = process_uploaded_file(file, annotation_option, is_multi_tab=False, implement_video_frames=implement_video_frames)
    
    if 'error' in result:
        flash(result['error'], 'error')
        return redirect(request.url)

    return render_template('result.html',
                           ppt_file=result['ppt_file'],
                           folder_count=result['folder_count'],
                           slide_count=result['slide_count'],
                           video_folder_found=result['video_folder_found'],
                           implement_video_frames=implement_video_frames,
                           annotation_option=annotation_option)


@app.route('/download/<filename>')
def download_file(filename):
    """Download generated files."""
    try:
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        if os.path.exists(file_path):
            return send_file(file_path, as_attachment=True)
        else:
            flash('File not found.', 'error')
            return redirect(url_for('index'))
    except Exception as e:
        flash(f'Error downloading file: {str(e)}', 'error')
        return redirect(url_for('index'))


@app.route('/convert-to-pdf/<filename>', methods=['GET', 'POST'])
def convert_to_pdf(filename):
    """Convert PowerPoint to PDF by converting slides to images first."""
    try:
        ppt_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        if not os.path.exists(ppt_path):
            return jsonify({'error': 'PowerPoint file not found'}), 404
        
        # Generate PDF filename
        pdf_filename = filename.replace('.pptx', '.pdf')
        pdf_path = os.path.join(app.config['OUTPUT_FOLDER'], pdf_filename)
        
        # Convert slides to images and create PDF
        success = convert_ppt_slides_to_pdf(ppt_path, pdf_path)
        
        if success:
            return send_file(pdf_path, as_attachment=True, download_name=pdf_filename)
        else:
            return jsonify({'error': 'Failed to convert PowerPoint to PDF'}), 500
            
    except Exception as e:
        logging.error(f"Error converting to PDF: {str(e)}")
        return jsonify({'error': f'Error converting to PDF: {str(e)}'}), 500


def convert_ppt_slides_to_pdf(ppt_path, pdf_path, custom_width=33.86, custom_height=19.05, dpi=300):
    """Convert PowerPoint slides to PDF with custom dimensions and DPI."""
    try:
        logger.info(f"Starting PDF conversion: {ppt_path} -> {pdf_path}")
        # In serverless environment, we skip win32com and go directly to pptx method
        logging.info("Using pptx method for PDF conversion in serverless environment")
        result = convert_with_pptx(ppt_path, pdf_path)
        logger.info(f"PDF conversion result: {result}")
        return result
    except Exception as e:
        logger.error(f"Error in convert_ppt_slides_to_pdf: {str(e)}")
        return False


# Commented out Windows-specific function for serverless compatibility
# def convert_with_win32com(ppt_path, pdf_path):
#     """Convert using win32com (Windows only)."""
#     # This function is disabled in serverless environment
#     return False


def convert_with_pptx(ppt_path, pdf_path):
    """Convert using LibreOffice headless mode or fallback to image-based PDF."""
    try:
        # Try LibreOffice headless conversion first
        if try_libreoffice_conversion(ppt_path, pdf_path):
            return True
        
        # Fallback: Create image-based PDF using slide thumbnails
        return convert_slides_to_image_pdf(ppt_path, pdf_path)
        
    except Exception as e:
        logging.error(f"Error in convert_with_pptx: {str(e)}")
        return False


def convert_slides_to_image_pdf(ppt_path, pdf_path):
    """Convert PowerPoint slides to images and create PDF."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from PIL import Image, ImageDraw, ImageFont
        import io
        
        prs = Presentation(ppt_path)
        
        # Create PDF
        c = canvas.Canvas(pdf_path, pagesize=A4)
        width, height = A4
        
        for i, slide in enumerate(prs.slides):
            # Create a slide image representation
            slide_img = create_slide_image(slide, i+1)
            
            if slide_img:
                # Convert PIL image to ReportLab ImageReader
                img_buffer = io.BytesIO()
                slide_img.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                img_reader = ImageReader(img_buffer)
                
                # Calculate dimensions to fit the page while maintaining aspect ratio
                img_width, img_height = slide_img.size
                aspect_ratio = img_width / img_height
                
                # Use most of the page, leaving margins
                max_width = width - 100  # 50px margin on each side
                max_height = height - 100  # 50px margin on top and bottom
                
                if aspect_ratio > max_width / max_height:
                    # Image is wider, fit to width
                    draw_width = max_width
                    draw_height = max_width / aspect_ratio
                else:
                    # Image is taller, fit to height
                    draw_height = max_height
                    draw_width = max_height * aspect_ratio
                
                # Center the image on the page
                x = (width - draw_width) / 2
                y = (height - draw_height) / 2
                
                c.drawImage(img_reader, x, y, width=draw_width, height=draw_height)
            else:
                # Fallback: text-based slide if image creation fails
                c.setFont("Helvetica-Bold", 16)
                c.drawString(50, height - 50, f"Slide {i+1}")
                
                c.setFont("Helvetica", 12)
                y_position = height - 100
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_lines = shape.text.split('\n')
                        for line in text_lines:
                            if y_position > 50:
                                c.drawString(50, y_position, line[:100])  # Limit line length
                                y_position -= 15
            
            c.showPage()  # New page for next slide
        
        c.save()
        return True
        
    except Exception as e:
        logging.error(f"Error in convert_slides_to_image_pdf: {str(e)}")
        return False


def create_slide_image(slide, slide_number):
    """Create a visual representation of a slide as an image."""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        # Create a white background image (16:9 aspect ratio)
        img_width, img_height = 1920, 1080
        img = Image.new('RGB', (img_width, img_height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Try to load a font, fallback to default if not available
        try:
            title_font = ImageFont.truetype("arial.ttf", 48)
            text_font = ImageFont.truetype("arial.ttf", 32)
            small_font = ImageFont.truetype("arial.ttf", 24)
        except:
            title_font = ImageFont.load_default()
            text_font = ImageFont.load_default()
            small_font = ImageFont.load_default()
        
        # Draw slide number
        draw.text((50, 50), f"Slide {slide_number}", fill='black', font=small_font)
        
        y_position = 150
        
        # Process shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # Determine if this is likely a title (first text or large text)
                is_title = y_position == 150 or len(text) < 100
                font_to_use = title_font if is_title else text_font
                
                # Split text into lines that fit the image width
                words = text.split()
                lines = []
                current_line = ""
                
                for word in words:
                    test_line = current_line + (" " if current_line else "") + word
                    # Rough estimation of text width (more accurate would use textbbox)
                    if len(test_line) * 20 < img_width - 100:  # Leave margins
                        current_line = test_line
                    else:
                        if current_line:
                            lines.append(current_line)
                        current_line = word
                
                if current_line:
                    lines.append(current_line)
                
                # Draw the text lines
                for line in lines:
                    if y_position < img_height - 100:  # Leave bottom margin
                        draw.text((50, y_position), line, fill='black', font=font_to_use)
                        y_position += 60 if is_title else 40
                
                y_position += 20  # Extra space between text blocks
        
        return img
        
    except Exception as e:
        logging.error(f"Error creating slide image: {str(e)}")
        return None


def try_libreoffice_conversion(ppt_path, pdf_path):
    """Try to convert using LibreOffice to PNG images first, then create PDF."""
    temp_dir = None
    try:
        import subprocess
        
        # Common LibreOffice paths
        libreoffice_paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
            "libreoffice",
            "soffice"
        ]
        
        libreoffice_cmd = None
        for path in libreoffice_paths:
            try:
                # Test if LibreOffice is available
                result = subprocess.run([path, "--version"], 
                                       capture_output=True, 
                                       text=True, 
                                       timeout=10)
                if result.returncode == 0:
                    libreoffice_cmd = path
                    break
            except (subprocess.TimeoutExpired, FileNotFoundError, subprocess.SubprocessError):
                continue
        
        if not libreoffice_cmd:
            logging.info("LibreOffice not found")
            return False
        
        # Create temporary directory for PNG images
        temp_dir = tempfile.mkdtemp()
        
        # Convert to PNG images using LibreOffice
        cmd = [
            libreoffice_cmd,
            "--headless",
            "--convert-to", "png",
            "--outdir", temp_dir,
            ppt_path
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        
        if result.returncode == 0:
            # Find all generated PNG files
            png_files = []
            base_name = os.path.splitext(os.path.basename(ppt_path))[0]
            
            # LibreOffice creates files like: filename.png, filename-1.png, filename-2.png, etc.
            for file in os.listdir(temp_dir):
                if file.startswith(base_name) and file.endswith('.png'):
                    png_files.append(os.path.join(temp_dir, file))
            
            # Sort files to maintain slide order
            png_files.sort(key=lambda x: (
                int(x.split('-')[-1].split('.')[0]) if '-' in os.path.basename(x) else 0
            ))
            
            if png_files:
                # Create PDF from PNG images
                create_pdf_from_images(png_files, pdf_path)
                return os.path.exists(pdf_path)
            else:
                logging.error("No PNG files generated by LibreOffice")
                return False
        else:
            logging.error(f"LibreOffice PNG conversion failed: {result.stderr}")
            return False
            
    except Exception as e:
        logging.error(f"Error in LibreOffice conversion: {str(e)}")
        return False
    finally:
        # Clean up temporary directory
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
            except Exception as e:
                logging.warning(f"Failed to clean up temp directory: {str(e)}")


def create_pdf_from_images(image_paths, pdf_path):
    """Create PDF from list of image paths with custom dimensions (33.86cm x 19.05cm, 300 DPI)."""
    try:
        from reportlab.lib.units import cm
        
        # Custom page size: 33.86cm x 19.05cm (landscape)
        page_width = 33.86 * cm
        page_height = 19.05 * cm
        
        c = canvas.Canvas(pdf_path, pagesize=(page_width, page_height))
        
        for image_path in image_paths:
            if os.path.exists(image_path):
                # Get image dimensions
                img = Image.open(image_path)
                img_width, img_height = img.size
                
                # Calculate scaling to fit page while maintaining aspect ratio
                # Leave small margins (1cm on each side)
                available_width = page_width - (2 * cm)
                available_height = page_height - (2 * cm)
                
                scale_x = available_width / img_width
                scale_y = available_height / img_height
                scale = min(scale_x, scale_y)  # Use the smaller scale to fit within bounds
                
                # Calculate centered position
                scaled_width = img_width * scale
                scaled_height = img_height * scale
                x = (page_width - scaled_width) / 2
                y = (page_height - scaled_height) / 2
                
                # Draw image
                c.drawImage(image_path, x, y, width=scaled_width, height=scaled_height)
                c.showPage()
        
        c.save()
        
    except Exception as e:
        logging.error(f"Error creating PDF from images: {str(e)}")
        raise



if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)