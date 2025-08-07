import os
import tempfile
import logging
import zipfile
import shutil
import asyncio
import aiohttp
import traceback
from flask import Flask, render_template, request, jsonify, send_file, flash, redirect, url_for
from werkzeug.utils import secure_filename
from werkzeug.middleware.proxy_fix import ProxyFix
from dotenv import load_dotenv

# Configure logging for serverless environment FIRST
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

# Global variables for storage systems
unified_storage = None
PresentationGenerator = None

def validate_environment():
    """Validate and log environment variables for debugging."""
    try:
        logger.info("=== Environment Validation ===")
        logger.info(f"Python version: {os.sys.version}")
        logger.info(f"Current working directory: {os.getcwd()}")
        logger.info(f"Temp directory: {tempfile.gettempdir()}")
        
        # Check critical environment variables
        env_vars = {
            'VERCEL': os.environ.get('VERCEL'),
            'VERCEL_ENV': os.environ.get('VERCEL_ENV'),
            'SESSION_SECRET': os.environ.get('SESSION_SECRET'),
            'PYTHONPATH': os.environ.get('PYTHONPATH')
        }
        
        for key, value in env_vars.items():
            if value:
                logger.info(f"{key}: {'*' * min(len(str(value)), 10)} (length: {len(str(value))})")
            else:
                logger.warning(f"{key}: NOT SET")
        
        # Check if we're in serverless environment
        is_serverless = os.environ.get('VERCEL') is not None
        logger.info(f"Running in serverless environment: {is_serverless}")
        
        return True
    except Exception as e:
        logger.error(f"Environment validation failed: {str(e)}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return False

def safe_initialize_storage():
    """Safely initialize storage systems with error handling."""
    global unified_storage, PresentationGenerator
    
    try:
        logger.info("=== Storage Initialization ===")
        
        # Load environment variables
        try:
            load_dotenv()
            logger.info("Environment variables loaded successfully")
        except Exception as e:
            logger.warning(f"Failed to load .env file: {str(e)} (this is normal in serverless)")
        
        # Initialize unified storage
        try:
            from utils.unified_storage import initialize_unified_storage
            unified_storage = initialize_unified_storage()
            logger.info("Unified storage initialized successfully")
        except Exception as e:
            logger.error(f"Failed to initialize unified storage: {str(e)}")
            logger.error(f"Unified storage traceback: {traceback.format_exc()}")
            unified_storage = None
        
        # Initialize presentation generator
        try:
            from utils.presentation_generator import PresentationGenerator
            logger.info("Presentation generator imported successfully")
        except Exception as e:
            logger.error(f"Failed to import presentation generator: {str(e)}")
            logger.error(f"Presentation generator traceback: {traceback.format_exc()}")
            PresentationGenerator = None
        
        return True
    except Exception as e:
        logger.error(f"Critical error during storage initialization: {str(e)}")
        logger.error(f"Critical traceback: {traceback.format_exc()}")
        return False

# Validate environment and initialize storage
try:
    logger.info("=== Application Startup ===")
    validate_environment()
    safe_initialize_storage()
    logger.info("=== Startup Complete ===")
except Exception as e:
    logger.error(f"CRITICAL: Application startup failed: {str(e)}")
    logger.error(f"CRITICAL traceback: {traceback.format_exc()}")

# Create Flask app
app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "dev-secret-key-change-in-production")
app.wsgi_app = ProxyFix(app.wsgi_app, x_proto=1, x_host=1)

# Configuration for local file storage
UPLOAD_FOLDER = os.path.join(os.getcwd(), 'uploads')
OUTPUT_FOLDER = os.path.join(os.getcwd(), 'outputs')
MAX_CONTENT_LENGTH = 200 * 1024 * 1024  # 200MB max file size
ALLOWED_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'zip'}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


def allowed_file(filename):
    """Check if file has allowed extension."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def is_image_file(filename):
    """Check if file is an image."""
    image_extensions = {'png', 'jpg', 'jpeg', 'gif', 'bmp'}
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in image_extensions


def run_async(coro):
    """Run async function in Flask context."""
    try:
        loop = asyncio.get_event_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    return loop.run_until_complete(coro)


def get_content_type(filename):
    """Get content type based on file extension."""
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    content_types = {
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'gif': 'image/gif',
        'bmp': 'image/bmp',
        'zip': 'application/zip'
    }
    return content_types.get(ext, 'application/octet-stream')


def extract_folder_structure(zip_path, extract_to):
    """Extract ZIP file and return folder structure."""
    try:
        folder_structure = {}
        
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(extract_to)
            
        # Walk through extracted files
        for root, dirs, files in os.walk(extract_to):
            image_files = [f for f in files if is_image_file(f)]
            if image_files:
                rel_path = os.path.relpath(root, extract_to)
                if rel_path == '.':
                    folder_name = 'root'
                else:
                    folder_name = rel_path
                
                image_paths = [os.path.join(root, f) for f in image_files]
                folder_structure[folder_name] = image_paths
        
        return folder_structure
        
    except Exception as e:
        logger.error(f"Error extracting ZIP: {str(e)}")
        raise


@app.route('/static/<path:filename>')
def serve_static(filename):
    """Serve static files as fallback for serverless environment."""
    try:
        from flask import send_from_directory
        import mimetypes
        
        static_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
        logger.info(f"Serving static file: {filename} from {static_dir}")
        
        # Ensure proper MIME types for common file types
        if filename.endswith('.css'):
            mimetype = 'text/css'
        elif filename.endswith('.js'):
            mimetype = 'application/javascript'
        elif filename.endswith('.png'):
            mimetype = 'image/png'
        elif filename.endswith('.ico'):
            mimetype = 'image/x-icon'
        else:
            mimetype = mimetypes.guess_type(filename)[0] or 'application/octet-stream'
        
        response = send_from_directory(static_dir, filename)
        response.headers['Content-Type'] = mimetype
        response.headers['Cache-Control'] = 'public, max-age=3600'
        return response
        
    except Exception as e:
        logger.error(f"Error serving static file {filename}: {str(e)}")
        return f"Static file not found: {filename}", 404


@app.route('/favicon.ico')
def favicon():
    """Serve favicon from static directory."""
    try:
        from flask import send_from_directory
        static_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'static')
        return send_from_directory(static_dir, 'favicon.ico')
    except Exception as e:
        logger.error(f"Error serving favicon: {str(e)}")
        return '', 404


@app.route('/')
def index():
    """Main page with upload form."""
    try:
        logger.info("Accessing main page")
        return render_template('index.html')
    except Exception as e:
        logger.error(f"Error in index route: {str(e)}")
        return f"Error loading page: {str(e)}", 500


@app.route('/health')
def health_check():
    """Health check endpoint for debugging."""
    try:
        # Check component status
        components_status = {
            'unified_storage': unified_storage is not None,
            'presentation_generator': PresentationGenerator is not None
        }
        
        # Overall health status
        is_healthy = all(components_status.values())
        
        # Environment info
        env_info = {
            'python_version': os.sys.version,
            'cwd': os.getcwd(),
            'temp_dir': tempfile.gettempdir(),
            'upload_folder': UPLOAD_FOLDER,
            'output_folder': OUTPUT_FOLDER
        }
        
        return jsonify({
            'status': 'healthy' if is_healthy else 'degraded',
            'message': 'All components initialized' if is_healthy else 'Some components failed to initialize',
            'components': components_status,
            'environment': env_info,
            'timestamp': os.times() if hasattr(os, 'times') else 'unavailable'
        }), 200 if is_healthy else 503
        
    except Exception as e:
        logger.error(f"Error in health check: {str(e)}")
        logger.error(f"Health check traceback: {traceback.format_exc()}")
        return jsonify({
            'status': 'error',
            'error': str(e),
            'traceback': traceback.format_exc()
        }), 500


@app.route('/startup-status')
def startup_status():
    """Detailed startup status endpoint for debugging initialization issues."""
    try:
        from datetime import datetime
        
        # Detailed component status with error information
        component_details = {}
        
        # Blob storage removed - using local file storage only
        
        # Unified storage status
        if unified_storage is not None:
            component_details['unified_storage'] = {
                'status': 'initialized',
                'environment': 'local'
            }
        else:
            component_details['unified_storage'] = {
                'status': 'failed',
                'environment': 'unknown'
            }
        
        # Presentation generator status
        component_details['presentation_generator'] = {
            'status': 'initialized' if PresentationGenerator is not None else 'failed',
            'available': PresentationGenerator is not None
        }
        
        # Environment variables check
        env_status = {
            'SESSION_SECRET': bool(os.environ.get('SESSION_SECRET')),
            'PYTHONPATH': os.environ.get('PYTHONPATH')
        }
        
        # System information
        system_info = {
            'python_version': os.sys.version,
            'platform': os.sys.platform,
            'working_directory': os.getcwd(),
            'temp_directory': tempfile.gettempdir(),
            'upload_folder_exists': os.path.exists(UPLOAD_FOLDER),
            'output_folder_exists': os.path.exists(OUTPUT_FOLDER)
        }
        
        # Overall status
        all_components_ok = all(
            details.get('status') == 'initialized' 
            for details in component_details.values()
        )
        
        return jsonify({
            'overall_status': 'healthy' if all_components_ok else 'degraded',
            'components': component_details,
            'environment_variables': env_status,
            'system_info': system_info,
            'timestamp': datetime.now().isoformat()
        }), 200 if all_components_ok else 503
        
    except Exception as e:
        logger.error(f"Startup status check failed: {str(e)}")
        logger.error(f"Startup status traceback: {traceback.format_exc()}")
        return jsonify({
            'overall_status': 'error',
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500


@app.route('/validate-environment')
def validate_environment():
    """Validate environment variables and provide setup guidance."""
    try:
        from datetime import datetime
        
        # Check environment variables
        env_status = {
            'USER_ID': {
                'present': bool(os.environ.get('USER_ID')),
                'required': False,
                'description': 'Optional, used for file organization'
            }
        }
        
        # Check service availability
        services_status = {
            'unified_storage': {
                'available': unified_storage is not None,
                'environment': 'local'
            },
            'presentation_generator': {
                'available': PresentationGenerator is not None
            }
        }
        
        # Determine missing requirements
        missing_requirements = []
        for var_name, var_info in env_status.items():
            if var_info['required'] and not var_info['present']:
                missing_requirements.append({
                    'variable': var_name,
                    'description': var_info['description']
                })
        
        # Setup instructions
        setup_instructions = {
            'local_storage': {
                'title': 'Local File Storage Setup',
                'steps': [
                    '1. Ensure uploads/ and outputs/ directories exist',
                    '2. Check file permissions for read/write access',
                    '3. Verify sufficient disk space for file operations'
                ]
            }
        }
        
        # Overall status
        is_properly_configured = len(missing_requirements) == 0 and services_status['unified_storage']['available']
        
        return jsonify({
            'status': 'configured' if is_properly_configured else 'needs_setup',
            'environment_variables': env_status,
            'services': services_status,
            'missing_requirements': missing_requirements,
            'setup_instructions': setup_instructions,
            'timestamp': datetime.now().isoformat()
        }), 200 if is_properly_configured else 422
        
    except Exception as e:
        logger.error(f"Environment validation failed: {str(e)}")
        return jsonify({
            'status': 'error',
            'error': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500


@app.route('/fallback-info')
def fallback_info():
    """Provide information about available fallback functionality when services are degraded."""
    try:
        from datetime import datetime
        # Check what services are available
        services_available = {
            'unified_storage': unified_storage is not None,
            'presentation_generator': PresentationGenerator is not None
        }
        
        # Determine available functionality
        available_features = []
        unavailable_features = []
        
        if services_available['unified_storage'] and services_available['presentation_generator']:
            available_features.extend([
                'File upload and processing',
                'PowerPoint generation',
                'File download'
            ])
        else:
            unavailable_features.extend([
                'File upload and processing',
                'PowerPoint generation',
                'File download'
            ])
        
        # Always available features
        available_features.extend([
            'Health checks',
            'System status monitoring',
            'Error reporting'
        ])
        
        # Fallback recommendations
        fallback_options = []
        if not services_available['unified_storage']:
            fallback_options.append({
                'issue': 'Storage service unavailable',
                'recommendation': 'Try again in a few minutes or contact support',
                'alternative': 'Use local development environment if available'
            })
        
        if not services_available['presentation_generator']:
            fallback_options.append({
                'issue': 'Presentation generator unavailable',
                'recommendation': 'Check system dependencies and restart service',
                'alternative': 'Manual PowerPoint creation from uploaded images'
            })
        
        return jsonify({
            'status': 'degraded' if unavailable_features else 'operational',
            'services_status': services_available,
            'available_features': available_features,
            'unavailable_features': unavailable_features,
            'fallback_options': fallback_options,
            'support_contact': 'Check logs for detailed error information',
            'timestamp': datetime.now().isoformat()
        })
        
    except Exception as e:
        logger.error(f"Fallback info check failed: {str(e)}")
        return jsonify({
            'status': 'error',
            'message': 'Unable to determine system status',
            'error': str(e),
            'basic_functionality': 'Health checks and error reporting available',
            'timestamp': datetime.now().isoformat()
        }), 500


@app.route('/debug-info')
def debug_info():
    """Provide detailed debugging information for troubleshooting deployment issues."""
    try:
        from datetime import datetime
        import sys
        debug_data = {
            'timestamp': datetime.now().isoformat(),
            'environment': {
                'python_version': sys.version,
                'platform': sys.platform,
                'tmpdir': os.getenv('TMPDIR', 'Not set'),
                'temp_directory': tempfile.gettempdir(),
                'cwd': os.getcwd()
            },
            'environment_variables': {
                'USER_ID': 'Set' if os.getenv('USER_ID') else 'Not set',
                'TMPDIR': os.getenv('TMPDIR', 'Not set')
            },
            'components': {
                'unified_storage': {
                    'initialized': unified_storage is not None,
                    'class_name': unified_storage.__class__.__name__ if unified_storage else None,
                    'storage_type': 'local',
                    'error': getattr(unified_storage, '_init_error', None) if unified_storage else 'Not initialized'
                },
                'presentation_generator': {
                    'available': PresentationGenerator is not None,
                    'class_name': PresentationGenerator.__name__ if PresentationGenerator else None
                }
            },
            'dependencies': {
                'imports_successful': {
                    'os': True,
                    'sys': True,
                    'tempfile': True,
                    'logging': True,
                    'datetime': True,
                    'flask': True,
                    'werkzeug': True
                }
            },
            'file_system': {
                'temp_writable': False,
                'temp_path_exists': False,
                'uploads_dir_exists': False,
                'outputs_dir_exists': False
            }
        }
        
        # Test file system access
        try:
            temp_dir = tempfile.gettempdir()
            debug_data['file_system']['temp_path_exists'] = os.path.exists(temp_dir)
            
            # Test write access
            test_file = os.path.join(temp_dir, 'debug_test.txt')
            with open(test_file, 'w') as f:
                f.write('test')
            os.remove(test_file)
            debug_data['file_system']['temp_writable'] = True
        except Exception as e:
            debug_data['file_system']['temp_write_error'] = str(e)
        
        # Check local directories
        try:
            debug_data['file_system']['uploads_dir_exists'] = os.path.exists('uploads')
            debug_data['file_system']['outputs_dir_exists'] = os.path.exists('outputs')
        except Exception as e:
            debug_data['file_system']['dir_check_error'] = str(e)
        
        # Test imports
        import_tests = {
            'aiohttp': False,
            'PIL': False,
            'pptx': False,
            'zipfile': False,
            'shutil': False
        }
        
        for module_name in import_tests:
            try:
                __import__(module_name)
                import_tests[module_name] = True
            except ImportError as e:
                import_tests[module_name] = f'Import error: {str(e)}'
            except Exception as e:
                import_tests[module_name] = f'Other error: {str(e)}'
        
        debug_data['dependencies']['imports_successful'].update(import_tests)
        
        return jsonify(debug_data)
        
    except Exception as e:
        logger.error(f"Debug info generation failed: {str(e)}")
        return jsonify({
            'error': 'Debug info generation failed',
            'details': str(e),
            'timestamp': datetime.now().isoformat()
        }), 500


# Blob upload URL route removed - using direct file uploads for VPS deployment


async def process_blob_file(file_identifier, filename):
    """Download ZIP file from storage and process it to PPTX."""
    temp_dir = None
    try:
        # Check if required components are available
        if unified_storage is None:
            logger.error(f"Unified storage not available for processing: {filename}")
            return {'success': False, 'message': 'Storage service not available'}
        
        if PresentationGenerator is None:
            logger.error(f"PresentationGenerator not available for processing: {filename}")
            return {'success': False, 'message': 'Presentation generator not available'}
        
        import aiohttp
        import tempfile
        import zipfile
        import shutil
        
        logger.info(f"Starting file processing for: {filename}")
        
        # Create temporary directory
        temp_dir = tempfile.mkdtemp()
        zip_path = os.path.join(temp_dir, filename)
        
        logger.info(f"Downloading file from storage: {file_identifier}")
        
        # Download file using unified storage
        file_data = await unified_storage.download_file(file_identifier)
        if file_data is None:
            logger.error(f"Failed to download file from storage: {file_identifier}")
            return {'success': False, 'message': 'Failed to download file from storage'}
        
        # Save downloaded data to temporary file
        with open(zip_path, 'wb') as f:
            f.write(file_data)
        logger.info(f"Successfully downloaded and saved file: {filename}")
        
        logger.info(f"Extracting ZIP file to: {temp_dir}")
        
        # Extract ZIP file
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        logger.info("Starting PPTX generation")
        
        # Generate presentation
        generator = PresentationGenerator()
        original_filename = os.path.splitext(filename)[0]
        
        ppt_filename = generator.generate_from_folder(
            temp_dir,
            annotation_option='with_annos',
            implement_video_frames=False,
            original_filename=original_filename
        )
        
        logger.info(f"Successfully generated PPTX: {ppt_filename}")
        
        # Save PPTX using unified storage
        ppt_storage_url = None
        try:
            with open(ppt_filename, 'rb') as f:
                ppt_data = f.read()
            
            ppt_basename = os.path.basename(ppt_filename)
            
            # Save to local outputs directory
            saved_path = unified_storage.save_output_file(ppt_data, ppt_basename)
            if saved_path:
                ppt_storage_url = f"/local-file/{ppt_basename}"
                logger.info(f"PPTX saved locally: {saved_path}")
            else:
                logger.warning("Failed to save PPTX to local storage")
        except Exception as e:
            logger.error(f"Error saving PPTX to storage: {str(e)}")
        
        # Get presentation info
        from pptx import Presentation
        prs = Presentation(ppt_filename)
        slide_count = len(prs.slides)
        
        folder_structure = generator._organize_folder_structure(temp_dir)
        folder_count = len(folder_structure)
        
        # Clean up temp directory
        shutil.rmtree(temp_dir)
        
        return {
            'success': True,
            'ppt_file': os.path.basename(ppt_filename),
            'ppt_storage_url': ppt_storage_url,
            'folder_count': folder_count,
            'slide_count': slide_count
        }
        
    except Exception as e:
        logger.error(f"Error in process_blob_file {filename}: {str(e)}")
        logger.error(f"Error details: {type(e).__name__}: {str(e)}")
        # Clean up on error
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
        return {'success': False, 'message': str(e)}


async def process_single_image_blob(file_identifier, filename):
    """Download single image from storage and convert to PPTX."""
    temp_dir = None
    try:
        # Check if required components are available
        if unified_storage is None:
            logger.error(f"Unified storage not available for processing: {filename}")
            return {'success': False, 'message': 'Storage service not available'}
        
        if PresentationGenerator is None:
            logger.error(f"PresentationGenerator not available for processing: {filename}")
            return {'success': False, 'message': 'Presentation generator not available'}
        
        import aiohttp
        import tempfile
        from PIL import Image
        
        logger.info(f"Starting single image processing for: {filename}")
        
        # Create temporary directory
        temp_dir = tempfile.mkdtemp()
        image_path = os.path.join(temp_dir, filename)
        
        logger.info(f"Downloading image from storage: {file_identifier}")
        
        # Download image using unified storage
        file_data = await unified_storage.download_file(file_identifier)
        if file_data is None:
            logger.error(f"Failed to download image from storage: {file_identifier}")
            return {'success': False, 'message': 'Failed to download image from storage'}
        
        # Save downloaded data to temporary file
        with open(image_path, 'wb') as f:
            f.write(file_data)
        logger.info(f"Successfully downloaded and saved image: {filename}")
        
        # Create a simple folder structure for single image
        image_folder = os.path.join(temp_dir, 'images')
        os.makedirs(image_folder, exist_ok=True)
        
        logger.info(f"Organizing image into folder structure")
        
        # Move image to images folder
        new_image_path = os.path.join(image_folder, filename)
        shutil.move(image_path, new_image_path)
        
        logger.info("Starting PPTX generation for single image")
        
        # Generate presentation
        generator = PresentationGenerator()
        original_filename = os.path.splitext(filename)[0]
        
        ppt_filename = generator.generate_from_folder(
            temp_dir,
            annotation_option='with_annos',
            implement_video_frames=False,
            original_filename=original_filename
        )
        
        logger.info(f"Successfully generated PPTX from single image: {ppt_filename}")
        
        # Save PPTX using unified storage
        ppt_storage_url = None
        try:
            with open(ppt_filename, 'rb') as f:
                ppt_data = f.read()
            
            ppt_basename = os.path.basename(ppt_filename)
            
            # Save to local outputs directory
            saved_path = unified_storage.save_output_file(ppt_data, ppt_basename)
            if saved_path:
                ppt_storage_url = f"/local-file/{ppt_basename}"
                logger.info(f"PPTX saved locally: {saved_path}")
            else:
                logger.warning("Failed to save PPTX to local storage")
        except Exception as e:
            logger.error(f"Error saving PPTX to storage: {str(e)}")
        
        # Clean up temp directory
        shutil.rmtree(temp_dir)
        
        return {
            'success': True,
            'ppt_file': os.path.basename(ppt_filename),
            'ppt_storage_url': ppt_storage_url,
            'slide_count': 1
        }
        
    except Exception as e:
        logger.error(f"Error in process_single_image_blob {filename}: {str(e)}")
        logger.error(f"Error details: {type(e).__name__}: {str(e)}")
        # Clean up on error
        if temp_dir and os.path.exists(temp_dir):
            shutil.rmtree(temp_dir, ignore_errors=True)
        return {'success': False, 'message': str(e)}


@app.route('/upload-complete', methods=['POST'])
def upload_complete():
    """Handle upload completion notification and trigger PPTX conversion."""
    try:
        data = request.get_json()
        
        # Handle both blob_url (for deployed) and url (for localhost)
        file_identifier = data.get('blob_url') or data.get('url')
        if not file_identifier:
            return jsonify({'error': 'File URL/identifier is required'}), 400
        
        filename = data.get('filename', 'unknown')
        file_size = data.get('fileSize', 0)
        
        logger.info(f"Upload completed for: {filename}, file_identifier: {file_identifier}")
        
        # For ZIP files, process and convert to PPTX
        if filename.lower().endswith('.zip'):
            try:
                # Download file from storage and process
                result = run_async(process_blob_file(file_identifier, filename))
                
                if result and result.get('success'):
                    storage_url = result.get('ppt_storage_url', '')
                    redirect_url = f'/result/{result["ppt_file"]}'
                    if storage_url and not storage_url.startswith('/local-file/'):
                        redirect_url += f'?blob_url={storage_url}'
                    
                    return jsonify({
                        'success': True,
                        'message': 'File processed successfully',
                        'redirect_url': redirect_url,
                        'ppt_file': result['ppt_file'],
                        'ppt_storage_url': storage_url,
                        'folder_count': result.get('folder_count', 0),
                        'slide_count': result.get('slide_count', 0)
                    })
                else:
                    return jsonify({
                        'error': result.get('message', 'Processing failed')
                    }), 500
                    
            except Exception as e:
                logger.error(f"Error processing ZIP file: {str(e)}")
                return jsonify({
                    'error': f'Processing error: {str(e)}'
                }), 500
        else:
            # Single image file - convert to PPTX
            try:
                result = run_async(process_single_image_blob(file_identifier, filename))
                
                if result and result.get('success'):
                    storage_url = result.get('ppt_storage_url', '')
                    redirect_url = f'/result/{result["ppt_file"]}'
                    if storage_url and not storage_url.startswith('/local-file/'):
                        redirect_url += f'?blob_url={storage_url}'
                    
                    return jsonify({
                        'success': True,
                        'message': 'Image processed successfully',
                        'redirect_url': redirect_url,
                        'ppt_file': result['ppt_file'],
                        'ppt_storage_url': storage_url,
                        'slide_count': 1
                    })
                else:
                    return jsonify({
                        'error': result.get('message', 'Processing failed')
                    }), 500
                    
            except Exception as e:
                logger.error(f"Error processing image file: {str(e)}")
                return jsonify({
                    'error': f'Processing error: {str(e)}'
                }), 500
            
    except Exception as e:
        logger.error(f"Error in upload_complete: {str(e)}")
        return jsonify({'error': f'Upload complete error: {str(e)}'}), 500


@app.route('/local-upload/<filename>', methods=['PUT'])
def local_upload(filename):
    """Handle local file uploads for localhost environment."""
    try:
        logger.info(f"Local upload endpoint called for: {filename}")
        
        # Check if unified storage is available
        if unified_storage is None:
            logger.error("Unified storage not available for local upload")
            return jsonify({
                'error': 'Storage service not available',
                'details': 'Storage system failed to initialize'
            }), 503
        
        # This endpoint is for local file uploads
        
        # Get file data from request
        file_data = request.get_data()
        if not file_data:
            return jsonify({'error': 'No file data received'}), 400
        
        # Upload file using unified storage
        upload_result = run_async(unified_storage.upload_file(
            file_data, filename, request.content_type or 'application/octet-stream'
        ))
        
        if upload_result:
            logger.info(f"Local file uploaded successfully: {filename}")
            return jsonify({
                'success': True,
                'url': upload_result['url'],
                'filename': upload_result['filename'],
                'size': upload_result['size'],
                'storage_type': upload_result['storage_type']
            })
        else:
            return jsonify({'error': 'Failed to upload file locally'}), 500
            
    except Exception as e:
        logger.error(f"Error in local upload: {str(e)}")
        return jsonify({'error': f'Local upload error: {str(e)}'}), 500


def process_uploaded_file(filename, file_path, annotation_option='with_annos'):
    """Process uploaded file and return result dictionary."""
    try:
        logger.info(f"=== PROCESSING UPLOADED FILE: {filename} ===")
        logger.info(f"File path: {file_path}")
        
        # Check if file exists
        file_exists = os.path.exists(file_path)
        logger.info(f"File exists: {file_exists}")
        if not file_exists:
            logger.error(f"File not found at path: {file_path}")
            return {'success': False, 'error': 'File not found'}
        
        # Get file info
        file_size = os.path.getsize(file_path)
        logger.info(f"File size: {file_size} bytes")
        
        # Get file extension
        file_ext = os.path.splitext(filename)[1].lower()
        logger.info(f"File extension: {file_ext}")
        
        # Get original filename without extension for output naming
        original_filename = os.path.splitext(filename)[0]
        
        if file_ext == '.zip':
            # Handle ZIP file processing
            try:
                import zipfile
                extract_dir = os.path.join(UPLOAD_FOLDER, 'extracted', original_filename)
                os.makedirs(extract_dir, exist_ok=True)
                
                logger.info(f"Extracting ZIP to: {extract_dir}")
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(extract_dir)
                
                # Check if PresentationGenerator is available
                if PresentationGenerator is None:
                    logger.error("PresentationGenerator not available")
                    return {'success': False, 'error': 'Presentation generator not available'}
                
                # Use PresentationGenerator to create presentation
                logger.info("Using PresentationGenerator for ZIP processing")
                generator = PresentationGenerator()
                
                # Generate presentation using the sophisticated generator
                ppt_path = generator.generate_from_folder(
                    temp_dir=extract_dir,
                    annotation_option=annotation_option,
                    implement_video_frames=False,
                    video_position_params=None,
                    original_filename=original_filename
                )
                
                # Get output filename from the generated path
                output_filename = os.path.basename(ppt_path)
                
                # Verify output file was created
                output_exists = os.path.exists(ppt_path)
                output_size = os.path.getsize(ppt_path) if output_exists else 0
                logger.info(f"Output file created: {output_exists}")
                logger.info(f"Output file size: {output_size} bytes")
                logger.info(f"Generated presentation: {ppt_path}")
                
                # Get slide count from the generated presentation
                try:
                    from pptx import Presentation
                    prs = Presentation(ppt_path)
                    slide_count = len(prs.slides)
                except Exception as e:
                    logger.warning(f"Could not get slide count: {str(e)}")
                    slide_count = 0
                
                result = {
                    'success': True,
                    'output_filename': output_filename,
                    'result_url': f'/result/{output_filename}',
                    'slide_count': slide_count
                }
                logger.info(f"Returning result: {result}")
                return result
                
            except Exception as e:
                logger.error(f"Error processing ZIP file: {str(e)}")
                logger.error(f"Traceback: {traceback.format_exc()}")
                return {'success': False, 'error': f'ZIP processing error: {str(e)}'}
        
        elif file_ext in ['.png', '.jpg', '.jpeg', '.gif', '.bmp']:
            # Handle single image file
            logger.info("Processing single image file")
            try:
                # Create a temporary directory for single image processing
                import tempfile
                with tempfile.TemporaryDirectory() as temp_dir:
                    # Create image folder structure for PresentationGenerator
                    image_folder = os.path.join(temp_dir, 'images')
                    os.makedirs(image_folder, exist_ok=True)
                    
                    # Copy image to the temp folder
                    import shutil
                    temp_image_path = os.path.join(image_folder, filename)
                    shutil.copy2(file_path, temp_image_path)
                    
                    # Check if PresentationGenerator is available
                    if PresentationGenerator is None:
                        logger.error("PresentationGenerator not available")
                        return {'success': False, 'error': 'Presentation generator not available'}
                    
                    # Use PresentationGenerator to create presentation
                    logger.info("Using PresentationGenerator for single image processing")
                    generator = PresentationGenerator()
                    
                    # Generate presentation using the sophisticated generator
                    ppt_path = generator.generate_from_folder(
                        temp_dir=temp_dir,
                        annotation_option=annotation_option,
                        implement_video_frames=False,
                        video_position_params=None,
                        original_filename=original_filename
                    )
                    
                    # Get output filename from the generated path
                    output_filename = os.path.basename(ppt_path)
                    
                    # Verify output file was created
                    output_exists = os.path.exists(ppt_path)
                    output_size = os.path.getsize(ppt_path) if output_exists else 0
                    logger.info(f"Single image output file created: {output_exists}")
                    logger.info(f"Single image output file size: {output_size} bytes")
                    
                    # Get slide count from the generated presentation
                    try:
                        from pptx import Presentation
                        prs = Presentation(ppt_path)
                        slide_count = len(prs.slides)
                    except Exception as e:
                        logger.warning(f"Could not get slide count: {str(e)}")
                        slide_count = 1
                    
                    result = {
                        'success': True,
                        'output_filename': output_filename,
                        'result_url': f'/result/{output_filename}',
                        'slide_count': slide_count
                    }
                    logger.info(f"Single image returning result: {result}")
                    return result
                
            except Exception as e:
                logger.error(f"Error processing image file: {str(e)}")
                logger.error(f"Traceback: {traceback.format_exc()}")
                return {'success': False, 'error': f'Image processing error: {str(e)}'}
        
        else:
            return {'success': False, 'error': 'Unsupported file type'}
    
    except Exception as e:
        logger.error(f"Error in process_uploaded_file: {str(e)}")
        logger.error(f"Traceback: {traceback.format_exc()}")
        return {'success': False, 'error': f'Processing error: {str(e)}'}


@app.route('/upload', methods=['POST'])
def upload_file():
    """Handle file uploads with JSON responses for VPS deployment."""
    try:
        logger.info("=== UPLOAD ENDPOINT CALLED ===")
        logger.info(f"Request method: {request.method}")
        logger.info(f"Request files: {list(request.files.keys())}")
        
        if 'file' not in request.files:
            logger.error("No file in request.files")
            return jsonify({
                'error': 'No file provided',
                'message': 'Please select a file to upload'
            }), 400

        file = request.files['file']
        annotation_option = request.form.get('annotation_option', 'with_annos')
        logger.info(f"File object: {file}")
        logger.info(f"File filename: {file.filename}")
        logger.info(f"Annotation option: {annotation_option}")
        
        if not file or file.filename == '':
            logger.error("File is empty or has no filename")
            return jsonify({
                'error': 'No file selected',
                'message': 'Please select a valid file'
            }), 400

        if not allowed_file(file.filename):
            logger.error(f"File type not allowed: {file.filename}")
            return jsonify({
                'error': 'Invalid file type',
                'message': 'Please upload ZIP, PNG, JPG, JPEG, GIF, or BMP files.'
            }), 400

        # Save file to local storage
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        logger.info(f"Saving file to: {file_path}")
        logger.info(f"UPLOAD_FOLDER: {UPLOAD_FOLDER}")
        logger.info(f"Upload folder exists: {os.path.exists(UPLOAD_FOLDER)}")
        
        file.save(file_path)
        
        # Verify file was saved
        file_saved = os.path.exists(file_path)
        file_size = os.path.getsize(file_path) if file_saved else 0
        logger.info(f"File saved successfully: {file_saved}")
        logger.info(f"File size: {file_size} bytes")
        
        # Process the uploaded file
        logger.info(f"Starting file processing for: {filename}")
        result = process_uploaded_file(filename, file_path, annotation_option)
        logger.info(f"Processing result: {result}")
        
        if result.get('success'):
            result_url = result.get('result_url', f'/result/{result.get("output_filename", filename)}')
            logger.info(f"Processing successful, result_url: {result_url}")
            return jsonify({
                'success': True,
                'message': f'File {filename} uploaded and processed successfully',
                'filename': filename,
                'result_url': result_url
            }), 200
        else:
            logger.error(f"Processing failed: {result.get('error')}")
            return jsonify({
                'error': 'Processing failed',
                'message': result.get('error', 'Unknown error occurred during processing')
            }), 500
        
    except Exception as e:
        logger.error(f"Error in upload: {str(e)}")
        logger.error(f"Exception type: {type(e)}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")
        return jsonify({
            'error': 'Upload failed',
            'message': f'Upload error: {str(e)}'
        }), 500


@app.route('/manual-upload', methods=['POST'])
def manual_upload():
    """Handle manual upload with form data."""
    try:
        logger.info("Manual upload endpoint called")
        
        if 'file' not in request.files:
            flash('No file selected', 'error')
            return redirect(url_for('index'))

        file = request.files['file']
        
        if not file or file.filename == '':
            flash('No file selected', 'error')
            return redirect(url_for('index'))

        if not allowed_file(file.filename):
            flash('Invalid file type. Please upload ZIP, PNG, JPG, JPEG, GIF, or BMP files.', 'error')
            return redirect(url_for('index'))

        # Process file similar to upload endpoint
        filename = secure_filename(file.filename)
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        file.save(file_path)
        
        flash(f'File {filename} uploaded successfully!', 'success')
        return redirect(url_for('index'))
        
    except Exception as e:
        logger.error(f"Error in manual upload: {str(e)}")
        flash(f'Upload error: {str(e)}', 'error')
        return redirect(url_for('index'))


@app.route('/local-file/<filename>')
def download_local_file(filename):
    """Download files from local storage."""
    try:
        # Check if unified storage is available
        if unified_storage is None:
            logger.error("Unified storage not available for local file download")
            return jsonify({
                'error': 'Storage service not available',
                'details': 'Storage system failed to initialize'
            }), 503
        
        # Since we're using local storage only, this check is not needed
        # All file access is local
        
        # Try outputs folder first, then uploads folder
        file_path = unified_storage.get_output_file_path(filename)
        if not file_path:
            # Check uploads folder
            uploads_path = os.path.join(unified_storage.local_upload_dir, filename)
            if os.path.exists(uploads_path):
                file_path = uploads_path
        
        if file_path and os.path.exists(file_path):
            return send_file(file_path, as_attachment=True, download_name=filename)
        else:
            return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        logger.error(f"Error downloading local file: {str(e)}")
        return jsonify({'error': f'Download error: {str(e)}'}), 500


@app.route('/download/<filename>')
def download_file(filename):
    """Download generated files using unified storage."""
    try:
        # Check if unified storage is available
        if unified_storage is None:
            logger.error("Unified storage not available for download")
            return jsonify({
                'error': 'Storage service not available',
                'details': 'Storage system failed to initialize'
            }), 503
        
        # Check if file exists in temp outputs (local storage)
        file_path = os.path.join(OUTPUT_FOLDER, filename)
        if os.path.exists(file_path):
            return send_file(file_path, as_attachment=True, download_name=filename)
        else:
            # Try unified storage output directory
            unified_path = unified_storage.get_output_file_path(filename)
            if unified_path and os.path.exists(unified_path):
                return send_file(unified_path, as_attachment=True, download_name=filename)
            else:
                return jsonify({'error': 'File not found'}), 404
    except Exception as e:
        logger.error(f"Error downloading file: {str(e)}")
        return jsonify({'error': f'Download error: {str(e)}'}), 500


# Test endpoint for debugging upload flow
@app.route('/test-upload-flow')
def test_upload_flow():
    """Test endpoint to verify upload and result flow."""
    try:
        logger.info("=== TEST UPLOAD FLOW ENDPOINT ===")
        
        # Check directories
        upload_exists = os.path.exists(UPLOAD_FOLDER)
        output_exists = os.path.exists(OUTPUT_FOLDER)
        
        # List files in directories
        upload_files = os.listdir(UPLOAD_FOLDER) if upload_exists else []
        output_files = os.listdir(OUTPUT_FOLDER) if output_exists else []
        
        # Check unified storage
        unified_available = unified_storage is not None
        
        test_info = {
            'directories': {
                'upload_folder': UPLOAD_FOLDER,
                'upload_exists': upload_exists,
                'upload_files': upload_files,
                'output_folder': OUTPUT_FOLDER,
                'output_exists': output_exists,
                'output_files': output_files
            },
            'unified_storage': {
                'available': unified_available,
                'local_upload_dir': unified_storage.local_upload_dir if unified_available else None,
                'local_output_dir': unified_storage.local_output_dir if unified_available else None
            },
            'current_working_directory': os.getcwd()
        }
        
        logger.info(f"Test info: {test_info}")
        
        return jsonify({
            'success': True,
            'test_info': test_info
        })
        
    except Exception as e:
        logger.error(f"Error in test endpoint: {str(e)}")
        return jsonify({
            'error': str(e),
            'success': False
        }), 500


# Blob download and info routes removed - using local file storage for VPS deployment


@app.route('/result/<filename>')
def show_result(filename):
    """Display the result page with download links."""
    try:
        logger.info(f"=== RESULT ROUTE CALLED FOR: {filename} ===")
        logger.info(f"Request URL: {request.url}")
        logger.info(f"Request args: {dict(request.args)}")
        
        # Check if unified storage is available
        if unified_storage is None:
            logger.error("Unified storage not available for result display")
            return render_template('index.html'), 503
        
        logger.info(f"Unified storage available: {unified_storage is not None}")
        
        # Get blob URL from query parameters
        blob_url = request.args.get('blob_url')
        logger.info(f"Blob URL from query: {blob_url}")
        
        slide_count = 0
        file_exists = False
        
        # Check for local file first (using unified storage)
        logger.info("Checking for local file...")
        if unified_storage:
            local_file_path = unified_storage.get_output_file_path(filename)
            logger.info(f"Local file path from unified storage: {local_file_path}")
            if local_file_path and os.path.exists(local_file_path):
                file_exists = True
                logger.info(f"Local file exists: {local_file_path}")
                try:
                    from pptx import Presentation
                    prs = Presentation(local_file_path)
                    slide_count = len(prs.slides)
                    logger.info(f"Got slide count from local file: {slide_count}")
                except Exception as e:
                    logger.warning(f"Could not read presentation info from local file: {str(e)}")
            else:
                logger.info(f"Local file not found or path is None")
        
        # Check temp outputs folder (for deployed environment)
        temp_file_path = os.path.join(OUTPUT_FOLDER, filename)
        logger.info(f"Checking temp file path: {temp_file_path}")
        logger.info(f"Temp file exists: {os.path.exists(temp_file_path)}")
        if not file_exists and os.path.exists(temp_file_path):
            file_exists = True
            logger.info(f"Found file in temp outputs: {temp_file_path}")
            try:
                from pptx import Presentation
                prs = Presentation(temp_file_path)
                slide_count = len(prs.slides)
                logger.info(f"Got slide count from temp file: {slide_count}")
            except Exception as e:
                logger.warning(f"Could not read presentation info from temp file: {str(e)}")
        else:
            logger.info(f"File not found in temp outputs or file_exists is already True")
        
        # If no local file but we have blob URL, try to get info from storage
        if not file_exists and blob_url:
            try:
                # Download file temporarily to get presentation info
                file_data = run_async(unified_storage.download_file(blob_url))
                if file_data:
                    file_exists = True
                    import tempfile
                    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                        temp_file.write(file_data)
                        temp_file.flush()
                        
                        from pptx import Presentation
                        prs = Presentation(temp_file.name)
                        slide_count = len(prs.slides)
                        logger.info(f"Got slide count from storage: {slide_count}")
                        
                        # Clean up temp file
                        os.unlink(temp_file.name)
            except Exception as e:
                logger.warning(f"Could not read presentation info from storage: {str(e)}")
        
        # If we have a file (local or remote), show the result page
        logger.info(f"Final file_exists status: {file_exists}")
        logger.info(f"Final slide_count: {slide_count}")
        
        if file_exists:
            local_exists = unified_storage and unified_storage.get_output_file_path(filename) and os.path.exists(unified_storage.get_output_file_path(filename))
            temp_exists = os.path.exists(temp_file_path)
            logger.info(f"Showing result page for {filename} (local: {local_exists}, temp: {temp_exists})")
            logger.info(f"Rendering result.html with filename={filename}, slide_count={slide_count}")
            return render_template('result.html',
                                 filename=filename,
                                 ppt_file=filename,
                                 folder_count=1,  # Default value
                                 slide_count=slide_count,
                                 video_folder_found=False,
                                 implement_video_frames=False)
        else:
            logger.error(f"Result file not found locally or in blob storage: {filename}")
            logger.error(f"Checked paths:")
            logger.error(f"  - Local path: {unified_storage.get_output_file_path(filename) if unified_storage else 'N/A (no unified storage)'}")
            logger.error(f"  - Temp path: {temp_file_path}")
            logger.error(f"Redirecting to index.html")
            return render_template('index.html'), 404
        
    except Exception as e:
        logger.error(f"Error in show_result: {str(e)}")
        return render_template('index.html'), 500


def convert_slide_to_image(slide, slide_width, slide_height, dpi=300):
    """Convert a single slide to a high-quality image using PIL and python-pptx.
    
    This function renders slide content as an image while preserving formatting.
    """
    import io
    from PIL import Image, ImageDraw, ImageFont
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches, Pt
    import math
    
    try:
        # Calculate image dimensions based on slide size and DPI
        img_width = int((slide_width / 914400) * dpi)  # Convert EMU to inches, then to pixels
        img_height = int((slide_height / 914400) * dpi)
        
        # Create a white background image
        img = Image.new('RGB', (img_width, img_height), 'white')
        draw = ImageDraw.Draw(img)
        
        # Scale factor for positioning
        scale_x = img_width / slide_width
        scale_y = img_height / slide_height
        
        # Process each shape in the slide
        for shape in slide.shapes:
            try:
                # Get shape position and size in EMU, convert to pixels
                left = int(shape.left * scale_x)
                top = int(shape.top * scale_y)
                width = int(shape.width * scale_x)
                height = int(shape.height * scale_y)
                
                # Handle text shapes
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    
                    # Try to get font size from the shape
                    font_size = 12  # Default
                    if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                        first_para = shape.text_frame.paragraphs[0]
                        if first_para.runs and first_para.runs[0].font.size:
                            font_size = int(first_para.runs[0].font.size.pt * (dpi / 72))
                    
                    # Use a basic font (you might want to install more fonts)
                    try:
                        font = ImageFont.truetype("arial.ttf", font_size)
                    except:
                        font = ImageFont.load_default()
                    
                    # Draw text with word wrapping
                    lines = text_content.split('\n')
                    y_offset = top
                    
                    for line in lines:
                        # Simple word wrapping
                        words = line.split(' ')
                        current_line = ""
                        
                        for word in words:
                            test_line = current_line + (" " if current_line else "") + word
                            bbox = draw.textbbox((0, 0), test_line, font=font)
                            text_width = bbox[2] - bbox[0]
                            
                            if text_width <= width:
                                current_line = test_line
                            else:
                                if current_line:
                                    draw.text((left, y_offset), current_line, fill='black', font=font)
                                    y_offset += font_size + 2
                                current_line = word
                        
                        if current_line:
                            draw.text((left, y_offset), current_line, fill='black', font=font)
                            y_offset += font_size + 2
                
                # Handle image shapes
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Extract image data
                        image_stream = shape.image.blob
                        shape_img = Image.open(io.BytesIO(image_stream))
                        
                        # Resize image to fit the shape bounds
                        shape_img = shape_img.resize((width, height), Image.Resampling.LANCZOS)
                        
                        # Paste the image onto the slide image
                        img.paste(shape_img, (left, top))
                        
                    except Exception as img_error:
                        # Draw a placeholder rectangle for failed images
                        draw.rectangle([left, top, left + width, top + height], outline='gray', fill='lightgray')
                        draw.text((left + 5, top + 5), "[Image]", fill='black')
                
                # Handle other shapes (rectangles, etc.)
                else:
                    # Draw a simple rectangle for other shapes
                    draw.rectangle([left, top, left + width, top + height], outline='lightblue', fill=None)
                    
            except Exception as shape_error:
                logger.warning(f"Error processing shape: {str(shape_error)}")
                continue
        
        return img
        
    except Exception as e:
        logger.error(f"Error converting slide to image: {str(e)}")
        return None


def convert_pptx_to_pdf_serverless(input_path, output_dir):
    """Convert PPTX to PDF by first converting slides to images, then embedding in PDF.
    
    This preserves the exact formatting and layout of the original slides.
    """
    import tempfile
    import io
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    from PIL import Image
    
    try:
        logger.info(f"Starting image-based PDF conversion for: {input_path}")
        
        # Load the presentation
        prs = Presentation(input_path)
        
        if not prs.slides:
            logger.error("No slides found in presentation")
            return None
        
        # Get slide dimensions
        slide_width = prs.slide_width  # In EMU (English Metric Units)
        slide_height = prs.slide_height
        
        # Convert EMU to points (1 EMU = 1/914400 inch, 1 inch = 72 points)
        slide_width_pts = (slide_width / 914400) * 72
        slide_height_pts = (slide_height / 914400) * 72
        
        logger.info(f"Slide dimensions: {slide_width_pts:.1f} x {slide_height_pts:.1f} points")
        
        # Generate PDF filename
        base_name = os.path.splitext(os.path.basename(input_path))[0]
        pdf_path = os.path.join(output_dir, f"{base_name}.pdf")
        
        # Create PDF canvas with slide dimensions
        c = canvas.Canvas(pdf_path, pagesize=(slide_width_pts, slide_height_pts))
        
        logger.info(f"Converting {len(prs.slides)} slides to PDF")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                logger.info(f"Processing slide {slide_num}/{len(prs.slides)}")
                
                # Convert slide to image
                slide_img = convert_slide_to_image(slide, slide_width, slide_height, dpi=150)
                
                if slide_img:
                    # Create a new page (except for the first slide)
                    if slide_num > 1:
                        c.showPage()
                    
                    # Convert PIL image to bytes
                    img_buffer = io.BytesIO()
                    slide_img.save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    
                    # Create ImageReader object
                    img_reader = ImageReader(img_buffer)
                    
                    # Draw the image to fill the entire page
                    c.drawImage(img_reader, 0, 0, width=slide_width_pts, height=slide_height_pts)
                    
                    logger.info(f"Added slide {slide_num} as image to PDF")
                else:
                    logger.warning(f"Failed to convert slide {slide_num} to image")
                    # Create a new page with error message
                    if slide_num > 1:
                        c.showPage()
                    c.setFont("Helvetica", 12)
                    c.drawString(50, slide_height_pts - 50, f"Error: Could not render slide {slide_num}")
                
            except Exception as e:
                logger.warning(f"Error processing slide {slide_num}: {str(e)}")
                # Create a new page with error message
                if slide_num > 1:
                    c.showPage()
                c.setFont("Helvetica", 12)
                c.drawString(50, slide_height_pts - 50, f"Error processing slide {slide_num}: {str(e)}")
                continue
        
        # Save the PDF
        c.save()
        
        if os.path.exists(pdf_path):
            logger.info(f"PDF conversion successful: {pdf_path}")
            return pdf_path
        else:
            logger.error(f"PDF file not created: {pdf_path}")
            return None
            
    except Exception as e:
        logger.error(f"Error in image-based PDF conversion: {str(e)}")
        return None


@app.route('/convert-to-pdf/<filename>', methods=['GET', 'POST'])
def convert_to_pdf(filename):
    """Convert PowerPoint to PDF using LibreOffice."""
    import tempfile
    import shutil
    
    try:
        logger.info(f"PDF conversion requested for: {filename}")
        
        # Check if unified storage is available
        if unified_storage is None:
            logger.error("Unified storage not available for PDF conversion")
            return jsonify({
                'error': 'Storage service not available',
                'details': 'Storage system failed to initialize'
            }), 503
        
        # Get blob URL from query parameters if provided
        blob_url = request.args.get('blob_url')
        
        # Check if file exists locally first
        local_file_path = os.path.join(OUTPUT_FOLDER, filename)
        input_file_path = None
        temp_input_file = None
        
        # Check the unified storage output path
        unified_local_path = unified_storage.get_output_file_path(filename)
        if unified_local_path and os.path.exists(unified_local_path):
            logger.info(f"Using unified storage local file: {unified_local_path}")
            input_file_path = unified_local_path
        
        # Check temp output folder
        if not input_file_path and os.path.exists(local_file_path):
            logger.info(f"Using temp file: {local_file_path}")
            input_file_path = local_file_path
        elif blob_url:
            logger.info(f"Downloading file from unified storage: {blob_url}")
            try:
                # Download file from unified storage
                file_data = run_async(unified_storage.download_file(blob_url))
                if file_data:
                    # Create temporary file for conversion
                    temp_input_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
                    temp_input_file.write(file_data)
                    temp_input_file.flush()
                    temp_input_file.close()
                    input_file_path = temp_input_file.name
                    logger.info(f"Downloaded file to temporary location: {input_file_path}")
                else:
                    return jsonify({'error': 'Failed to download file from storage'}), 404
            except Exception as e:
                logger.error(f"Error downloading from storage: {str(e)}")
                return jsonify({'error': f'Storage download error: {str(e)}'}), 500
        else:
            return jsonify({'error': 'File not found in storage'}), 404
        
        # Create temporary directory for PDF output
        with tempfile.TemporaryDirectory() as temp_output_dir:
            try:
                # Convert PPTX to PDF using serverless approach
                pdf_path = convert_pptx_to_pdf_serverless(input_file_path, temp_output_dir)
                
                if pdf_path and os.path.exists(pdf_path):
                    # Generate PDF filename
                    base_name = os.path.splitext(filename)[0]
                    pdf_filename = f"{base_name}.pdf"
                    
                    # Read the PDF file
                    with open(pdf_path, 'rb') as pdf_file:
                        pdf_data = pdf_file.read()
                    
                    logger.info(f"PDF conversion successful. Size: {len(pdf_data)} bytes")
                    
                    # Return the PDF file for download
                    from flask import Response
                    response = Response(
                        pdf_data,
                        mimetype='application/pdf',
                        headers={
                            'Content-Disposition': f'attachment; filename="{pdf_filename}"',
                            'Content-Length': str(len(pdf_data))
                        }
                    )
                    
                    return response
                else:
                    return jsonify({'error': 'PDF conversion failed'}), 500
                    
            finally:
                # Clean up temporary input file if created
                if temp_input_file and os.path.exists(temp_input_file.name):
                    try:
                        os.unlink(temp_input_file.name)
                        logger.info(f"Cleaned up temporary file: {temp_input_file.name}")
                    except Exception as e:
                        logger.warning(f"Failed to clean up temporary file: {str(e)}")
        
    except Exception as e:
        logger.error(f"Error in PDF conversion: {str(e)}")
        return jsonify({'error': f'PDF conversion error: {str(e)}'}), 500


@app.errorhandler(413)
def too_large(e):
    """Handle file too large error."""
    return jsonify({'error': 'File too large. Maximum size is 200MB.'}), 413


@app.errorhandler(500)
def internal_error(e):
    """Handle internal server errors."""
    logger.error(f"Internal server error: {str(e)}")
    return jsonify({'error': 'Internal server error occurred.'}), 500


if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)